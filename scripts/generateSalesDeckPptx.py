#!/usr/bin/env python3
# ============================================================
# CORE Continuum — 6サービス 営業資料 PPTX
# 出力: ~/Desktop/CORE_Continuum_6サービス_営業資料.pptx
# ダークモード×ゴールド、Cinzel(英字見出し)×Noto Serif JP(日本語見出し)
# ============================================================
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── 色定義 (CORE Continuum ブランド: 金 #C9A24B) ───────────
BG = RGBColor(0x0A, 0x0A, 0x0E)
BG_SOFT = RGBColor(0x12, 0x12, 0x18)
INK = RGBColor(0xF7, 0xF4, 0xEC)
INK_DIM = RGBColor(0xA8, 0xA0, 0x94)
INK_FAINT = RGBColor(0x5A, 0x55, 0x4E)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
GOLD_SOFT = RGBColor(0xE3, 0xC9, 0x8A)
LINE = RGBColor(0x2A, 0x28, 0x24)

# 製品別アクセント色
C_RESONANCE = RGBColor(0x4E, 0x9E, 0x82)   # 翡翠系(LINE)
C_GUILD = RGBColor(0x6F, 0xA8, 0xFF)       # 青(DAO)
C_IRIS = RGBColor(0xE1, 0x5B, 0x9E)        # ピンク(Instagram)
C_LUME = RGBColor(0xF2, 0x9B, 0x4A)        # オレンジ(継続支援)
C_PRISM = RGBColor(0xB0, 0x7B, 0xD9)       # 紫(パーソナルAI)
C_CRYSTAL = RGBColor(0x8F, 0xB4, 0xE0)     # スチールブルー(コンシェルジュ)

FONT_DISPLAY = 'Noto Serif JP'
FONT_SANS = 'Noto Sans JP'
FONT_EN = 'Cinzel'
FONT_NUM = 'Inter'

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)
BLANK_LAYOUT = prs.slide_layouts[6]

TOTAL = 14


# ── ヘルパー ───────────────────────────
def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg


def add_text(slide, text, left, top, width, height,
             font=FONT_SANS, size=14, color=INK, bold=False, italic=False,
             align='center', anchor='middle', letter_spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
        else:
            p2 = tf.add_paragraph()
            p2.alignment = p.alignment
            if line_spacing:
                p2.line_spacing = line_spacing
            run = p2.add_run()
        if i == 0 and line_spacing:
            p.line_spacing = line_spacing
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if letter_spacing is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(letter_spacing))
    return tb


def add_circle(slide, cx, cy, r, fill_color=None, alpha=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill_color:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_color
        if alpha is not None:
            _set_alpha(sh, alpha)
    else:
        sh.fill.background()
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _set_alpha(sh, alpha):
    sp = sh.fill._xPr
    solidFill = sp.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            a = srgb.find(qn('a:alpha'))
            if a is None:
                a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_w=None, alpha=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_color
        if alpha is not None:
            _set_alpha(sh, alpha)
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_gradient_bar(slide, x, y, w, h):
    seg_w = w // 3
    add_rect(slide, x, y, seg_w, h, fill_color=GOLD)
    add_rect(slide, x + seg_w, y, seg_w, h, fill_color=GOLD_SOFT)
    add_rect(slide, x + 2 * seg_w, y, w - 2 * seg_w, h, fill_color=INK_FAINT)


def add_footer(slide, page_num):
    add_text(slide, 'CORE Continuum — 6 Services, One Ecosystem',
             Inches(0.6), Inches(SLIDE_H_IN - 0.45), Inches(6), Inches(0.3),
             font=FONT_SANS, size=8, color=INK_FAINT, align='left')
    add_text(slide, f'{page_num:02d} / {TOTAL}',
             Inches(SLIDE_W_IN - 1.6), Inches(SLIDE_H_IN - 0.45), Inches(1.0), Inches(0.3),
             font=FONT_NUM, size=9, color=INK_FAINT, align='right')


def draw_table(slide, x, y, col_widths, rows, row_h=Inches(0.5), header_color=GOLD,
               font_size=13, header_size=13):
    """シンプルな枠線テーブル描画。rows[0] はヘッダ。"""
    total_w = sum(col_widths)
    n_rows = len(rows)
    # 外枠
    add_rect(slide, x, y, total_w, row_h * n_rows, fill_color=BG_SOFT, line_color=LINE, alpha=0.6)
    for r, row in enumerate(rows):
        ry = y + row_h * r
        is_header = (r == 0)
        if is_header:
            add_rect(slide, x, ry, total_w, row_h, fill_color=GOLD, alpha=0.16)
        cx = x
        for c, cell in enumerate(row):
            cw = col_widths[c]
            add_text(slide, str(cell), cx + Inches(0.15), ry, cw - Inches(0.25), row_h,
                     font=FONT_SANS, size=header_size if is_header else font_size,
                     color=GOLD if is_header else INK, bold=is_header, align='left')
            cx += cw
        if r > 0:
            add_rect(slide, x, ry, total_w, Emu(1), fill_color=LINE)


# ── テンプレ ───────────────────────────
def slide_cover():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_circle(s, Emu(int(SLIDE_W_IN * 914400 / 2)), Emu(int(SLIDE_H_IN * 914400 / 2)),
               Inches(3.6), fill_color=GOLD, alpha=0.08)
    add_circle(s, Inches(2.2), Inches(1.4), Inches(1.6), fill_color=GOLD, alpha=0.10)
    add_circle(s, Inches(11.2), Inches(6.2), Inches(2.0), fill_color=GOLD_SOFT, alpha=0.08)
    add_text(s, 'CORE CONTINUUM', Inches(0.7), Inches(1.5), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=FONT_EN, size=16, color=GOLD, bold=True, align='left', letter_spacing=600)
    add_text(s, '6つのAIサービスで、\n事業のあらゆる接点を自動化する。', Inches(0.7), Inches(2.2),
             Inches(SLIDE_W_IN - 1.4), Inches(2.6), font=FONT_DISPLAY, size=48, color=INK, bold=True,
             align='left', line_spacing=1.15)
    add_gradient_bar(s, Inches(0.72), Inches(4.9), Inches(2.4), Inches(0.045))
    add_text(s, 'Resonance ・ Guild ・ Iris ・ Lume ・ Prism ・ Crystal',
             Inches(0.7), Inches(5.15), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=FONT_EN, size=15, color=INK_DIM, align='left', letter_spacing=200)
    add_text(s, '接客・コミュニティ・SNS運用・継続支援・秘書業務・顧客対応 ── その一つひとつに、専属のAIを。',
             Inches(0.7), Inches(5.7), Inches(SLIDE_W_IN - 1.4), Inches(0.6),
             font=FONT_SANS, size=14, color=INK_DIM, align='left')
    add_text(s, 'CORE (設立準備中) — 導入企業様向けご提案資料',
             Inches(0.7), Inches(6.7), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_SANS, size=11, color=INK_FAINT, align='left')
    add_footer(s, 1)
    return s


PRODUCTS = [
    dict(key='RESONANCE', name_jp='Resonance', color=C_RESONANCE, no='01',
         concept='公式LINEに届いたメッセージへ、AIが個別の自然な返信を下書き。\n確認して送るだけで、会話が続き関係が深まる。',
         area='LINE公式アカウント運用',
         features=[
             'AI返信 ── 相手の口調・会話の流れを読んで個別に返信案を自動生成',
             '返信キュー ── 熱量スコアで優先順位づけ、ワンクリック承認送信',
             'Googleカレンダー連携 ── 予約を自動認識し、確定でカレンダー＆Meet発行',
             '紹介機能 ── 満足したユーザーの口コミ流入を実データで可視化',
         ],
         diff='最大の差別化＝LINE無料枠(月200通)を1通も消費しない自動返信。何通返しても追加料金ゼロ。',
         pricing=[('Free', '¥0'), ('Solo', '¥1,980/月'), ('Pro', '¥6,980/月'), ('Business', '¥14,800/月')]),
    dict(key='GUILD', name_jp='Guild', color=C_GUILD, no='02',
         concept='提案→投票→決定のすべてが透明・改ざん検知可能。\n初対面のメンバーでも、今日から安心して協力できる。',
         area='DAO型コミュニティ運営OS',
         features=[
             '提案→投票 ── クアドラティック投票で決定、締切自動集計',
             '改ざん検知タイムライン ── ハッシュ連鎖で全記録を公開・独立検証可能',
             '順位表・レベル解放 ── 貢献ポイントで段階的にプロフィール装飾が解放',
             '共創フィードバック ── ユーザーの声を採用すると謝礼トークンを配布',
         ],
         diff='競合(Skool/Circle/Discord)にない「改ざん検知の透明性」が信頼の土台。',
         pricing=[('Free', '¥0'), ('Community', '¥980/月'), ('Pro', '¥3,980/月'), ('Enterprise', '個別見積')]),
    dict(key='IRIS', name_jp='Iris', color=C_IRIS, no='03',
         concept='アカウント分析から企画・台本・投稿・案件管理までを1つのAI相棒に統合。\nインフルエンサーの"仕事"を丸ごと引き受ける。',
         area='Instagram運用支援AI',
         features=[
             '企画・台本スタジオ ── テンプレ起点で白紙ゼロから台本を自動生成',
             '9マスグリッドプレビュー ── 公開前にフィード全体の統一感を確認・並べ替え',
             'アカウント分析 ── 反応率・伸びるテーマを学習し最適投稿時間を提案',
             '案件管理 ── 受注→下書き→投稿→レポートを一気通貫でサポート',
         ],
         diff='Agencyプランは運用代行会社向けにクライアント別ワークスペースを分離、30アカウント同時運用に対応。',
         pricing=[('Lite', '¥2,980/月'), ('Standard', '¥6,980/月'), ('Pro', '¥12,800/月'), ('Agency', '¥29,800/月')]),
    dict(key='LUME', name_jp='Lume', color=C_LUME, no='04',
         concept='プロフィールを30秒で作成し、そのページから直接「販売・予約・投げ銭」で収益化。\n訪問者分析とAIコーチングで、フォロワーを売上に変える。',
         area='プロフィール収益化 × 継続支援',
         features=[
             'ライブプレビュー編集 ── 保存なしで即反映、12種のテーマ',
             'クリックヒートマップ ── よく押されたリンクほど赤く光り可視化',
             '販売・予約・投げ銭カード ── URLを貼るだけでEC・予約に直結',
             '選べるストリーク ── 2/5/7/14日から自分で選択、罰なしで継続',
         ],
         diff='リンク集(Linktree等)で終わらず、プロフィールからその場で収益化できる点が独自。',
         pricing=[('無料', '7日間'), ('Pro', '¥1,480/月'), ('Business', '¥3,480/月')]),
    dict(key='PRISM', name_jp='Prism', color=C_PRISM, no='05',
         concept='蓄積した資料・メモ・数字を学習し、毎朝「今日の一手」を能動的に提案する専属エージェント。\n知識が増えるほど、提案の精度が上がる。',
         area='パーソナルAIエージェント(専属秘書)',
         features=[
             'ナレッジベース ── テキスト・URL・PDF・画像を保存しAIが自動要約',
             'プロアクティブ提案 ── 朝夜の巡回で「事業・財務・判断」に絞った一手を提示',
             '根拠チップ ── 実際に使った知識件数だけを正直に表示(嘘数字禁止)',
             'コマンドパレット(Cmd+K) ── ナレッジ検索・提案・実行を1本に集約',
         ],
         diff='「知識を足すほど提案が賢くなる」自己強化ループを、根拠チップで体感として可視化。',
         pricing=[('Starter', '¥2,980/月'), ('Standard', '¥9,800/月'), ('Exclusive', '¥29,800/月')]),
    dict(key='CRYSTAL', name_jp='Crystal', color=C_CRYSTAL, no='06',
         concept='話しかけるだけの24時間AIコンシェルジュ。\nブランドの言葉づかいのまま、夜間・休日の取りこぼしをゼロに。',
         area='AIコンシェルジュ・接客エージェント',
         features=[
             '24時間・365日の音声応対 ── マイクに話しかけるだけで字幕と声で返答',
             'ナレッジ貼るだけ学習＋FAQ自動生成 ── 会社案内を貼るだけでAIが学習',
             'AI SDR ── 見込み客の条件を設定し、会話の流れで自然に確認',
             '設置は3通り ── 専用リンク・タグ1行・設置代行、HTML不要',
         ],
         diff='Intercom Fin(解決1件$0.99)やQualified(年間約$68,000〜)と同じ中核機能を、月額定額・従量課金なしで提供。',
         pricing=[('Standard', '¥29,800/月 + 初期¥98,000'), ('Luxury', '¥49,800/月 + 初期¥298,000')]),
]


def slide_overview():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'THE ECOSYSTEM', Inches(0.7), Inches(0.55), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_EN, size=13, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, 'なぜ、6つなのか。', Inches(0.7), Inches(1.0), Inches(SLIDE_W_IN - 1.4), Inches(0.9),
             font=FONT_DISPLAY, size=34, color=INK, bold=True, align='left')
    add_text(s, '事業には、6つの接点がある。それぞれに専属のAIを置けば、人手を増やさずに全部が回りだす。',
             Inches(0.7), Inches(1.85), Inches(SLIDE_W_IN - 1.4), Inches(0.5),
             font=FONT_SANS, size=14, color=INK_DIM, align='left')

    card_w = Inches(3.95)
    card_h = Inches(1.95)
    gap_x = Inches(0.25)
    gap_y = Inches(0.25)
    start_x = Inches(0.7)
    start_y = Inches(2.55)
    for i, p in enumerate(PRODUCTS):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        add_rect(s, x, y, card_w, card_h, fill_color=p['color'], alpha=0.10, rounded=True)
        add_rect(s, x, y, Inches(0.07), card_h, fill_color=p['color'])
        add_text(s, p['no'], x + Inches(0.25), y + Inches(0.18), Inches(0.8), Inches(0.4),
                 font=FONT_NUM, size=12, color=p['color'], bold=True, align='left')
        add_text(s, p['name_jp'], x + Inches(0.25), y + Inches(0.42), card_w - Inches(0.5), Inches(0.5),
                 font=FONT_EN, size=22, color=INK, bold=True, align='left')
        add_text(s, p['area'], x + Inches(0.25), y + Inches(0.95), card_w - Inches(0.5), Inches(0.4),
                 font=FONT_SANS, size=11, color=p['color'], bold=True, align='left')
        add_text(s, p['concept'].split('\n')[0], x + Inches(0.25), y + Inches(1.35), card_w - Inches(0.5), Inches(0.5),
                 font=FONT_SANS, size=10.5, color=INK_DIM, align='left', anchor='top')
    add_footer(s, 2)
    return s


def slide_product(p, page):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    color = p['color']
    # 左帯アクセント
    add_rect(s, Inches(0), Inches(0), Inches(0.12), Inches(SLIDE_H_IN), fill_color=color)
    add_text(s, f"{p['no']}  /  {p['area']}", Inches(0.7), Inches(0.55), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_SANS, size=12, color=color, bold=True, letter_spacing=300, align='left')
    add_text(s, p['name_jp'].upper(), Inches(0.7), Inches(0.95), Inches(8.5), Inches(1.0),
             font=FONT_EN, size=48, color=INK, bold=True, align='left')
    add_text(s, p['concept'], Inches(0.7), Inches(1.95), Inches(7.6), Inches(1.1),
             font=FONT_SANS, size=15, color=INK_DIM, align='left', anchor='top', line_spacing=1.3)

    # 主要機能 4件
    add_text(s, 'KEY FEATURES', Inches(0.7), Inches(3.15), Inches(6), Inches(0.35),
             font=FONT_EN, size=11, color=color, bold=True, letter_spacing=400, align='left')
    fy = Inches(3.55)
    for i, feat in enumerate(p['features']):
        y = fy + Inches(0.62) * i
        add_circle(s, Inches(0.95), y + Inches(0.15), Inches(0.05), fill_color=color)
        add_text(s, feat, Inches(1.15), y, Inches(7.1), Inches(0.55),
                 font=FONT_SANS, size=12.5, color=INK, align='left', anchor='top')

    # 差別化
    diff_y = Inches(6.15)
    add_rect(s, Inches(0.7), diff_y, Inches(7.6), Inches(0.85), fill_color=color, alpha=0.10, rounded=True)
    add_text(s, '差別化ポイント', Inches(0.95), diff_y + Inches(0.08), Inches(2), Inches(0.3),
             font=FONT_SANS, size=10, color=color, bold=True, align='left')
    add_text(s, p['diff'], Inches(0.95), diff_y + Inches(0.34), Inches(7.1), Inches(0.48),
             font=FONT_SANS, size=11.5, color=INK, align='left', anchor='top')

    # 料金パネル(右側)
    px = Inches(8.7)
    pw = Inches(3.9)
    add_rect(s, px, Inches(1.95), pw, Inches(5.0), fill_color=BG_SOFT, line_color=LINE, rounded=True)
    add_text(s, 'PRICING', px + Inches(0.3), Inches(2.2), pw - Inches(0.6), Inches(0.35),
             font=FONT_EN, size=11, color=color, bold=True, letter_spacing=400, align='left')
    py = Inches(2.65)
    row_h = Inches(4.1) / max(len(p['pricing']), 1)
    for i, (plan, price) in enumerate(p['pricing']):
        ry = py + row_h * i
        if i > 0:
            add_rect(s, px + Inches(0.3), ry, pw - Inches(0.6), Emu(1), fill_color=LINE)
        add_text(s, plan, px + Inches(0.3), ry + Inches(0.12), pw - Inches(0.6), Inches(0.4),
                 font=FONT_SANS, size=14, color=INK_DIM, bold=True, align='left')
        add_text(s, price, px + Inches(0.3), ry + Inches(0.55), pw - Inches(0.6), Inches(0.55),
                 font=FONT_NUM, size=22, color=GOLD, bold=True, align='left')
    add_footer(s, page)
    return s


def slide_value():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'THE VALUE', Inches(0.7), Inches(0.6), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_EN, size=13, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, '6つに共通する、ひとつの価値。', Inches(0.7), Inches(1.05), Inches(SLIDE_W_IN - 1.4), Inches(0.9),
             font=FONT_DISPLAY, size=32, color=INK, bold=True, align='left')

    cards = [
        ('24時間、休まない', '深夜も休日も、応対・返信・提案が止まらない。夜間の機会損失をゼロに。'),
        ('人手を増やさない', '採用・教育コストなしで、対応量をそのまま増やせる。固定費を上げずに事業を伸ばす。'),
        ('月額定額・低コスト', '各サービス数千円〜。海外の従量課金AIエージェント(例: 解決1件$0.99〜)より予測しやすい。'),
        ('導入は最短1分〜', 'リンクを貼る・タグを1行入れるだけで動き出す。専門知識・HTML編集は不要。'),
    ]
    card_w = Inches(5.7)
    card_h = Inches(2.2)
    gap = Inches(0.3)
    start_x = Inches(0.7)
    start_y = Inches(2.3)
    accents = [GOLD, C_IRIS, C_GUILD, C_LUME]
    for i, (t, b) in enumerate(cards):
        row = i // 2
        col = i % 2
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        add_rect(s, x, y, card_w, card_h, fill_color=accents[i], alpha=0.10, rounded=True)
        add_circle(s, x + Inches(0.55), y + Inches(0.55), Inches(0.22), fill_color=accents[i])
        add_text(s, t, x + Inches(0.95), y + Inches(0.3), card_w - Inches(1.2), Inches(0.6),
                 font=FONT_DISPLAY, size=19, color=INK, bold=True, align='left')
        add_text(s, b, x + Inches(0.5), y + Inches(1.0), card_w - Inches(1.0), card_h - Inches(1.2),
                 font=FONT_SANS, size=12.5, color=INK_DIM, align='left', anchor='top', line_spacing=1.25)
    add_footer(s, 9)
    return s


def slide_usecase_matrix():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'WHO NEEDS WHAT', Inches(0.7), Inches(0.55), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_EN, size=13, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, '業種別・導入マップ', Inches(0.7), Inches(1.0), Inches(SLIDE_W_IN - 1.4), Inches(0.8),
             font=FONT_DISPLAY, size=32, color=INK, bold=True, align='left')

    rows = [
        ['業種・立場', 'おすすめサービス', '効く理由'],
        ['店舗・サロン・整体', 'Resonance / Crystal', '公式LINE返信の自動化＋サイトの接客をAIが24時間代行'],
        ['オンラインサロン・コミュニティ運営', 'Guild', '提案・投票・貢献の可視化で「関わりたくなる」場をつくる'],
        ['インフルエンサー・クリエイター', 'Iris / Lume', '投稿制作から収益化(販売・予約・投げ銭)までを一気通貫'],
        ['経営者・士業・フリーランス', 'Prism', '知識を蓄積するほど提案が賢くなる専属AI秘書'],
        ['高級不動産・ホテル・接客業', 'Crystal', 'ブランドの言葉づかいで24時間、見込み客を逃さない'],
        ['SNS運用代行会社', 'Iris (Agency)', 'クライアント別ワークスペースで複数アカウントを一元運用'],
    ]
    col_widths = [Inches(3.7), Inches(2.8), Inches(5.4)]
    draw_table(s, Inches(0.7), Inches(2.0), col_widths, rows, row_h=Inches(0.66), font_size=12.5, header_size=12.5)
    add_footer(s, 10)
    return s


def slide_pricing_all():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'PRICING OVERVIEW', Inches(0.7), Inches(0.5), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_EN, size=13, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, '全サービス 料金一覧', Inches(0.7), Inches(0.92), Inches(SLIDE_W_IN - 1.4), Inches(0.7),
             font=FONT_DISPLAY, size=30, color=INK, bold=True, align='left')

    rows = [
        ['サービス', 'エントリー', 'ミドル', 'ハイエンド'],
        ['Resonance', 'Solo ¥1,980/月', 'Pro ¥6,980/月', 'Business ¥14,800/月'],
        ['Guild', 'Free ¥0', 'Community ¥980/月', 'Pro ¥3,980/月'],
        ['Iris', 'Lite ¥2,980/月', 'Standard ¥6,980/月', 'Agency ¥29,800/月'],
        ['Lume', '無料(7日間)', 'Pro ¥1,480/月', 'Business ¥3,480/月'],
        ['Prism', 'Starter ¥2,980/月', 'Standard ¥9,800/月', 'Exclusive ¥29,800/月'],
        ['Crystal', '─', 'Standard ¥29,800/月', 'Luxury ¥49,800/月'],
    ]
    col_widths = [Inches(2.3), Inches(3.1), Inches(3.1), Inches(3.4)]
    draw_table(s, Inches(0.7), Inches(1.85), col_widths, rows, row_h=Inches(0.68), font_size=12.5, header_size=12.5)
    add_text(s, '※全サービス、初回は無料または7日間トライアルからお試しいただけます(Crystalのみ初期費用あり)。',
             Inches(0.7), Inches(6.75), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_SANS, size=10.5, color=INK_FAINT, align='left')
    add_footer(s, 11)
    return s


def slide_flow():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'HOW TO START', Inches(0.7), Inches(0.6), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_EN, size=13, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, '導入までの3ステップ', Inches(0.7), Inches(1.05), Inches(SLIDE_W_IN - 1.4), Inches(0.9),
             font=FONT_DISPLAY, size=32, color=INK, bold=True, align='left')

    steps = [
        ('01', 'まず無料で試す', '各サービスとも無料プランまたは7日間トライアルで、実際の画面・機能をすぐ体験できます。'),
        ('02', 'ブランドに合わせて設定', 'ブランド名・トーン・ナレッジ(会社案内など)を貼るだけで、その場でAIの応対が変わります。'),
        ('03', '設置してすぐ稼働', 'リンクを貼る・タグを1行入れるだけ。専門知識やHTML編集は不要、最短1分で動き始めます。'),
    ]
    card_w = Inches(3.85)
    gap = Inches(0.25)
    start_x = Inches(0.7)
    y = Inches(2.5)
    h = Inches(3.6)
    for i, (num, t, b) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        add_rect(s, x, y, card_w, h, fill_color=BG_SOFT, line_color=LINE, rounded=True)
        add_text(s, num, x + Inches(0.35), y + Inches(0.3), Inches(1.5), Inches(1.0),
                 font=FONT_NUM, size=44, color=GOLD, bold=True, align='left')
        add_text(s, t, x + Inches(0.35), y + Inches(1.35), card_w - Inches(0.7), Inches(0.7),
                 font=FONT_DISPLAY, size=20, color=INK, bold=True, align='left')
        add_text(s, b, x + Inches(0.35), y + Inches(2.1), card_w - Inches(0.7), Inches(1.3),
                 font=FONT_SANS, size=12.5, color=INK_DIM, align='left', anchor='top', line_spacing=1.3)
        if i < 2:
            add_text(s, '→', x + card_w + Inches(0.02), y + h / 2 - Inches(0.3), gap - Inches(0.04), Inches(0.6),
                     font=FONT_SANS, size=20, color=GOLD, align='center')
    add_footer(s, 12)
    return s


def slide_cta():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_circle(s, Emu(int(SLIDE_W_IN * 914400 / 2)), Emu(int(SLIDE_H_IN * 914400 / 2)),
               Inches(3.2), fill_color=GOLD, alpha=0.10)
    add_text(s, 'GET STARTED', Inches(0.7), Inches(1.6), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
             font=FONT_EN, size=14, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, 'まずは、無料でお試しください。', Inches(0.7), Inches(2.1), Inches(SLIDE_W_IN - 1.4), Inches(1.1),
             font=FONT_DISPLAY, size=40, color=INK, bold=True, align='left')
    add_text(s, '気になるサービスから1つ、あるいは事業の接点に合わせて複数を組み合わせてご導入いただけます。\nご不明点・お見積りのご相談は、メール1通でお気軽にどうぞ。',
             Inches(0.7), Inches(3.35), Inches(SLIDE_W_IN - 1.4), Inches(1.1),
             font=FONT_SANS, size=15, color=INK_DIM, align='left', anchor='top', line_spacing=1.4)
    add_gradient_bar(s, Inches(0.72), Inches(4.7), Inches(2.4), Inches(0.045))
    add_rect(s, Inches(0.7), Inches(5.1), Inches(6.2), Inches(0.9), fill_color=BG_SOFT, line_color=LINE, rounded=True)
    add_text(s, 'お問い合わせ', Inches(1.0), Inches(5.25), Inches(2), Inches(0.3),
             font=FONT_SANS, size=10, color=GOLD, bold=True, align='left')
    add_text(s, 'core.guild.inc@gmail.com', Inches(1.0), Inches(5.55), Inches(5.6), Inches(0.4),
             font=FONT_NUM, size=18, color=INK, bold=True, align='left')
    add_footer(s, 13)
    return s


def slide_closing():
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_circle(s, Emu(int(SLIDE_W_IN * 914400 / 2)), Emu(int(SLIDE_H_IN * 914400 / 2)),
               Inches(3.8), fill_color=GOLD, alpha=0.07)
    add_text(s, 'CORE CONTINUUM', Inches(0.5), Inches(2.7), Inches(SLIDE_W_IN - 1), Inches(0.9),
             font=FONT_EN, size=40, color=INK, bold=True, letter_spacing=300)
    add_text(s, '事業のあらゆる接点に、途切れない知性を。', Inches(0.5), Inches(3.75), Inches(SLIDE_W_IN - 1), Inches(0.7),
             font=FONT_DISPLAY, size=22, color=GOLD_SOFT, italic=True)
    add_gradient_bar(s, Inches(SLIDE_W_IN / 2 - 1.2), Inches(4.7), Inches(2.4), Inches(0.045))
    add_footer(s, 14)
    return s


# ── 生成 ───────────────────────────
slide_cover()
slide_overview()
for i, p in enumerate(PRODUCTS):
    slide_product(p, 3 + i)
slide_value()
slide_usecase_matrix()
slide_pricing_all()
slide_flow()
slide_cta()
slide_closing()

out_path = os.path.expanduser('~/Desktop/CORE_Continuum_6サービス_営業資料.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
print(f'Slides: {len(prs.slides.__iter__.__self__._sldIdLst)}')
