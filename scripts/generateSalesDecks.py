#!/usr/bin/env python3
# ============================================================
# CORE 営業資料 — 4プロダクト個別 + 統合 (計5ファイル)
# コーポレートサイト (/corp) のデザイン言語:
#   黒基調・明朝の巨大見出し・グラデアクセント・光と核のモチーフ・1スライド1メッセージ
# 出力: ~/Desktop/CORE_<name>_営業資料_2026-06-09.pptx
# ============================================================
import math, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

DESKTOP = os.path.expanduser('~/Desktop/CORE 営業資料')
DATE = '2026-06-09'

SW, SH = 13.333, 7.5  # 16:9 inches

# ── フォント (macOS 標準・確実に描画) ──
F_DISP = 'Hiragino Mincho ProN'   # 明朝 = 編集的な大見出し
F_SANS = 'Hiragino Sans'          # 本文・ラベル
F_NUM  = 'Helvetica Neue'         # 数字・英字

def hx(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

BG       = hx('#06060e')
INK      = hx('#ffffff')
INK_DIM  = hx('#aeb4c2')
INK_FAINT= hx('#5a6172')
WHITE    = hx('#ffffff')

# ── プロダクト別パレット ──
PAL = {
    'prism': dict(primary='#a78bfa', bright='#c4b5fd', deep='#7c3aed',
                  grad=['#ff5757', '#ff9842', '#fbbf24', '#4ade80', '#60a5fa', '#a78bfa', '#f472b6'],
                  en='CORE PRISM', mark='prism'),
    'iris': dict(primary='#E1306C', bright='#F472B6', deep='#833AB4',
                 grad=['#FCB045', '#E1306C', '#833AB4'], en='CORE IRIS', mark='iris'),
    'resonance': dict(primary='#06C755', bright='#34D399', deep='#0EA5E9',
                      grad=['#34D399', '#06C755', '#0EA5E9'], en='CORE RESONANCE', mark='resonance'),
    'lume': dict(primary='#FFA42A', bright='#FFD86B', deep='#FF7A18',
                 grad=['#FFD86B', '#FFA42A', '#FF7A18'], en='CORE LUME', mark='lume'),
    'core': dict(primary='#38BDF8', bright='#7DD3FC', deep='#0EA5E9',
                 grad=['#7DD3FC', '#E0F2FE', '#38BDF8'], en='CORE INC.', mark='core'),
}

# ── ビルド状態 (デッキごとに差し替え) ──
prs = None
P = None      # current palette (hex strings)
PAGE = 0
TOTAL = 0
FOOTER_L = '株式会社コア  ·  CORE Inc.'


def newdeck(palette_key):
    global prs, P, PAGE
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    P = {k: (hx(v) if k in ('primary', 'bright', 'deep') else v) for k, v in PAL[palette_key].items()}
    P['grad_rgb'] = [hx(c) for c in PAL[palette_key]['grad']]
    PAGE = 0


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── 図形ヘルパー (座標は inch) ──
def bg(s, color=BG):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _alpha(sh, alpha):
    sp = sh.fill._xPr
    sf = sp.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            a = srgb.find(qn('a:alpha'))
            if a is None:
                a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))


def rect(s, x, y, w, h, fill=None, line=None, lw=None, alpha=None, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None: _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw or 1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def circle(s, cx, cy, r, fill=None, line=None, lw=None, alpha=None):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(r * 2), Inches(r * 2))
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None: _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw or 1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def ellipse(s, cx, cy, w, h, fill=None, line=None, lw=None, alpha=None, rot=0):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - w / 2), Inches(cy - h / 2), Inches(w), Inches(h))
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None: _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(lw or 1)
    else:
        sh.line.fill.background()
    sh.rotation = rot
    sh.shadow.inherit = False
    return sh


def tri(s, cx, cy, w, h, fill=None, alpha=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(cx - w / 2), Inches(cy - h / 2), Inches(w), Inches(h))
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None: _alpha(sh, alpha)
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def line(s, x1, y1, x2, y2, color, w=1.2, dash=None):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = color; cn.line.width = Pt(w)
    cn.shadow.inherit = False
    if dash:
        ln = cn.line._get_or_add_ln()
        d = etree.SubElement(ln, qn('a:prstDash')); d.set('val', dash)
    return cn


def txt(s, text, x, y, w, h, font=F_SANS, size=14, color=INK, bold=False, italic=False,
        align='center', anchor='middle', spc=None, leading=None):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    al = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.CENTER)
    for i, ln in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if leading is not None:
            p.line_spacing = leading
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic
        if spc is not None:
            r._r.get_or_add_rPr().set('spc', str(spc))
    return tb


def grad_bar(s, x, y, w, h, colors=None):
    colors = colors or P['grad_rgb']
    seg = w / len(colors)
    for i, c in enumerate(colors):
        rect(s, x + i * seg, y, seg + 0.01, h, fill=c)


def footer(s):
    txt(s, FOOTER_L, 0.5, SH - 0.46, 6, 0.3, font=F_SANS, size=8, color=INK_FAINT, align='left')
    txt(s, f'{PAGE:02d} / {TOTAL}', SW - 1.6, SH - 0.46, 1.1, 0.3, font=F_NUM, size=9, color=INK_FAINT, align='right')


def glow(s, cx, cy, r, color, alpha=0.16):
    circle(s, cx, cy, r, fill=color, alpha=alpha)


# ── ブランドマーク描画 ──
def mark(s, kind, cx, cy, size):
    r = size / 2
    pr = P['primary']; br = P['bright']; dp = P['deep']
    glow(s, cx, cy, size * 0.95, pr, 0.14)
    if kind == 'prism':
        tri(s, cx, cy + size * 0.04, size * 1.02, size * 0.92, fill=pr, alpha=0.92)
        line(s, cx, cy - r * 0.85, cx, cy + r * 0.42, WHITE, 1.0)
        line(s, cx, cy - r * 0.85, cx + r * 0.42, cy + r * 0.06, WHITE, 0.8)
    elif kind == 'iris':
        for k in range(6):
            rot = k * 60
            d = size * 0.24
            px = cx + math.sin(math.radians(rot)) * d
            py = cy - math.cos(math.radians(rot)) * d
            ellipse(s, px, py, size * 0.17, size * 0.52, line=pr, lw=2.0, rot=rot)
        circle(s, cx, cy, size * 0.05, fill=pr)
    elif kind == 'resonance':
        # 同心の波紋 + 源の点
        for i, rr in enumerate([0.95, 0.66, 0.37]):
            circle(s, cx, cy, r * rr, line=P['grad_rgb'][min(i, len(P['grad_rgb']) - 1)], lw=2.4 - i * 0.5)
        circle(s, cx, cy, size * 0.08, fill=pr)
    elif kind == 'lume':
        rect(s, cx - r, cy - r, size, size, fill=dp, rounded=True)
        rect(s, cx - r, cy - r, size, size, fill=pr, alpha=0.5, rounded=True)
        circle(s, cx, cy - size * 0.03, size * 0.27, fill=WHITE)
        circle(s, cx - size * 0.09, cy - size * 0.13, size * 0.055, fill=WHITE)
    elif kind == 'core':
        circle(s, cx, cy, r * 0.95, line=br, lw=2.4)
        ellipse(s, cx, cy, size * 0.95, size * 0.36, line=pr, lw=2.0, rot=-24)
        ellipse(s, cx, cy, size * 0.95, size * 0.36, line=pr, lw=1.6, rot=34)
        circle(s, cx, cy, size * 0.16, fill=br)
        circle(s, cx - size * 0.05, cy - size * 0.05, size * 0.05, fill=WHITE)


# ============================================================
#  スライド・テンプレート
# ============================================================
def _adv():
    global PAGE
    PAGE += 1


def s_cover(eyebrow, title, sub, note, mark_kind):
    _adv(); s = blank(); bg(s)
    glow(s, SW / 2, SH / 2, 4.0, P['primary'], 0.13)
    glow(s, 3.0, 2.0, 2.2, P['bright'], 0.08)
    glow(s, 10.5, 5.6, 2.4, P['deep'], 0.10)
    mark(s, mark_kind, SW / 2, 2.15, 1.5)
    txt(s, eyebrow, 0.5, 3.15, SW - 1, 0.4, font=F_NUM, size=12.5, color=P['bright'], bold=True, spc=600)
    txt(s, title, 0.5, 3.65, SW - 1, 1.9, font=F_DISP, size=52, color=INK, bold=True, leading=1.1)
    txt(s, sub, 0.5, 5.55, SW - 1, 0.6, font=F_DISP, size=20, color=INK_DIM, italic=True)
    grad_bar(s, SW / 2 - 1.1, 6.25, 2.2, 0.045)
    txt(s, note, 0.5, 6.55, SW - 1, 0.5, font=F_SANS, size=12.5, color=INK_DIM)
    footer(s)


def s_section(num, en, jp):
    _adv(); s = blank(); bg(s)
    glow(s, 3.4, 3.0, 3.0, P['primary'], 0.10)
    txt(s, f'{num:02d}', 0.5, 0.7, SW - 1, 2.6, font=F_NUM, size=200, color=INK_FAINT, bold=True)
    txt(s, en, 0.5, 3.55, SW - 1, 0.5, font=F_NUM, size=13, color=P['bright'], bold=True, spc=700)
    txt(s, jp, 0.5, 4.05, SW - 1, 1.5, font=F_DISP, size=62, color=INK, bold=True)
    grad_bar(s, SW / 2 - 1.0, 6.1, 2.0, 0.05)
    footer(s)


def s_keyword(big, small, color=None, size=130):
    _adv(); s = blank(); bg(s)
    glow(s, SW / 2, SH / 2, 3.6, P['primary'], 0.10)
    txt(s, big, 0.4, 1.7, SW - 0.8, 3.4, font=F_DISP, size=size, color=color or INK, bold=True, leading=1.08)
    if small:
        txt(s, small, 1.0, 5.55, SW - 2, 1.0, font=F_SANS, size=19, color=INK_DIM)
    footer(s)


def s_title_body(eyebrow, title, body):
    _adv(); s = blank(); bg(s)
    txt(s, eyebrow, 0.9, 1.15, SW - 1.8, 0.4, font=F_NUM, size=11.5, color=P['bright'], bold=True, spc=600, align='left')
    txt(s, title, 0.9, 1.6, SW - 1.8, 1.5, font=F_DISP, size=42, color=INK, bold=True, align='left', leading=1.15)
    grad_bar(s, 0.92, 3.15, 1.4, 0.04)
    txt(s, body, 0.9, 3.55, SW - 1.8, 3.0, font=F_SANS, size=20, color=INK_DIM, align='left', anchor='top', leading=1.7)
    footer(s)


def s_cards4(eyebrow, title, cards):
    _adv(); s = blank(); bg(s)
    txt(s, eyebrow, 0.9, 0.85, SW - 1.8, 0.4, font=F_NUM, size=11.5, color=P['bright'], bold=True, spc=600, align='left')
    txt(s, title, 0.9, 1.3, SW - 1.8, 1.0, font=F_DISP, size=36, color=INK, bold=True, align='left')
    gx, gy, gw, gh, gap = 0.9, 2.7, 5.65, 1.75, 0.35
    for i, (t, b) in enumerate(cards):
        x = gx + (i % 2) * (gw + gap)
        y = gy + (i // 2) * (gh + gap)
        rect(s, x, y, gw, gh, fill=hx('#0e0e1a'), line=P['primary'], lw=0.75, rounded=True)
        circle(s, x + 0.55, y + 0.62, 0.12, fill=P['primary'])
        txt(s, t, x + 0.95, y + 0.28, gw - 1.2, 0.7, font=F_DISP, size=17.5, color=INK, bold=True, align='left', anchor='middle')
        txt(s, b, x + 0.4, y + 1.02, gw - 0.8, 0.6, font=F_SANS, size=12.5, color=INK_DIM, align='left', anchor='top', leading=1.45)
    footer(s)


def s_steps(eyebrow, title, steps):
    _adv(); s = blank(); bg(s)
    txt(s, eyebrow, 0.9, 0.9, SW - 1.8, 0.4, font=F_NUM, size=11.5, color=P['bright'], bold=True, spc=600, align='left')
    txt(s, title, 0.9, 1.35, SW - 1.8, 1.0, font=F_DISP, size=36, color=INK, bold=True, align='left')
    n = len(steps); gw = 3.55; gap = (SW - 1.8 - gw * n) / (n - 1) if n > 1 else 0
    y = 3.0; gh = 2.7
    for i, (num, t, b) in enumerate(steps):
        x = 0.9 + i * (gw + gap)
        rect(s, x, y, gw, gh, fill=hx('#0c0c16'), line=P['primary'], lw=0.7, rounded=True)
        txt(s, num, x + 0.35, y + 0.3, 1.5, 0.9, font=F_NUM, size=44, color=P['bright'], bold=True, align='left')
        txt(s, t, x + 0.38, y + 1.25, gw - 0.7, 0.6, font=F_DISP, size=18, color=INK, bold=True, align='left')
        txt(s, b, x + 0.38, y + 1.85, gw - 0.7, 0.7, font=F_SANS, size=12.5, color=INK_DIM, align='left', anchor='top', leading=1.45)
        if i < n - 1:
            txt(s, '→', x + gw, y + gh / 2 - 0.3, gap, 0.6, font=F_SANS, size=22, color=P['primary'])
    footer(s)


def s_twocol(title, ll, lb, rl, rb):
    _adv(); s = blank(); bg(s)
    txt(s, title, 0.9, 1.0, SW - 1.8, 1.0, font=F_DISP, size=34, color=INK, bold=True, align='left')
    midx = SW / 2
    rect(s, 0.9, 2.5, 5.4, 3.6, fill=hx('#0c0c16'), line=INK_FAINT, lw=0.6, rounded=True)
    rect(s, midx + 0.3, 2.5, 5.4, 3.6, fill=hx('#0e0e1a'), line=P['primary'], lw=1.0, rounded=True)
    txt(s, ll, 1.25, 2.85, 4.8, 0.5, font=F_SANS, size=12.5, color=INK_FAINT, bold=True, spc=400, align='left')
    txt(s, lb, 1.25, 3.5, 4.8, 2.4, font=F_SANS, size=16.5, color=INK_DIM, align='left', anchor='top', leading=1.6)
    txt(s, rl, midx + 0.65, 2.85, 4.8, 0.5, font=F_SANS, size=12.5, color=P['bright'], bold=True, spc=400, align='left')
    txt(s, rb, midx + 0.65, 3.5, 4.8, 2.4, font=F_SANS, size=16.5, color=INK, align='left', anchor='top', leading=1.6)
    footer(s)


def s_stat(big, label, sub):
    _adv(); s = blank(); bg(s)
    glow(s, SW / 2, 3.2, 3.0, P['primary'], 0.12)
    txt(s, label, 0.5, 1.7, SW - 1, 0.5, font=F_SANS, size=15, color=INK_DIM, bold=True, spc=500)
    txt(s, big, 0.5, 2.2, SW - 1, 2.6, font=F_NUM, size=200, color=P['bright'], bold=True)
    if sub:
        txt(s, sub, 1.0, 5.35, SW - 2, 1.0, font=F_SANS, size=19, color=INK, leading=1.5)
    footer(s)


def s_pricing(title, tiers, foot):
    _adv(); s = blank(); bg(s)
    txt(s, title, 0.9, 0.95, SW - 1.8, 1.0, font=F_DISP, size=36, color=INK, bold=True, align='left')
    n = len(tiers); gap = 0.3
    gw = (SW - 1.8 - gap * (n - 1)) / n
    y = 2.55; gh = 3.1
    for i, (name, price, note, hot) in enumerate(tiers):
        x = 0.9 + i * (gw + gap)
        rect(s, x, y, gw, gh, fill=hx('#0e0e1a') if hot else hx('#0b0b14'),
             line=P['primary'] if hot else INK_FAINT, lw=1.4 if hot else 0.6, rounded=True)
        if hot:
            rect(s, x, y, gw, 0.06, fill=P['primary'])
            txt(s, '人気', x, y + 0.18, gw, 0.3, font=F_SANS, size=10, color=P['bright'], bold=True, spc=300)
        txt(s, name, x, y + 0.6, gw, 0.45, font=F_DISP, size=17, color=INK, bold=True)
        txt(s, price, x, y + 1.15, gw, 0.7, font=F_NUM, size=27, color=P['bright'] if hot else INK, bold=True)
        txt(s, note, x + 0.25, y + 1.95, gw - 0.5, 1.0, font=F_SANS, size=11.5, color=INK_DIM, anchor='top', leading=1.4)
    txt(s, foot, 0.9, 5.95, SW - 1.8, 0.5, font=F_SANS, size=14, color=P['bright'], bold=True)
    footer(s)


def s_cta(big, sub, mark_kind):
    _adv(); s = blank(); bg(s)
    glow(s, SW / 2, SH / 2, 4.2, P['primary'], 0.16)
    mark(s, mark_kind, SW / 2, 2.2, 1.15)
    txt(s, big, 0.5, 3.25, SW - 1, 1.8, font=F_DISP, size=46, color=INK, bold=True, leading=1.12)
    grad_bar(s, SW / 2 - 1.1, 5.1, 2.2, 0.045)
    txt(s, sub, 0.5, 5.4, SW - 1, 1.0, font=F_NUM, size=15, color=INK_DIM, leading=1.6)
    footer(s)


# ============================================================
#  プロダクト個別デッキ (共通テンプレ)
# ============================================================
def build_product(key, data):
    global TOTAL, PAGE, FOOTER_L
    newdeck(key)
    FOOTER_L = f"株式会社コア  ·  {data['en']}"
    slides = data['slides']
    TOTAL = len(slides)
    PAGE = 0
    mk = PAL[key]['mark']
    for kind, args in slides:
        if kind == 'cover':   s_cover(*args, mk)
        elif kind == 'kw':    s_keyword(*args)
        elif kind == 'tb':    s_title_body(*args)
        elif kind == 'cards': s_cards4(*args)
        elif kind == 'steps': s_steps(*args)
        elif kind == '2col':  s_twocol(*args)
        elif kind == 'stat':  s_stat(*args)
        elif kind == 'price': s_pricing(*args)
        elif kind == 'cta':   s_cta(args[0], args[1], mk)
    out = f"{DESKTOP}/CORE_{data['file']}_営業資料_{DATE}.pptx"
    prs.save(out)
    print('saved', out, f'({TOTAL}p)')


# ── 各プロダクトの中身 ──
PRODUCTS = {
'prism': dict(en='CORE PRISM', file='Prism', slides=[
    ('cover', ('CORE PRISM  /  事業家のためのAIエージェントOS',
               'すべての事業を、\nひとつの頭脳で。',
               'One mind for your whole business.',
               '経営の司令塔となる、13名のAIエージェント。')),
    ('kw', ('社長は、ひとりで\n全部やっている。', '営業も、財務も、契約も、議事録も。')),
    ('tb', ('THE PROBLEM', 'その仕事、AIに分担できます。',
            '経営者の中には、いくつもの役割が同居しています。決断する者、売る者、数字を読む者、つくる者。\nPrism は、その一つひとつを専属の AI エージェントとして外に取り出す経営 OS です。')),
    ('cards', ('FEATURES', '七つの役割に、七人の専属AI。', [
        ('七つの役割・七人のAI', '経営・営業・財務・創造… 役割ごとに専門のエージェントが担当'),
        ('商談から契約まで一気通貫', '議事録・提案書・財務・契約レビューを下書きまで自動化'),
        ('ひとつの横断検索', 'すべての文脈に、ひとことでアクセス'),
        ('朝と晩の能動提案', '「次にやるべきこと」を、AIの方から差し出す'),
    ])),
    ('steps', ('HOW IT WORKS', '使い方は、三つだけ。', [
        ('01', '話す・渡す', '音声・ファイル・画像を、そのまま投げる'),
        ('02', 'AIが仕上げる', '議事録も提案も契約レビューも、下書きまで完成'),
        ('03', '確認して送る', 'あなたは、目を通して送り出すだけ'),
    ])),
    ('2col', ('「ばらばら」から、「ひとつ」へ。',
              '従来', 'コンサル¥30万＋秘書＋ChatGPT＋Notion＋議事録ツール…\nそれぞれ別契約・別管理。',
              'CORE Prism', '13名のAI役員が、ひとつのOSで並走。\n月¥4,800から、7日間無料。')),
    ('stat', ('1/7', 'コンサル代を', '顧問・秘書・各種SaaSを、ひとつのOSに置き換える')),
    ('tb', ('CONNECTED', 'Prism は、司令塔。',
            'Instagram の Iris、LINE の Resonance、リンクの Lume。\n三つの SNS ツールが集めたお客様の動きは、最後にすべて Prism へ集まり、13 名の AI 役員が次の一手まで提案します。')),
    ('price', ('料金 — プランごとの「できること」。', [
        ('Free', '¥0', '7日間 全機能おためし', False),
        ('Starter', '¥4,800', '基本AI・1人格・個人向け', False),
        ('Standard', '¥9,800', '全AI＋商談AI・人格無制限', True),
        ('Exclusive', '¥29,800', '専任CS・法人連携・導入伴走', False),
    ], '全プラン 7日間無料・クレジットカード登録不要 ／ 法人Studioは要問合せ')),
    ('cta', ('あなたは、\n決めるだけ。', 'core-prism-app.vercel.app    ·    hello@core-inc.jp')),
]),

'iris': dict(en='CORE IRIS', file='Iris', slides=[
    ('cover', ('CORE IRIS  /  インフルエンサーのためのAI',
               'Instagramを、\nAIと育てる。',
               'Run Instagram with an AI agent.',
               '投稿も、分析も、案件も。あなたの6人のAIチーム。')),
    ('kw', ('フォロワー数より、\n案件数。', '「いいね」は、仕事にならない。')),
    ('tb', ('THE PROBLEM', 'Instagram運用の全部を、AIに。',
            '届いた反応、寄せられた感情、流れていく数字。そのすべてを受けとめ、次の一手に変える。\nインフルエンサーが抱える六つの仕事を、Iris はひとつのアプリにまとめます。')),
    ('cards', ('FEATURES', '六つの仕事を、ひとつに。', [
        ('投稿AI', '構成・テロップ・キャプション・ハッシュタグを丸ごと下書き'),
        ('Instagram解析', '数字を踏まえて、AIが次の一手まで決める'),
        ('案件管理', 'スクショを渡せば、AIが3秒で自動入力'),
        ('交渉文AI', '返信文・断り文・カウンターオファーを生成'),
    ])),
    ('steps', ('HOW IT WORKS', '渡して、確認して、投稿する。', [
        ('01', 'ネタを渡す', 'スクショ・写真・思いつきを、そのまま'),
        ('02', 'AIが下書き', '投稿も返信も戦略も、AIが先に書く'),
        ('03', '確認して出す', '整えて、投稿・送信するだけ'),
    ])),
    ('2col', ('「見る」だけでは、伸びない。',
              '従来の分析ツール', '数字を見るだけ。\n案件管理も、交渉も、投稿制作もなし。',
              'CORE Iris', '数字を踏まえて、AIが次の一手まで決める。\n月¥2,800から、7日間無料。')),
    ('tb', ('CONNECTED', 'Iris の発見は、流れていく。',
            'ここで掴んだファンの反応は、そのまま Resonance の LINE 配信や、司令塔 Prism の経営判断へとつながります。\n発見を、行動に変えるまでが一続きです。')),
    ('price', ('料金 — プランごとの「できること」。', [
        ('Free', '¥0', '7日間 全機能おためし', False),
        ('Lite', '¥2,800', 'AI相談30回・案件管理無制限', False),
        ('Standard', '¥6,800', 'AI相談/解析ほぼ無制限', True),
        ('Pro', '¥9,800', 'チーム5名・ブランドマッチ', False),
        ('Studio', '¥29,800', '無制限・API・ホワイトラベル', False),
    ], '全プラン 7日間無料・クレジットカード登録不要')),
    ('cta', ('数字に、\n振り回されない。', 'core-prism-app.vercel.app/iris    ·    hello@core-inc.jp')),
]),

'resonance': dict(en='CORE RESONANCE', file='Resonance', slides=[
    ('cover', ('CORE RESONANCE  /  店舗・サロン・教室のためのAI',
               'LINEのご縁を、\nAIが温める。',
               'Let it resonate.',
               '一斉配信なのに、一人ひとりに。')),
    ('kw', ('同じ文を、全員に。\nもう、やめにする。', 'お客様は、一人ひとり違う。')),
    ('tb', ('THE PROBLEM', '一人ひとりに、AIが書き分ける。',
            '名簿の一人ひとりに、その人のための一文を AI が書き分け、LINE で手紙のように届ける個別配信。\n「また会いたい」を、静かに育てます。')),
    ('cards', ('FEATURES', '手紙のような、一斉配信。', [
        ('個別文面AI', '一人ひとりに、別の言葉を書き分ける'),
        ('承認制', '送る前に、必ず全件を確認できる安心設計'),
        ('LINE公式に接続', 'お持ちのLINE公式アカウントに、そのまま'),
        ('AIレター', '会話の履歴を読み、一人ずつ下書きを用意'),
    ])),
    ('steps', ('HOW IT WORKS', 'つないで、確認して、届く。', [
        ('01', 'LINEをつなぐ', 'お持ちのLINE公式アカウントを接続'),
        ('02', 'AIが書き分け', '名簿ごとに、最適な文面を下書き'),
        ('03', '全件確認して送る', '一人ひとりに、手紙のように届く'),
    ])),
    ('2col', ('「同報」では、心は動かない。',
              '従来の配信ツール', 'セグメントに同じ文。\nLステップ¥5,000〜・エルメ¥10,780〜。',
              'CORE Resonance', '一人ひとり別の文面＋送信前の全件確認。\n月¥980から、7日間無料。')),
    ('tb', ('CONNECTED', '誰に届けるかは、仲間が教える。',
            'Iris（Instagram）や Lume（リンク）が見つけた「いま関心のある人」へ、最適なタイミングで届きます。\nそして、その結果はすべて Prism に集まります。')),
    ('price', ('料金 — プランごとの「できること」。', [
        ('7日間無料', '¥0', '全機能を7日間おためし', False),
        ('Pro', '¥980', 'AI個別配信・1アカ・月3,000通', True),
        ('Business', '¥2,980', '3アカウント・月15,000通', False),
        ('Premium', '¥4,980', '10アカウント・月50,000通', False),
    ], '全プラン 7日間無料・自分のLINE/Claude鍵で原価ほぼ0（BYOK）')),
    ('cta', ('売り込みでは、ない。\nご縁を、温める。', 'resonancebot-ivory.vercel.app    ·    hello@core-inc.jp')),
]),

'lume': dict(en='CORE LUME', file='Lume', slides=[
    ('cover', ('CORE LUME  /  クリエイターのためのリンクハブ',
               'すべてのリンクを、\nひとつに。',
               'Every link, in one place.',
               'あなたのリンクを、いちばん美しく。')),
    ('kw', ('プロフィールは、\n一行しかない。', 'その一行に、全部を込める。')),
    ('tb', ('THE PROBLEM', 'その一行に、全部をまとめる。',
            'あなたのすべてのリンクを束ねるハブ。誰が、どこから、どのリンクに触れたのか。\nクリエイターの「いま」を、色と熱で映し出す、リンクまとめの新しい基準です。')),
    ('cards', ('FEATURES', '美しさと、データの両立。', [
        ('30秒で完成', '5つのテーマで、すぐに整うプロフィール'),
        ('クリックヒートマップ', '押された比率を、熱で可視化する'),
        ('流入元クロス分析', 'どこから来て、何を踏んだかが分かる'),
        ('時間帯・傾向', '踏まれる時間まで、ひと目で見える'),
    ])),
    ('steps', ('HOW IT WORKS', '並べるだけで、見えてくる。', [
        ('01', 'リンクを並べる', 'すべてのリンクを、ひとつに集約'),
        ('02', '美しく仕上がる', '5つのテーマで、プロフィール完成'),
        ('03', '色で、分かる', '誰がどこを踏んだかを、熱で可視化'),
    ])),
    ('2col', ('「並べる」だけでは、もったいない。',
              '従来のリンクまとめ', 'ただ並べるだけ。\n誰が踏んだのか、分からない。',
              'CORE Lume', '美しく、そしてクリックを熱で可視化。\n月¥980から、7日間無料。')),
    ('tb', ('CONNECTED', 'そのクリックは、財産になる。',
            '誰がどのリンクを踏んだか。その流れは、Iris・Resonance・司令塔 Prism すべての判断材料になります。\nリンクは、ただの入口ではありません。')),
    ('price', ('料金 — プランごとの「できること」。', [
        ('7日間無料', '¥0', '全機能を7日間おためし', False),
        ('Pro', '¥980', 'ヒートマップ・流入元・時間帯解析', True),
        ('Business', '¥2,980', '複数プロフィール管理', False),
    ], '全プラン 7日間無料・クレジットカード登録不要')),
    ('cta', ('あなたのリンクを、\nいちばん美しく。', 'lume-deploy-five.vercel.app    ·    hello@core-inc.jp')),
]),
}

for k, d in PRODUCTS.items():
    build_product(k, d)


# ============================================================
#  統合デッキ — すべてがつながる
# ============================================================
def build_combined():
    global TOTAL, PAGE, FOOTER_L, P
    newdeck('core')
    FOOTER_L = '株式会社コア  ·  CORE Inc.  ·  すべての時代の、核となるものを。'
    TOTAL = 20
    PAGE = 0

    # 1 cover
    s_cover('CORE INC.  /  EST. 2026',
            'あなたの仕事も、SNSも、\nひとつの流れに。',
            'In the core of every era, there is a CORE.',
            'AIエージェントで、一気通貫に動かす。', 'core')

    # 2 problem keyword
    s_keyword('ひとりで、\n全部。', '経営も、Instagramも、LINEも、リンクも。')

    # 3 quote
    s_keyword('道具は増えた。\nでも、つながっていない。', 'アプリの数だけ、管理が増えていく。', size=80)

    # 4 section 01
    s_section(1, 'FOUR SPECIALISTS', '四つの専門。')

    # 5-8 product intros (mark + name + role + price)
    def s_product_intro(key, role, oneline, price):
        global PAGE
        PAGE += 1; s = blank(); bg(s)
        pal = {k: (hx(v) if k in ('primary', 'bright', 'deep') else v) for k, v in PAL[key].items()}
        pal['grad_rgb'] = [hx(c) for c in PAL[key]['grad']]
        global P
        keep = P; P = pal
        glow(s, 3.2, 3.7, 2.6, pal['primary'], 0.12)
        mark(s, PAL[key]['mark'], 3.0, 3.4, 1.7)
        txt(s, PAL[key]['en'], 6.0, 1.5, 6.5, 0.4, font=F_NUM, size=12, color=pal['bright'], bold=True, spc=500, align='left')
        txt(s, role, 6.0, 1.95, 6.7, 1.4, font=F_DISP, size=40, color=INK, bold=True, align='left', leading=1.15)
        txt(s, oneline, 6.0, 3.75, 6.7, 1.6, font=F_SANS, size=17, color=INK_DIM, align='left', anchor='top', leading=1.7)
        txt(s, price, 6.0, 5.55, 6.7, 0.5, font=F_NUM, size=15, color=pal['bright'], bold=True, align='left')
        P = keep
        footer(s)

    s_product_intro('prism', 'Prism — 全事業の司令塔',
                    '営業・財務・契約・議事録。経営のすべてを13名のAIエージェントが引き受ける、あなたの頭脳。', '月 ¥4,800〜 ・ 7日間無料')
    s_product_intro('iris', 'Iris — Instagram',
                    '投稿・分析・案件・DM返信・交渉。Instagram運用の六つの仕事を、ひとつのアプリに。', '月 ¥2,800〜 ・ 7日間無料')
    s_product_intro('resonance', 'Resonance — LINE',
                    '名簿の一人ひとりにAIが文面を書き分け、LINEで手紙のように届ける個別配信。', '月 ¥980〜 ・ 7日間無料')
    s_product_intro('lume', 'Lume — リンク',
                    'すべてのリンクをひとつに束ね、誰がどこを踏んだかを色と熱で可視化するハブ。', '月 ¥980〜 ・ 7日間無料')

    # 9 section 02
    s_section(2, 'ONE FLOW', 'つながり。')

    # 10 keyword
    s_keyword('バラバラの道具では、\nない。', 'ひとつの核で、つながっている。', size=78)

    # 11 connection diagram
    PAGE += 1; s = blank(); bg(s)
    glow(s, SW / 2, 3.5, 3.4, hx('#a78bfa'), 0.10)
    txt(s, 'Prism を中心に、三つのSNSがつながる。', 0.5, 0.7, SW - 1, 0.7, font=F_DISP, size=26, color=INK, bold=True)
    cx, cy = SW / 2, 4.1
    sats = [('iris', 'Instagram', cx, 2.25, '#E1306C'),
            ('resonance', 'LINE', cx - 3.6, 5.5, '#06C755'),
            ('lume', 'リンク', cx + 3.6, 5.5, '#FFA42A')]
    for key, role, sx, sy, col in sats:
        line(s, cx, cy, sx, sy, hx(col), 1.6, dash='sysDot')
    # center prism
    pkeep = P
    P = {k: (hx(v) if k in ('primary', 'bright', 'deep') else v) for k, v in PAL['prism'].items()}
    P['grad_rgb'] = [hx(c) for c in PAL['prism']['grad']]
    circle(s, cx, cy, 1.05, fill=hx('#0e0e1a'), line=hx('#a78bfa'), lw=1.5)
    mark(s, 'prism', cx, cy - 0.12, 0.95)
    txt(s, '司令塔 Prism', cx - 1.3, cy + 0.62, 2.6, 0.4, font=F_DISP, size=14, color=INK, bold=True)
    for key, role, sx, sy, col in sats:
        P = {k: (hx(v) if k in ('primary', 'bright', 'deep') else v) for k, v in PAL[key].items()}
        P['grad_rgb'] = [hx(c) for c in PAL[key]['grad']]
        circle(s, sx, sy, 0.78, fill=hx('#0b0b14'), line=hx(col), lw=1.2)
        mark(s, PAL[key]['mark'], sx, sy - 0.1, 0.62)
        txt(s, role, sx - 1.0, sy + 0.5, 2.0, 0.35, font=F_SANS, size=12, color=hx(col), bold=True)
    P = pkeep
    footer(s)

    # 12 flow scenario (4 steps using horizontal cards, core accent)
    s_steps('A DAY, CONNECTED', '横串にすると、こんなことが起きる。', [
        ('01', 'Lume', 'ファンが、どのリンクを踏んだか分かる'),
        ('02', 'Iris', 'その人のInstagramの反応を、AIが解析'),
        ('03', 'Resonance', 'いま響く一文を、LINEでその人だけに'),
    ])
    # 13 step 4 + payoff
    s_keyword('あなたは、\n確認して送るだけ。', 'Prism がすべてを記録し、13名のAI役員が次の一手を出す。', size=72)

    # 14 stat 13 officers
    s_stat('13', 'AI OFFICERS', '経営・営業・財務・創造・データ・人材・法務。\n13名の専門エージェントが、すべてのプロダクトに乗っています。')

    # 15 section 03
    s_section(3, 'TRUST & PRICING', '信頼と、料金。')

    # 16 pricing overview (4 products)
    PAGE += 1; s = blank(); bg(s)
    txt(s, '四つとも、7日間無料・クレカ不要。', 0.9, 0.95, SW - 1.8, 1.0, font=F_DISP, size=34, color=INK, bold=True, align='left')
    rows = [('Prism', '全事業の司令塔', '月 ¥4,800〜', '#a78bfa'),
            ('Iris', 'Instagram運用', '月 ¥2,800〜', '#E1306C'),
            ('Resonance', 'LINE個別配信', '月 ¥980〜', '#06C755'),
            ('Lume', 'リンクハブ＋解析', '月 ¥980〜', '#FFA42A')]
    y0 = 2.5
    for i, (nm, role, pr, col) in enumerate(rows):
        y = y0 + i * 0.95
        rect(s, 0.9, y, SW - 1.8, 0.82, fill=hx('#0c0c16'), line=hx(col), lw=0.8, rounded=True)
        circle(s, 1.45, y + 0.41, 0.13, fill=hx(col))
        txt(s, nm, 1.85, y + 0.12, 3.0, 0.55, font=F_DISP, size=19, color=INK, bold=True, align='left')
        txt(s, role, 4.8, y + 0.16, 4.5, 0.5, font=F_SANS, size=14, color=INK_DIM, align='left')
        txt(s, pr, SW - 4.0, y + 0.12, 3.1, 0.55, font=F_NUM, size=18, color=hx(col), bold=True, align='right')
    footer(s)

    # 17 creed (3 promises)
    P = {k: (hx(v) if k in ('primary', 'bright', 'deep') else v) for k, v in PAL['core'].items()}
    P['grad_rgb'] = [hx(c) for c in PAL['core']['grad']]
    s_cards4('OUR CREED', 'わたしたちが守る、三つの約束。', [
        ('偽りの数字は、載せない', 'まだ実績がないものには「—」と記す'),
        ('やさしい言葉で、語る', '専門用語は、できるかぎり日常の言葉に'),
        ('使ったぶんだけ、いただく', '気づかぬうちに高額にならない料金設計'),
        ('人を、ど真ん中に', '技術は手段。いつも中心にあるのは人'),
    ])

    # 18 mission keyword
    s_keyword('すべての時代の、\n核となるものを。', 'In the core of every era, there is a CORE.', size=64)

    # 19 company
    PAGE += 1; s = blank(); bg(s)
    txt(s, 'ABOUT', 0.9, 1.0, 5, 0.4, font=F_NUM, size=12, color=P['bright'], bold=True, spc=600, align='left')
    txt(s, '株式会社コア', 0.9, 1.5, SW - 1.8, 1.0, font=F_DISP, size=40, color=INK, bold=True, align='left')
    info = [('会社名', '株式会社コア / CORE Inc.'),
            ('代表取締役', '井出 直毅 (Naoki Ide)'),
            ('事業内容', 'エージェントAIを中心とした SaaS の開発・運営'),
            ('提供サービス', 'Prism（事業家）／ Iris（Instagram）／ Resonance（LINE）／ Lume（リンク）'),
            ('所在地', '兵庫県神戸市東灘区魚崎南町7丁目11-7')]
    for i, (k, v) in enumerate(info):
        y = 3.0 + i * 0.7
        txt(s, k, 0.9, y, 2.6, 0.5, font=F_SANS, size=13, color=P['bright'], bold=True, align='left')
        txt(s, v, 3.7, y, SW - 4.6, 0.5, font=F_SANS, size=14.5, color=INK, align='left')
    footer(s)

    # 20 CTA
    s_cta('核を、共に。', 'core-prism-app.vercel.app    ·    hello@core-inc.jp', 'core')

    out = f"{DESKTOP}/CORE_統合_営業資料_{DATE}.pptx"
    prs.save(out)
    print('saved', out, f'({PAGE}p)')


build_combined()
print('ALL DONE')
