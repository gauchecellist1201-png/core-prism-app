#!/usr/bin/env python3
# ============================================================
# CORE Identity OS — 60 分キーノート プレゼン資料生成
# 出力: ~/Desktop/CORE_キーノート_2026-05-11.pptx
# 約 130-140 枚、ダークモード、グラデアクセント
# ============================================================
import os
import qrcode
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── 色定義 ───────────────────────────
BG = RGBColor(0x0A, 0x0A, 0x14)
INK = RGBColor(0xFF, 0xFA, 0xF5)
INK_DIM = RGBColor(0x9A, 0x90, 0xAC)
INK_FAINT = RGBColor(0x55, 0x4D, 0x68)
GOLD = RGBColor(0xFC, 0xB0, 0x45)
PINK = RGBColor(0xE1, 0x30, 0x6C)
PURPLE = RGBColor(0xB0, 0x7B, 0xD9)
DEEP_PURPLE = RGBColor(0x5B, 0x2C, 0x8A)
PRISM_BLUE = RGBColor(0x6F, 0xA8, 0xFF)
LILAC = RGBColor(0x7A, 0x8A, 0xFF)
WARN_RED = RGBColor(0xFF, 0x5C, 0x5C)

FONT_DISPLAY = 'Noto Serif JP'
FONT_SERIF = 'Noto Serif JP'
FONT_SANS = 'Noto Sans JP'
FONT_NUM = 'Inter'

# ── スライドサイズ (16:9) ───────────────────────
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# ── プレゼン作成 ───────────────────────
prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)

BLANK_LAYOUT = prs.slide_layouts[6]  # blank


def add_bg(slide, color=BG):
    """背景色をフィル"""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, text, left, top, width, height,
             font=FONT_SANS, size=14, color=INK, bold=False, italic=False,
             align='center', anchor='middle', letter_spacing=None):
    """テキストボックスを追加"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    if anchor == 'top':
        tf.vertical_anchor = MSO_ANCHOR.TOP
    elif anchor == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    else:
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if align == 'left':
        p.alignment = PP_ALIGN.LEFT
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.CENTER
    # 改行対応
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            run = p.add_run()
            run.text = line
        else:
            p2 = tf.add_paragraph()
            p2.alignment = p.alignment
            run = p2.add_run()
            run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if letter_spacing is not None:
            # XML で文字間隔調整
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(letter_spacing))
    return tb


def add_circle(slide, cx, cy, r, fill_color=None, line_color=None, line_w=None, alpha=None):
    """円。位置は中心 (Emu)、r も Emu"""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2
    )
    if fill_color:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_color
        if alpha is not None:
            # transparency XML
            sp = sh.fill._xPr
            solidFill = sp.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    a = srgb.find(qn('a:alpha'))
                    if a is None:
                        a = etree.SubElement(srgb, qn('a:alpha'))
                    a.set('val', str(int(alpha * 100000)))
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = Pt(line_w)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, alpha=None, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill_color:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill_color
        if alpha is not None:
            sp = sh.fill._xPr
            solidFill = sp.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    a = srgb.find(qn('a:alpha'))
                    if a is None:
                        a = etree.SubElement(srgb, qn('a:alpha'))
                    a.set('val', str(int(alpha * 100000)))
    else:
        sh.fill.background()
    if line_color is not None:
        sh.line.color.rgb = line_color
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def add_gradient_bar(slide, x, y, w, h):
    """ゴールド→ピンク→パープルの 3 段グラデバー (3 矩形で擬似)"""
    # 3 分割
    seg_w = w // 3
    add_rect(slide, x, y, seg_w, h, fill_color=GOLD)
    add_rect(slide, x + seg_w, y, seg_w, h, fill_color=PINK)
    add_rect(slide, x + 2 * seg_w, y, w - 2 * seg_w, h, fill_color=PURPLE)


def add_footer(slide, page_num, total):
    """全スライド共通のフッター"""
    # 左下: 製作者
    add_text(slide, 'Core Identity OS  ·  株式会社コアプリズム (仮)',
             Inches(0.5), Inches(SLIDE_H_IN - 0.45), Inches(5), Inches(0.3),
             font=FONT_SANS, size=8, color=INK_FAINT, align='left')
    # 右下: ページ番号
    add_text(slide, f'{page_num:03d} / {total}',
             Inches(SLIDE_W_IN - 1.5), Inches(SLIDE_H_IN - 0.45), Inches(1.0), Inches(0.3),
             font=FONT_NUM, size=9, color=INK_FAINT, align='right')


# ── スライド種別ヘルパー ───────────────────────

def slide_cover(title_main, title_sub, presenter, eyebrow, page, total):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    # 中心のオーロラ風グロー
    add_circle(s, Emu(SLIDE_W_IN * 914400 / 2), Emu(SLIDE_H_IN * 914400 / 2),
               Inches(3.5), fill_color=PURPLE, alpha=0.15)
    add_circle(s, Inches(4.0), Inches(2.5),
               Inches(2.0), fill_color=GOLD, alpha=0.10)
    add_circle(s, Inches(9.5), Inches(5.5),
               Inches(2.2), fill_color=PINK, alpha=0.10)
    # eyebrow
    add_text(s, eyebrow, Inches(0.5), Inches(1.5), Inches(SLIDE_W_IN - 1), Inches(0.4),
             font=FONT_SANS, size=12, color=GOLD, bold=True, letter_spacing=600)
    # main title
    add_text(s, title_main, Inches(0.5), Inches(2.4), Inches(SLIDE_W_IN - 1), Inches(2.2),
             font=FONT_DISPLAY, size=64, color=INK, bold=True)
    # subtitle
    add_text(s, title_sub, Inches(0.5), Inches(4.6), Inches(SLIDE_W_IN - 1), Inches(0.8),
             font=FONT_SERIF, size=22, color=INK_DIM, italic=True)
    # gradient bar
    add_gradient_bar(s, Inches(5.66), Inches(5.7), Inches(2.0), Inches(0.04))
    # presenter
    add_text(s, presenter, Inches(0.5), Inches(5.95), Inches(SLIDE_W_IN - 1), Inches(0.5),
             font=FONT_SANS, size=14, color=INK_DIM)
    add_footer(s, page, total)
    return s


def slide_section(num, jp_title, en_title, page, total):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    # 巨大な番号
    add_text(s, f'{num:02d}', Inches(0.5), Inches(0.8), Inches(SLIDE_W_IN - 1), Inches(2.5),
             font=FONT_NUM, size=200, color=INK_FAINT, bold=True)
    # 英タイトル (eyebrow)
    add_text(s, en_title, Inches(0.5), Inches(3.4), Inches(SLIDE_W_IN - 1), Inches(0.5),
             font=FONT_SANS, size=14, color=GOLD, bold=True, letter_spacing=700)
    # 日本語タイトル
    add_text(s, jp_title, Inches(0.5), Inches(4.0), Inches(SLIDE_W_IN - 1), Inches(1.6),
             font=FONT_DISPLAY, size=72, color=INK, bold=True)
    # gradient bar
    add_gradient_bar(s, Inches(SLIDE_W_IN/2 - 1.0), Inches(6.0), Inches(2.0), Inches(0.05))
    add_footer(s, page, total)
    return s


def slide_keyword(big, small, page, total, big_color=INK, big_size=180, small_color=INK_DIM):
    """中央に巨大キーワード + 下に補足"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, big, Inches(0.5), Inches(2.0), Inches(SLIDE_W_IN - 1), Inches(3.0),
             font=FONT_DISPLAY, size=big_size, color=big_color, bold=True)
    if small:
        add_text(s, small, Inches(1.0), Inches(5.4), Inches(SLIDE_W_IN - 2), Inches(1.0),
                 font=FONT_SANS, size=20, color=small_color)
    add_footer(s, page, total)
    return s


def slide_quote(quote, attribution, page, total):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    # 開きクォート
    add_text(s, '"', Inches(0.5), Inches(0.8), Inches(SLIDE_W_IN - 1), Inches(1.5),
             font='Garamond', size=160, color=GOLD, bold=True)
    add_text(s, quote, Inches(1.0), Inches(2.4), Inches(SLIDE_W_IN - 2), Inches(3.5),
             font=FONT_DISPLAY, size=44, color=INK, italic=True)
    if attribution:
        add_text(s, '— ' + attribution, Inches(1.0), Inches(6.0), Inches(SLIDE_W_IN - 2), Inches(0.5),
                 font=FONT_SANS, size=14, color=INK_DIM)
    add_footer(s, page, total)
    return s


def slide_title_body(title, body, page, total, eyebrow=None, accent_color=GOLD):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    if eyebrow:
        add_text(s, eyebrow, Inches(0.7), Inches(0.7), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
                 font=FONT_SANS, size=11, color=accent_color, bold=True, letter_spacing=600, align='left')
    add_text(s, title, Inches(0.7), Inches(1.3), Inches(SLIDE_W_IN - 1.4), Inches(1.5),
             font=FONT_DISPLAY, size=46, color=INK, bold=True, align='left')
    add_text(s, body, Inches(0.7), Inches(3.3), Inches(SLIDE_W_IN - 1.4), Inches(3.5),
             font=FONT_SANS, size=20, color=INK_DIM, align='left', anchor='top')
    add_footer(s, page, total)
    return s


def slide_two_col(title, left_label, left_body, right_label, right_body, page, total, eyebrow=None,
                  left_color=PRISM_BLUE, right_color=PINK):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    if eyebrow:
        add_text(s, eyebrow, Inches(0.7), Inches(0.7), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
                 font=FONT_SANS, size=11, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, title, Inches(0.7), Inches(1.2), Inches(SLIDE_W_IN - 1.4), Inches(1.0),
             font=FONT_DISPLAY, size=36, color=INK, bold=True, align='left')
    # 左カラム
    col_w = Inches(5.5)
    col_h = Inches(4.2)
    col_y = Inches(2.8)
    add_rect(s, Inches(0.7), col_y, col_w, col_h,
             fill_color=left_color, alpha=0.08, rounded=True)
    add_text(s, left_label, Inches(1.0), col_y + Inches(0.3), col_w, Inches(0.4),
             font=FONT_SANS, size=11, color=left_color, bold=True, letter_spacing=500, align='left', anchor='top')
    add_text(s, left_body, Inches(1.0), col_y + Inches(1.0), col_w - Inches(0.6), col_h - Inches(1.3),
             font=FONT_SANS, size=18, color=INK, align='left', anchor='top')
    # 右カラム
    add_rect(s, Inches(7.1), col_y, col_w, col_h,
             fill_color=right_color, alpha=0.08, rounded=True)
    add_text(s, right_label, Inches(7.4), col_y + Inches(0.3), col_w, Inches(0.4),
             font=FONT_SANS, size=11, color=right_color, bold=True, letter_spacing=500, align='left', anchor='top')
    add_text(s, right_body, Inches(7.4), col_y + Inches(1.0), col_w - Inches(0.6), col_h - Inches(1.3),
             font=FONT_SANS, size=18, color=INK, align='left', anchor='top')
    add_footer(s, page, total)
    return s


def slide_4cards(title, cards, page, total, eyebrow=None):
    """2x2 のカード"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    if eyebrow:
        add_text(s, eyebrow, Inches(0.7), Inches(0.6), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
                 font=FONT_SANS, size=11, color=GOLD, bold=True, letter_spacing=600, align='left')
    add_text(s, title, Inches(0.7), Inches(1.1), Inches(SLIDE_W_IN - 1.4), Inches(1.0),
             font=FONT_DISPLAY, size=36, color=INK, bold=True, align='left')
    # 4 cards
    card_w = Inches(5.7)
    card_h = Inches(2.2)
    gap = Inches(0.3)
    start_x = Inches(0.7)
    start_y = Inches(2.7)
    accents = [GOLD, PINK, PURPLE, PRISM_BLUE]
    for i, (title_c, body_c) in enumerate(cards[:4]):
        row = i // 2
        col = i % 2
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        add_rect(s, x, y, card_w, card_h, fill_color=accents[i], alpha=0.10, rounded=True)
        # アイコン円
        add_circle(s, x + Inches(0.55), y + Inches(0.55), Inches(0.25),
                   fill_color=accents[i])
        add_text(s, title_c, x + Inches(1.0), y + Inches(0.25), card_w - Inches(1.2), Inches(0.6),
                 font=FONT_DISPLAY, size=20, color=INK, bold=True, align='left', anchor='middle')
        add_text(s, body_c, x + Inches(0.5), y + Inches(1.0), card_w - Inches(1.0), card_h - Inches(1.2),
                 font=FONT_SANS, size=13, color=INK_DIM, align='left', anchor='top')
    add_footer(s, page, total)
    return s


def slide_stat(big, label, sub, page, total, big_color=GOLD, big_font=FONT_NUM):
    """巨大な数字と短い説明"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, label, Inches(0.5), Inches(1.7), Inches(SLIDE_W_IN - 1), Inches(0.5),
             font=FONT_SANS, size=14, color=INK_DIM, bold=True, letter_spacing=600)
    add_text(s, big, Inches(0.5), Inches(2.4), Inches(SLIDE_W_IN - 1), Inches(3.0),
             font=big_font, size=260, color=big_color, bold=True)
    if sub:
        add_text(s, sub, Inches(1.0), Inches(5.7), Inches(SLIDE_W_IN - 2), Inches(1.0),
                 font=FONT_SANS, size=22, color=INK)
    add_footer(s, page, total)
    return s


def slide_full_quote_dark(big_text, sub, page, total, big_color=INK, big_size=66, italic=True):
    """大きな引用 (映画的)"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_circle(s, Inches(SLIDE_W_IN/2), Inches(SLIDE_H_IN/2),
               Inches(4.0), fill_color=PURPLE, alpha=0.10)
    add_text(s, big_text, Inches(0.7), Inches(2.0), Inches(SLIDE_W_IN - 1.4), Inches(3.5),
             font=FONT_DISPLAY, size=big_size, color=big_color, italic=italic, bold=True)
    if sub:
        add_text(s, sub, Inches(1.0), Inches(5.7), Inches(SLIDE_W_IN - 2), Inches(1.0),
                 font=FONT_SANS, size=18, color=INK_DIM)
    add_footer(s, page, total)
    return s


def slide_list(title, items, page, total, eyebrow=None, accent_color=GOLD):
    """タイトル + 縦リスト"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    if eyebrow:
        add_text(s, eyebrow, Inches(0.7), Inches(0.7), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
                 font=FONT_SANS, size=11, color=accent_color, bold=True, letter_spacing=600, align='left')
    add_text(s, title, Inches(0.7), Inches(1.2), Inches(SLIDE_W_IN - 1.4), Inches(1.2),
             font=FONT_DISPLAY, size=40, color=INK, bold=True, align='left')
    # 縦リスト
    item_h = Inches(0.7)
    start_y = Inches(3.0)
    for i, it in enumerate(items):
        y = start_y + i * (item_h + Inches(0.1))
        # 番号円
        add_circle(s, Inches(1.0), y + Inches(0.35), Inches(0.18),
                   fill_color=accent_color)
        add_text(s, str(i+1), Inches(0.7), y + Inches(0.05), Inches(0.6), Inches(0.6),
                 font=FONT_NUM, size=14, color=BG, bold=True)
        add_text(s, it, Inches(1.5), y, Inches(SLIDE_W_IN - 2.2), item_h,
                 font=FONT_SANS, size=20, color=INK, align='left')
    add_footer(s, page, total)
    return s


def slide_seven_agents(page, total):
    """7 つのエージェント (Prism) を虹色で表示"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'SEVEN AGENTS, ONE OS', Inches(0.5), Inches(0.6), Inches(SLIDE_W_IN - 1), Inches(0.4),
             font=FONT_SANS, size=12, color=GOLD, bold=True, letter_spacing=700)
    add_text(s, '7 つのあなたに、7 つのエージェント。', Inches(0.5), Inches(1.1), Inches(SLIDE_W_IN - 1), Inches(1.0),
             font=FONT_DISPLAY, size=36, color=INK, bold=True)
    agents = [
        ('🧭', '経営', 'CEO Agent', '#FF5757'),
        ('💼', '営業', 'Sales Agent', '#FF9842'),
        ('📊', '財務', 'CFO Agent', '#FBBF24'),
        ('✨', '創造', 'Creative Agent', '#4ADE80'),
        ('📚', '学び', 'Knowledge Agent', '#60A5FA'),
        ('👥', '人材', 'People Agent', '#A78BFA'),
        ('💗', '生活', 'Life Agent', '#F472B6'),
    ]
    # 7 つを横並び
    card_w = Inches(1.62)
    card_h = Inches(3.5)
    gap = Inches(0.05)
    total_w = card_w * 7 + gap * 6
    start_x = Emu((prs.slide_width - total_w) / 2)
    start_y = Inches(2.8)
    for i, (icon, name, role, hex_c) in enumerate(agents):
        x = start_x + i * (card_w + gap)
        c = RGBColor.from_string(hex_c[1:])
        add_rect(s, x, start_y, card_w, card_h, fill_color=c, alpha=0.12, rounded=True)
        add_text(s, icon, x, start_y + Inches(0.4), card_w, Inches(0.8),
                 font=FONT_SANS, size=36, color=INK)
        add_text(s, name, x, start_y + Inches(1.4), card_w, Inches(0.6),
                 font=FONT_DISPLAY, size=22, color=INK, bold=True)
        add_text(s, role, x, start_y + Inches(2.1), card_w, Inches(0.5),
                 font=FONT_SANS, size=10, color=c, bold=True, letter_spacing=400)
    add_footer(s, page, total)
    return s


def slide_six_facets(page, total):
    """Iris の 6 ファセット"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    add_text(s, 'SIX FACETS OF LIGHT', Inches(0.5), Inches(0.6), Inches(SLIDE_W_IN - 1), Inches(0.4),
             font=FONT_SANS, size=12, color=GOLD, bold=True, letter_spacing=700)
    add_text(s, '光は、6 つの色を持つ。', Inches(0.5), Inches(1.1), Inches(SLIDE_W_IN - 1), Inches(1.0),
             font=FONT_DISPLAY, size=36, color=INK, bold=True, italic=True)
    facets = [
        ('✦', '案件', '受注 → 下書き → 投稿 → レポートまで AI が一気通貫'),
        ('◐', '分析', 'Instagram 解析: 投稿時間・反応率・伸びるテーマ'),
        ('✶', '創作', 'キャプション・サムネ・OG 画像を即生成'),
        ('◇', '交渉', '料金交渉ロープレ・媒体資料・ブランド提案文'),
        ('✣', 'ブランド', 'フォント・カラー・トーンをパーソナル AI が提案'),
        ('❋', '仲間', '招待制クリエイターコミュニティ + コラボ募集'),
    ]
    card_w = Inches(3.9)
    card_h = Inches(2.0)
    gap = Inches(0.25)
    start_x = Inches(0.7)
    start_y = Inches(2.5)
    accents = [PINK, PURPLE, GOLD, RGBColor(0xFF, 0xA9, 0x4D), RGBColor(0xC7, 0x7D, 0xFF), RGBColor(0xFD, 0x7C, 0xB8)]
    for i, (icon, name, desc) in enumerate(facets):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        c = accents[i]
        add_rect(s, x, y, card_w, card_h, fill_color=c, alpha=0.10, rounded=True)
        add_text(s, icon, x + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7),
                 font=FONT_DISPLAY, size=30, color=c, bold=True, align='left')
        add_text(s, name, x + Inches(1.0), y + Inches(0.25), card_w - Inches(1.2), Inches(0.6),
                 font=FONT_DISPLAY, size=22, color=INK, bold=True, italic=True, align='left')
        add_text(s, desc, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.6), Inches(0.9),
                 font=FONT_SANS, size=12, color=INK_DIM, align='left', anchor='top')
    add_footer(s, page, total)
    return s


def slide_qr(page, total):
    """巨大 QR コード"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s)
    # QR 生成
    url = 'https://core-prism-app.vercel.app/keynote'
    qr = qrcode.QRCode(version=4, error_correction=qrcode.constants.ERROR_CORRECT_M,
                       box_size=20, border=2)
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color='white', back_color='#0a0a14')
    buf = BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    # 右半分に QR
    qr_size = Inches(4.5)
    qr_x = Inches(SLIDE_W_IN - 5.5)
    qr_y = Inches((SLIDE_H_IN - 4.5) / 2)
    s.shapes.add_picture(buf, qr_x, qr_y, qr_size, qr_size)
    # 左半分にテキスト
    add_text(s, 'CLOSED BETA INVITATION', Inches(0.7), Inches(1.5), Inches(6.5), Inches(0.4),
             font=FONT_SANS, size=12, color=GOLD, bold=True, letter_spacing=700, align='left')
    add_text(s, 'これが、\n片道切符です。', Inches(0.7), Inches(2.0), Inches(6.5), Inches(2.5),
             font=FONT_DISPLAY, size=54, color=INK, bold=True, align='left')
    add_text(s, '右の QR を、今、読み取ってください。', Inches(0.7), Inches(4.7), Inches(6.5), Inches(0.6),
             font=FONT_SANS, size=18, color=INK_DIM, align='left')
    add_text(s, 'core-prism-app.vercel.app/keynote', Inches(0.7), Inches(5.4), Inches(6.5), Inches(0.5),
             font=FONT_NUM, size=14, color=PURPLE, align='left')
    add_text(s, '講演会限定コード:  KEYNOTE30  →  トライアル 30 日延長', Inches(0.7), Inches(6.0), Inches(6.5), Inches(0.5),
             font=FONT_SANS, size=14, color=GOLD, bold=True, align='left')
    add_footer(s, page, total)
    return s


def slide_silence(page, total):
    """完全な静寂を表現する黒に近いスライド (沈黙の 15 秒用)"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_bg(s, color=RGBColor(0, 0, 0))
    add_text(s, '·', Inches(0.5), Inches(3.5), Inches(SLIDE_W_IN - 1), Inches(0.5),
             font=FONT_SERIF, size=60, color=INK_FAINT)
    add_footer(s, page, total)
    return s


# ========================================================================
# === ここから本番のスライド構成 (139 枚) ================================
# ========================================================================

# 各スライドを順次追加する前に、合計枚数を仮に置く (後で更新)
TOTAL = 139

slides_def = []  # (callable, args) のリスト
def add(fn, *args, **kwargs):
    slides_def.append((fn, args, kwargs))

# --- COVER (1) ---
add(slide_cover,
    title_main='人類史から読み解く\n自己の拡張',
    title_sub='— 脳科学、アート、AI が交差する「分身」の創造',
    presenter='井出直毅 / Naoki Ide  ·  2026.05.12',
    eyebrow='CORE IDENTITY OS — KEYNOTE')

# --- SECTION 1: 静寂と異常な履歴書 (10) ---
add(slide_section, 1, '静寂と異常な履歴書', 'SILENCE & THE ABNORMAL RESUME')
add(slide_silence)
add(slide_keyword, '15 秒', '皆さんの頭の中には、何が浮かんでいましたか?',
    big_color=INK, big_size=240)
add(slide_full_quote_dark,
    '現代を生きる私たちは、\nわずかな空白の時間すら\n耐えられなくなっている。',
    '水中で息を止めるようにして、ギリギリで日々を生きている。', big_size=44)
add(slide_title_body, '本日のスピーカー',
    '井出直毅 (Naoki Ide)\n\n本日はそんな皆さんの「息苦しさ」を完全に終わらせるために来ました。\n本日はよろしくお願いいたします。',
    eyebrow='SPEAKER')
add(slide_4cards, '異常な履歴書', [
    ('🏢  グランドリアルエステート COO', '不動産の数億のディールを動かす'),
    ('🎻  プロチェリスト GAUCHE', 'チェロミュージックスクールを経営'),
    ('🦷  大阪大学 歯学部 在籍', '医療の最前線で人を救う'),
    ('💻  AI エンジニア', 'コードで未来の常識を作る'),
], eyebrow='THE FOUR DOMAINS')
add(slide_quote, '「なんだこいつ、絶対に嘘ついてるだろ」',
    '会場の皆さんが、まず思うこと。')
add(slide_full_quote_dark,
    '僕は天才でもなければ、\nショートスリーパーの超人でもない。\n何なら毎日 12 時間は寝る。',
    '異常なほど「好奇心を抑えられない」、ちょっと厄介で不器用な人間なだけ。', big_size=42)
add(slide_keyword, '全部、やりたかった。', '一つに絞りたくなかった。',
    big_color=INK, big_size=110)
add(slide_full_quote_dark,
    '今日お話しするのは、\n血塗れになりながら見つけた\n自己拡張という究極のサバイバル術。',
    'これは僕の成功体験の自慢話ではありません。', big_size=40)

# --- SECTION 2: 時間の奴隷 (12) ---
add(slide_section, 2, '時間の奴隷', 'SLAVES OF TIME')
add(slide_keyword, '数十万年', 'のスケールで、視野を広げてみましょう。',
    big_color=GOLD, big_size=200)
add(slide_full_quote_dark,
    'ホモ・サピエンスの歴史の 99%、\n私たちには「時間」という概念が\nそのものがありませんでした。',
    '太陽が昇れば目を覚まし、お腹が空けばマンモスを追いかけ、日が沈めば焚き火を囲んで眠る。', big_size=40)
add(slide_4cards, '狩猟採集時代の生活', [
    ('☀️  太陽', '昇れば目を覚ます。それが朝の合図。'),
    ('🦣  マンモス', 'お腹が空けば追いかける。'),
    ('🔥  焚き火', '日が沈めば囲んで眠る。'),
    ('🌙  自然のサイクル', 'と完全に同化していた。'),
], eyebrow='LIFE BEFORE THE CLOCK')
add(slide_keyword, '13 — 14 世紀', '機械式時計の発明',
    big_color=PURPLE, big_size=170)
add(slide_title_body, 'ベネディクト派の修道士たち',
    'キリスト教の修道士たちは、1 日に 7 回、神に対して正確な時間に祈りを捧げる必要があった。\n\n太陽が出ていようが曇っていようが、絶対に決まった時間に祈らなければならない。\n\nそのために彼らは、歯車を組み合わせ、一定の間隔で鐘を鳴らす機械を作った。',
    eyebrow='THE INVENTORS')
add(slide_full_quote_dark,
    '神と繋がるために作った機械が、\n人間を「分」や「秒」という\n見えない概念の奴隷にした。',
    '歴史の最も皮肉な瞬間。', big_size=42)
add(slide_keyword, '18 世紀', '産業革命',
    big_color=PINK, big_size=180)
add(slide_full_quote_dark,
    '人間は、\n「労働者」という名前の\n一つの「部品」になった。',
    '工場のサイレンが鳴れば一斉に働き、サイレンが鳴れば帰る。', big_size=44)
add(slide_keyword, '肉の歯車', '時計に合わせて動く、人間の姿。',
    big_color=WARN_RED, big_size=140)
add(slide_full_quote_dark,
    '「選択と集中」という現代の美徳は、\n産業革命時代の工場長にとって\n都合が良かっただけのシステムの名残。',
    'お前は今日から、この機械のネジを締める職人だ。', big_size=38)

# --- SECTION 3: ルネサンスと現代のバグ (9) ---
add(slide_section, 3, 'ルネサンスと現代のバグ', 'RENAISSANCE & THE MODERN BUG')
add(slide_keyword, 'Leonardo da Vinci', 'ルネサンスを思い出してください。',
    big_color=GOLD, big_size=80)
add(slide_4cards, '彼は、5 つの顔を持っていた。', [
    ('🎨  画家', 'モナ・リザ。最後の晩餐。'),
    ('🗿  彫刻家 / 建築家', '空間そのものを設計する。'),
    ('🩺  解剖学者', '人体の内側を、絵で記録する。'),
    ('⚙️  軍事技術エンジニア', '戦車・空気銃・潜水服。'),
], eyebrow='DA VINCI HAS FIVE FACES')
add(slide_full_quote_dark,
    '誰も彼に、\n「絵を描くか解剖するか、\nどっちかに絞れ」とは言わなかった。',
    '人間は本来、複数の顔を持つ生き物。', big_size=44)
add(slide_keyword, 'マルチ・アイデンティティ', 'それが、人間の自然な姿。',
    big_color=PURPLE, big_size=80)
add(slide_4cards, '今、私たちもダ・ヴィンチに戻れる。', [
    ('💼  経営者', '戦略を練る。'),
    ('📱  発信者', '言葉を紡ぐ。'),
    ('🎬  クリエイター', '動画を作る。'),
    ('👨‍👩‍👧  親', '家族を抱きしめる。'),
], eyebrow='THE SECOND RENAISSANCE')
add(slide_keyword, 'BUG.', 'しかし、決定的な不具合がある。',
    big_color=WARN_RED, big_size=240)
add(slide_full_quote_dark,
    'やりたいことはルネサンス時代に戻り、\nツールは現代の最高峰を手に入れた。\nしかし、脳という生物学的ハードウェアは、\n狩猟採集時代から全くアップデートされていない。',
    'これが、皆さんが日々感じている「時間が足りない」の正体。', big_size=34)
add(slide_keyword, '一つしかない。', '自分という人間の肉体と脳が。',
    big_color=INK, big_size=110)

# --- SECTION 4: 脳科学 (12) ---
add(slide_section, 4, '脳科学が暴く隙間時間の罠', 'NEUROSCIENCE OF THE INTERSTICE')
add(slide_quote,
    '「できる人は隙間時間を活用している」\n「移動の 5 分でメールをさばけ」',
    'ビジネス書のコーナーで、目にする言葉。')
add(slide_keyword, '合法的な\n脳の破壊行為。', '今日から絶対に、やめてください。',
    big_color=WARN_RED, big_size=110)
add(slide_full_quote_dark,
    '視床下部腹内側核・弓状核尾側部における\n逆行性トレーサーを用いた神経回路の\nトレースとその機能的同定に関する研究',
    '— 大学時代、僕が血眼になって取り組んでいたテーマ。', big_size=24, italic=False)
add(slide_quote,
    '「あんた、日本語を喋りなさい」',
    '研究テーマを母に見せた時の反応。')
add(slide_full_quote_dark,
    '脳の配線図を逆から辿る。\n情報の「源 (ソース)」を突き止める。',
    'この技術が、自分の分身を AI で作るすべての基礎になっている。', big_size=46)
add(slide_title_body, '前頭前野とフロー状態',
    '皆さんが新規事業のアイデアを練っている時、脳は「前頭前野」に深く潜り、太い神経回路を繋ぎながら、フロー状態という海の底へ潜っていく。\n\nそこに、ポケットのスマホが「ブルッ」と震える。',
    eyebrow='THE DEEP DIVE')
add(slide_keyword, 'コンテキスト・スイッチ', '別の文脈の回路を、無理やり繋ぎ直す瞬間。',
    big_color=PINK, big_size=80)
# 巨大 23 分
add(slide_stat, '23', '分', '深い集中状態 (フロー) に戻るために脳が支払う、たった 1 通の通知の代償。',
    big_color=GOLD)
add(slide_full_quote_dark,
    'たった 1 分の作業のために、\n脳は 23 分のペナルティを払う。',
    'これを 1 日に何十回も繰り返したら、どうなるか。', big_size=46)
add(slide_keyword, '機能不全。', '夕方には、何もしていないのにただ疲れている。',
    big_color=WARN_RED, big_size=140)
add(slide_full_quote_dark,
    '「忙しい、時間が足りない」の本当の正体は、\n時間が物理的にないのではなく、\n細かな文脈の切り替えで脳がショートしている状態。',
    '時間ではなく、注意。', big_size=34)

# --- SECTION 5: 破滅 23 歳 (7) ---
add(slide_section, 5, '破滅へのカウントダウン', 'COUNTDOWN TO COLLAPSE')
add(slide_4cards, '当時の僕の 1 日', [
    ('🦷  医療の脳', '大阪大学 歯学部の実習で、人の歯を精密に削る。'),
    ('🎓  学生の脳', '夜中はオンライン授業を受ける。'),
    ('💼  経営者の脳', '隙間時間でメディア事業のトラブル対応。'),
    ('🚇  移動中のスマホ', '電車・トイレ・廊下、すべてが業務時間。'),
], eyebrow='A DAY IN THE LIFE')
add(slide_keyword, '5 分単位', '脳をカチャカチャ切り替えていた。',
    big_color=WARN_RED, big_size=170)
add(slide_full_quote_dark,
    'ある日、パソコンの画面の文字が、\n突然ゲシュタルト崩壊を起こして\n読めなくなりました。',
    '「隙間時間を活用する」という言葉に取り憑かれた、その代償。', big_size=42)
add(slide_4cards, '僕の体に起きた症状', [
    ('🌀  ゲシュタルト崩壊', '画面の文字が、文字に見えなくなる。'),
    ('💓  動悸', '心臓のリズムが、止まらない。'),
    ('🌡️  原因不明の微熱', '何週間も、続いた。'),
    ('🌫️  焦燥感', '何もできていない、という焦り。'),
], eyebrow='SYMPTOMS')
add(slide_full_quote_dark,
    '僕の脳は、\n限界を超えて煙を上げていた。',
    '完全に狂っていました。', big_size=54)

# --- SECTION 6: 虚無の EXIT (7) ---
add(slide_section, 6, '虚無の EXIT', 'THE EMPTY EXIT')
add(slide_keyword, '23 歳', 'メディア事業を売却。',
    big_color=GOLD, big_size=200)
add(slide_quote,
    '「学生のうちに事業売却なんて、\n君は天才だ！」',
    '大人の経営者たちのチヤホヤ。')
add(slide_full_quote_dark,
    '事業譲渡の契約書にハンコを押して、\n一人でカフェに入り、コーヒーを飲んだ時。',
    '何が起きたか。', big_size=44)
add(slide_keyword, '1 ミリも、動かなかった。', '心は、空っぽだった。',
    big_color=INK, big_size=70)
add(slide_full_quote_dark,
    '労働集約型のビジネスモデルの先にあるのは、\n栄光でも自由でもなく、\n人間性の死だと思い知らされた。',
    'これ以上事業を大きくしようと思えば、睡眠時間をゼロにするしかない。', big_size=36)
add(slide_keyword, '人間性の死', 'これが、成功の呪い。',
    big_color=WARN_RED, big_size=130)

# --- SECTION 7: 100 カ国の放浪 (9) ---
add(slide_section, 7, '100 カ国の放浪と右脳の解放', 'WANDERING 100 COUNTRIES')
add(slide_keyword, '片道切符。', '事業売却で得たお金で、買ったもの。',
    big_color=GOLD, big_size=180)
add(slide_full_quote_dark,
    '高級時計やタワマンの代わりに、\nバックパック一つを背負って、\n世界 100 カ国を放浪する旅に出た。',
    'あの時の僕は、とにかく左脳を休ませたかった。', big_size=42)
add(slide_full_quote_dark,
    'スペイン・バルセロナ。\nサグラダ・ファミリア。',
    'ガウディが、効率もタイムマネジメントも完全に無視して建てた、神への祈り。', big_size=54)
add(slide_keyword, '涙が止まらなかった。', 'なぜか、わからないまま。',
    big_color=PURPLE, big_size=100)
add(slide_4cards, '右脳に流し込んだ風景', [
    ('🏰  サグラダ・ファミリア', '効率を完全に無視した、狂気の美しさ。'),
    ('🧒  東南アジアの路地裏', '言葉が通じない子供たちと、ボール一つで日が暮れる。'),
    ('🪨  ヨーロッパの石畳', '何百年も前の人々の息遣い、歴史の匂い。'),
    ('🌅  砂漠の夕日', '時計のない世界に、光だけが沈んでいく。'),
], eyebrow='100 COUNTRIES')
add(slide_full_quote_dark,
    'ビジネスパーソンとして、\n一番アクセルを踏むべき 20 代前半。\nその時間を、僕は感性に投資した。',
    '効率なんてどうでもいい。生産性なんて知ったことか。', big_size=38)
add(slide_keyword, '右脳の解放', 'ただ、心が震えるものだけを吸収し続けた。',
    big_color=GOLD, big_size=110)
add(slide_quote,
    '「美しい」と感じるもの、\n「心が震える」ものだけを、\nスポンジのように吸収し続けた。',
    '20 代前半、3 年間。')

# --- SECTION 8: チェロの旋律 (8) ---
add(slide_section, 8, 'チェロと、人間本来の周波数', 'CELLO & THE HUMAN FREQUENCY')
add(slide_keyword, '🎻', '旅の中で、僕を救ってくれたもの。',
    big_color=GOLD, big_size=240)
add(slide_quote,
    'どれだけ完璧なエクセルを見せても、\n人は涙を流さない。',
    'ビジネスの世界で、痛感したこと。')
add(slide_full_quote_dark,
    'チェロから生まれる音色は、\n人間の声の周波数に最も近い。',
    'だから、論理を介さずに人と人が共鳴する。', big_size=44)
add(slide_keyword, '言葉が通じなくても。', 'たった 1 本の弦が震えるだけで、心が繋がる。',
    big_color=PURPLE, big_size=80)
add(slide_full_quote_dark,
    'メールの定型文を打つこと?\n領収書をエクセルに入力すること?',
    'いいえ、絶対に違う。', big_size=46)
add(slide_keyword, '創造。', '人間がやるべき仕事の、たった一つの答え。',
    big_color=GOLD, big_size=200)
add(slide_full_quote_dark,
    '左脳がやるべき「処理・作業・オペレーション」に、\n人間の尊い命の時間を奪わせてはいけない。',
    'この強烈な確信が、僕を究極のシステム開発へ駆り立てた。', big_size=34)

# --- SECTION 9: 予防医学からの最終警告 (10) ---
add(slide_section, 9, '予防医学からの最終警告', 'THE FINAL MEDICAL WARNING')
add(slide_full_quote_dark,
    '帰国した僕が見たのは、\nタスクの山に埋もれて溺れかけている、\n友人たちの「死んだ魚のような目」。',
    '優秀な経営者、才能あふれるクリエイターが、左脳的作業に殺されていた。', big_size=34)
add(slide_keyword, '緩やかな自殺。', 'タスクの過負荷とコンテキストスイッチが引き起こすもの。',
    big_color=WARN_RED, big_size=110)
add(slide_4cards, 'コルチゾールが起こすこと', [
    ('🧠  脳の物理的破壊', '記憶や学習を司る「海馬」の神経細胞が萎縮する。'),
    ('🦷  夜間の食いしばり', '体重ほどの圧力で、無意識に歯を削っている。'),
    ('💔  心筋への負担', '常時の緊張は、心臓にも蓄積する。'),
    ('🔥  慢性炎症', '免疫が下がり、あらゆる病気のリスクが上がる。'),
], eyebrow='WHAT STRESS DOES TO YOUR BODY')
add(slide_keyword, '海馬', 'が、文字通り萎縮 (いしゅく) して縮んでいく。',
    big_color=WARN_RED, big_size=200)
add(slide_keyword, 'ブラキシズム', '夜間の歯ぎしり。体重ほどの圧力で、歯を削る。',
    big_color=WARN_RED, big_size=110)
add(slide_full_quote_dark,
    '「時間が足りない」と焦りながら寝ている間、\nあなたの歯は、\nギリギリと削り取られている。',
    'ヒビが入り、最終的に割れて抜歯へ向かっている。', big_size=38)
add(slide_full_quote_dark,
    '事業をスケールさせたい。\nもっと影響力を持ちたい。\nその野心は、最高に美しい。',
    'でも、そのために脳を萎縮させ、歯を砕き、命の時間を削ってはいけない。', big_size=36)
add(slide_keyword, '気合根性論を、\n今日で捨ててください。', '昭和の精神論は、もう通用しない。',
    big_color=INK, big_size=64)
add(slide_full_quote_dark,
    '人間の肉体は、\n現代のこの異常な情報量とスピードを\n処理できるように進化していない。',
    '脱出しなければならない。今すぐ。物理的に。', big_size=38)

# --- SECTION 10: 道具から分身へ (10) ---
add(slide_section, 10, '道具から、分身へ', 'FROM TOOL TO TWIN')
add(slide_keyword, '人を雇う?', 'もう一つの罠が、待っている。',
    big_color=PINK, big_size=130)
add(slide_full_quote_dark,
    '採用コスト・教育コスト、そして何より\n「マネジメント」という巨大なタスクとストレスを\n新たに抱え込むことを意味する。',
    '「なぜ指示通りに動かないんだ」「辞められたらどうしよう」。これでは別のコンテキストスイッチ。', big_size=34)
add(slide_keyword, '究極のAIエージェント。', '中世の修道士が作った時間の奴隷から完全に解放されるために。',
    big_color=GOLD, big_size=64)
add(slide_two_col, 'ChatGPT と、僕が作ったもの。',
    'TOOL ・ 道具',
    'あなたがプロンプトを入力しないと動かない。\n\n結局、あなたが道具を持ち出し、電源を入れ、的を絞るという\n「時間と労力」を使っている。\n\n例えるなら、高機能な電動ドリル。',
    'AGENT ・ 分身',
    '僕の価値観をインストールされ、\n僕の代わりにシステムと連携し、\n\n僕が寝ている間も自律的に判断して動き続ける、\n\nもう一人の井出直毅。',
    eyebrow='PARADIGM SHIFT',
    left_color=INK_DIM, right_color=GOLD)
add(slide_keyword, 'もう一人の自分。', 'AI に丸投げできるのは、左脳的な作業のすべて。',
    big_color=PURPLE, big_size=100)
add(slide_full_quote_dark,
    '自分がやらなくていい左脳作業を AI に丸投げし、\n自分の右脳と命を守り抜く。',
    'これが、究極の自己拡張 OS。', big_size=44)
add(slide_keyword, 'Core Identity OS.', '本日、皆さんにお持ちした、究極の自己拡張 OS。',
    big_color=GOLD, big_size=110)
add(slide_full_quote_dark,
    'CORE = 核。\n中世の時計・産業革命・脳科学の限界、\nそのすべてを超えた、\n人間の本質を抽出する OS。',
    '', big_size=38)
add(slide_4cards, 'Core Identity OS の構造', [
    ('🌐  CORE Identity OS', '上位の OS。あなたという「核」を中心に置く。'),
    ('🔷  CORE Prism', '事業家・複数役割を持つ個人のための分身。'),
    ('🌸  CORE Iris', 'クリエイター・表現者のための分身。'),
    ('✨  さらに増殖する分身たち', 'Echo, Sage, Sanctuary, Atlas — あなたの数だけ。'),
], eyebrow='ARCHITECTURE')

# --- SECTION 11: CORE Prism (11) ---
add(slide_section, 11, 'CORE Prism — 事業家の左脳代行', 'PRISM · LEFT BRAIN OUTSOURCED')
add(slide_keyword, 'CORE Prism', 'すべての事業家、多角経営者のための AI エージェント。',
    big_color=PRISM_BLUE, big_size=110)
add(slide_seven_agents)
add(slide_full_quote_dark,
    'SF 映画のコックピットのようで、\nちょっとワクワクしませんか?',
    'でも、ただデザインがカッコいいだけではない。', big_size=46)
add(slide_full_quote_dark,
    'バラバラに散らばっていたすべてのデータが、\nプリズムという一つのレンズを通して統合される。',
    '不動産事業、スクール運営、店舗ビジネス。', big_size=38)
add(slide_keyword, '日中、僕が患者の歯型に集中している間。', 'スマホの裏側で、Prism は何をしているか。',
    big_color=INK, big_size=44)
add(slide_4cards, 'Prism が、自動でやっていること', [
    ('🏠  M&A 案件のスクレイピング', '利益率・立地条件で自動フィルタ + 要約。'),
    ('💳  Stripe 決済データの確認', '生徒さんの支払いをリアルタイムで把握。'),
    ('📇  HubSpot CRM 自動更新', '顧客情報・ステージ・ネクストアクション。'),
    ('📨  サンクスメール自動送信', '次のステップの案内まで完結。'),
], eyebrow='WHILE YOU SLEEP')
add(slide_keyword, '夜、画面を開く時。', 'やるべき作業は、1 ミリも残っていない。',
    big_color=GOLD, big_size=80)
add(slide_full_quote_dark,
    'あるのは、\n「明日、経営者として下すべき、\nたった一つの重要な決断」だけ。',
    '美しく整理されて、浮かび上がっている。', big_size=42)
add(slide_keyword, '左脳の完全代行。', 'あなたはコーヒーを飲みながら、加速していく事業を見届ければいい。',
    big_color=PRISM_BLUE, big_size=80)
add(slide_quote,
    '「人を増やさず、売上を伸ばす。」',
    'これが、Prism が事業家にもたらす答え。')

# --- SECTION 12: CORE Iris (10) ---
add(slide_section, 12, 'CORE Iris — 体温が届く、もう一人のあなた', 'IRIS · WARMTH AT SCALE')
add(slide_keyword, 'CORE Iris', 'すべてのインフルエンサー、表現者のための AI エージェント。',
    big_color=PINK, big_size=110)
add(slide_six_facets)
add(slide_4cards, 'フォロワー数の悲しい変遷', [
    ('💬  100 人', '全員のコメントに、心を込めて返信できた。'),
    ('🌊  1 万人', '物理的に、無理。一律「いいね」しか押せない。'),
    ('🌪️  10 万人', 'ファンは「自分はもうモブだ」と悟って離れていく。'),
    ('💔  影響力 ↑ = 距離 ↑', '現代の発信者が抱える、最も残酷なジレンマ。'),
], eyebrow='THE PARADOX OF INFLUENCE')
add(slide_keyword, 'Iris は、\nそれを終わらせる。', 'ボットではない。あなたの瞳になる。',
    big_color=INK, big_size=64)
add(slide_full_quote_dark,
    '過去の発信データ、よく使う言い回し、\n言葉の裏のニュアンス、そして何より\n「あなたの体温」を学習する。',
    '本物の「あなたが書いた」と区別がつかないレベルで。', big_size=34)
add(slide_keyword, '"あなただけ"', 'のメッセージが、数万人に届き続ける。',
    big_color=PINK, big_size=100)
add(slide_full_quote_dark,
    'ファンから見れば、\n大好きなあなたから、自分だけに向けた\n温かいメッセージが届き続ける。',
    'これ以上のエンゲージメントはない。', big_size=38)
add(slide_keyword, '24 / 7', 'Iris は、あなたが眠っている間も働き続ける。',
    big_color=PINK, big_size=200)
add(slide_full_quote_dark,
    'あなたは完全に解放された右脳を使って、\n誰にも真似できない最高のコンテンツを、\nただ作り続ければいい。',
    'DM 返信に、人生を奪わせない。', big_size=34)

# --- SECTION 13: AI が仕事を奪う錯覚 (9) ---
add(slide_section, 13, 'AI が仕事を奪う、という錯覚', 'THE GREATEST ILLUSION')
add(slide_quote,
    '「AI に仕事を奪われ、\n人間はただの抜け殻になるのではないか?」',
    '世間が連日、報道していること。')
add(slide_keyword, '完全な錯覚。', '断言します。',
    big_color=WARN_RED, big_size=160)
add(slide_full_quote_dark,
    '同じフォーマットで数値を打ち込み、\nコピペでメールを作り、\nひたすら同じような告知を投稿する。',
    'それ、すでに皆さんが「機械」になっているということ。', big_size=36)
add(slide_keyword, 'AI は、防波堤。', '人間が人間らしくあるための、最強の防波堤。',
    big_color=GOLD, big_size=80)
add(slide_full_quote_dark,
    'AI にチェロを弾かせるためでも、\nAI に世界を旅させるためでもない。\n\n僕が、自分の手でチェロを弾くため。\n僕が、自分の足でサグラダ・ファミリアに行くため。',
    '', big_size=34)
add(slide_keyword, '純度 100%。', 'あなたという人間の Core を、抽出する。',
    big_color=GOLD, big_size=140)
add(slide_full_quote_dark,
    'Core Identity OS は、\nあなたから自分を奪う装置ではない。\n不純物を取り除き、最も美しい部分だけを抽出する装置。',
    '抽出された Core が、世界を動かす。', big_size=32)

# --- SECTION 14: シンフォニー (9) ---
add(slide_section, 14, 'マルチ・アイデンティティの交響曲', 'THE SYMPHONY OF SELVES')
add(slide_full_quote_dark,
    'これはもはや、\n仕事 (ワーク) ではない。\n一つの交響曲 (シンフォニー)。',
    '左脳の Prism と、右脳の Iris が連携した時。', big_size=44)
add(slide_4cards, '一日の演奏', [
    ('🌅  朝', 'Prism が今日の決断を一つ報告。あなたはタクトで指示を出す。'),
    ('☀️  昼', '誰にも邪魔されない、最高にクリエイティブな時間。'),
    ('🌹  同時に', 'Iris は数万のファンと熱狂的に共鳴している。'),
    ('🌙  夜', '一日の演奏の余韻を、ゆっくり振り返る。'),
], eyebrow='A DAY AS THE CONDUCTOR')
add(slide_two_col, '楽器、それとも指揮台。',
    'BEFORE',
    '4 つの楽器を、全部一人で走り回りながら演奏。\n\n過呼吸になり、倒れかけた。',
    'AFTER',
    '指揮台に立ち、Prism と Iris という最強のオーケストラを従える。\n\nあなたが描くのは「どんな美しい音楽を奏でるか」だけ。',
    eyebrow='YOU BECOME THE CONDUCTOR',
    left_color=INK_DIM, right_color=GOLD)
add(slide_keyword, 'マルチ・アイデンティティ。', 'あなたの可能性は、文字通り「無限」になる。',
    big_color=PURPLE, big_size=70)
add(slide_full_quote_dark,
    'やりたいことの数だけ\nエージェントを生み出し、\n同時に走らせればいい。',
    '一つに縛られる必要は、もうない。', big_size=44)
add(slide_keyword, '∞', '可能性は、無限に。',
    big_color=GOLD, big_size=400)
add(slide_quote,
    '人類は、\nダ・ヴィンチを取り戻す。',
    '500 年ぶりの復活。')

# --- SECTION 15: 不可逆の扉 (8) ---
add(slide_section, 15, '不可逆の扉', 'THE IRREVERSIBLE DOOR')
add(slide_full_quote_dark,
    '今日、皆さんは知ってしまった。',
    '時間の呪縛から逃れ、命と健康を守りながら、事業をスケールさせる方法がすでに存在することを。', big_size=54)
add(slide_keyword, '人間の脳は、残酷。', '一度「高い基準」を知った脳は、二度と元の低い基準に戻れない。',
    big_color=INK, big_size=64)
add(slide_full_quote_dark,
    '明日、皆さんがオフィスに戻ります。\nメールの山を開き、\nキーボードを叩いて定型文を作り始めた時。',
    '想像してみてください。', big_size=38)
add(slide_quote,
    '「ああ、僕は一体何をやっているんだ。\nこれは、僕がやるべきことじゃない。」',
    '指が止まる。ため息が出る。')
add(slide_keyword, '気づいてしまった。', '無駄な作業に、命の時間がドブに捨てられていることを。',
    big_color=WARN_RED, big_size=80)
add(slide_keyword, '昨日までの自分には、\nもう、戻れない。', '知らなければ、幸せだったかもしれない。',
    big_color=INK, big_size=64)
add(slide_full_quote_dark,
    'でも、もう、\n扉は開いてしまった。',
    '', big_size=64)

# --- SECTION 16: グランドフィナーレ (12) ---
add(slide_section, 16, 'グランドフィナーレ', 'GRAND FINALE')
add(slide_keyword, '一般非公開。', 'Core Identity OS は、まだ市場に出していない。',
    big_color=GOLD, big_size=110)
add(slide_full_quote_dark,
    '単なる「業務効率化ツール」ではないから。\n人間の生き方そのものを変える、\n強大すぎる「OS」だから。',
    '現状維持で満足している人、楽をしたいだけの人に、この武器は渡せない。', big_size=34)
add(slide_keyword, '今日、ここにいる皆さんに。', '一瞬も目を離さず、最後まで耳を傾けてくれた皆さんに。',
    big_color=INK, big_size=54)
add(slide_qr)  # 巨大 QR
add(slide_keyword, '片道切符。', '過去の延長線上から引き剥がす、たった一つの鍵。',
    big_color=GOLD, big_size=130)
add(slide_two_col, '今、この場で、決断してください。',
    'OPTION A',
    '産業革命時代の名残のまま、\n自分の脳と命を削って、\n限界のある労働集約のビジネスを続ける。\n\n— 知らないフリをして生きる —',
    'OPTION B',
    '自分の「分身」を手に入れ、\n心身ともに健康なまま、\nあなたが本当にやりたかった創造の世界へ飛び立つ。\n\n— 扉を、開く —',
    eyebrow='THE CHOICE',
    left_color=INK_FAINT, right_color=GOLD)
add(slide_full_quote_dark,
    '未来のビジネスの価値は、\n「いかに苦労して作業したか」ではない。',
    '23 歳で絶望し、世界中を旅して、たどり着いた真実。', big_size=42)
add(slide_full_quote_dark,
    'いかに感性を研ぎ澄まし、\nAI と共に、誰も見たことがない\n美しい世界を描けたか。',
    'これが、未来の価値の決定式。', big_size=42)
add(slide_keyword, 'あちら側で、お会いしましょう。', '限界のない世界で。',
    big_color=GOLD, big_size=64)
add(slide_keyword, 'ありがとうございました。', '井出直毅 / 2026.05.12 / 60 分間、本当に。',
    big_color=INK, big_size=72)

# --- END (3) ---
add(slide_qr)  # 念押しの QR
add(slide_keyword, 'KEYNOTE30', '講演会限定コード — トライアル 30 日延長 (5/15 実装予定)',
    big_color=GOLD, big_size=130)
add(slide_keyword, 'Core Identity OS', 'produced by 株式会社コアプリズム (仮) · 2026',
    big_color=PURPLE, big_size=70)


# === 実行 ===
TOTAL = len(slides_def)
print(f'Generating {TOTAL} slides...')
for idx, (fn, args, kwargs) in enumerate(slides_def, 1):
    fn(*args, page=idx, total=TOTAL, **kwargs)

# 保存
out_path = os.path.expanduser('~/Desktop/CORE_キーノート_2026-05-11.pptx')
prs.save(out_path)
print(f'✓ Saved: {out_path}')
print(f'  Total slides: {TOTAL}')
