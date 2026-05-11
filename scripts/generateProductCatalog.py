#!/usr/bin/env python3
# ============================================================
# CORE Identity OS — 製品カタログ (Apple トーン)
# 出力: ~/Desktop/CORE_製品カタログ_2026-05-11.pptx
# 約 200 枚、ライト基調 + ダーク章扉のサンドイッチ
# ============================================================
import os
import qrcode
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── パレット (Apple トーン: 白 + 黒 + 戦略的アクセント) ──
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xFA, 0xFA, 0xFC)
BG_GRAY    = RGBColor(0xF2, 0xF2, 0xF7)   # Apple System Gray
BG_DARK    = RGBColor(0x0B, 0x0B, 0x10)
INK        = RGBColor(0x0A, 0x0A, 0x14)   # 純黒回避
INK_SOFT   = RGBColor(0x1D, 0x1D, 0x1F)
INK_MUTED  = RGBColor(0x6E, 0x6E, 0x73)
INK_FAINT  = RGBColor(0xA1, 0xA1, 0xA6)
DARK_INK   = RGBColor(0xF5, 0xF5, 0xF7)
DARK_MUTED = RGBColor(0xA1, 0xA1, 0xA6)

# ブランドアクセント
PRISM_BLUE = RGBColor(0x00, 0x71, 0xE3)   # Apple Blue
PRISM_DEEP = RGBColor(0x00, 0x33, 0xA0)
IRIS_PINK  = RGBColor(0xE1, 0x30, 0x6C)
IRIS_GOLD  = RGBColor(0xFC, 0xB0, 0x45)
IRIS_PURPLE= RGBColor(0xB0, 0x7B, 0xD9)
CORE_GOLD  = RGBColor(0xFC, 0xB0, 0x45)
GREEN      = RGBColor(0x34, 0xC7, 0x59)
RED        = RGBColor(0xFF, 0x3B, 0x30)

# 7 エージェントカラー (Apple System Colors 系)
AGENT_COLORS = {
    '経営':   RGBColor(0xFF, 0x3B, 0x30),  # Red
    '営業':   RGBColor(0xFF, 0x95, 0x00),  # Orange
    '財務':   RGBColor(0xFF, 0xCC, 0x00),  # Yellow
    '創造':   RGBColor(0x34, 0xC7, 0x59),  # Green
    '学び':   RGBColor(0x00, 0x7A, 0xFF),  # Blue
    '人材':   RGBColor(0xAF, 0x52, 0xDE),  # Purple
    '生活':   RGBColor(0xFF, 0x2D, 0x55),  # Pink
}

# フォント (英語系 + 日本語系)
F_DISPLAY = 'SF Pro Display'
F_TEXT    = 'SF Pro Text'
F_JP      = 'Hiragino Kaku Gothic ProN'
F_JP_MIN  = 'Hiragino Mincho ProN'
F_NUM     = 'SF Mono'

# 16:9
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)
BLANK = prs.slide_layouts[6]


# ─── ヘルパー ────────────────────────────────

def bg(slide, color=WHITE):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, t, left, top, width, height,
         font=F_JP, size=14, color=INK, bold=False, italic=False,
         align='left', anchor='top', letter_spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {
        'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM
    }[anchor]
    p = tf.paragraphs[0]
    p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}[align]
    for i, line in enumerate(t.split('\n')):
        if i == 0:
            run = p.add_run()
            run.text = line
        else:
            p2 = tf.add_paragraph()
            p2.alignment = p.alignment
            run = p2.add_run()
            run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        if letter_spacing is not None:
            rPr = run._r.get_or_add_rPr()
            rPr.set('spc', str(letter_spacing))
    return tb


def rect(slide, x, y, w, h, fill=None, line=None, alpha=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            sp = sh.fill._xPr
            solidFill = sp.find(qn('a:solidFill'))
            srgb = solidFill.find(qn('a:srgbClr')) if solidFill is not None else None
            if srgb is not None:
                a = srgb.find(qn('a:alpha'))
                if a is None:
                    a = etree.SubElement(srgb, qn('a:alpha'))
                a.set('val', str(int(alpha * 100000)))
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def circle(slide, cx, cy, r, fill=None, alpha=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            sp = sh.fill._xPr
            solidFill = sp.find(qn('a:solidFill'))
            srgb = solidFill.find(qn('a:srgbClr')) if solidFill is not None else None
            if srgb is not None:
                a = srgb.find(qn('a:alpha'))
                if a is None:
                    a = etree.SubElement(srgb, qn('a:alpha'))
                a.set('val', str(int(alpha * 100000)))
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def footer(slide, page, total, dark=False):
    fg = DARK_MUTED if dark else INK_FAINT
    # ロゴテキスト風 (左下)
    text(slide, 'CORE  Identity  OS', Inches(0.6), Inches(SLIDE_H_IN - 0.5), Inches(3.5), Inches(0.25),
         font=F_DISPLAY, size=8, color=fg, letter_spacing=400, align='left')
    # ページ番号 (右下)
    text(slide, f'{page:03d}', Inches(SLIDE_W_IN - 1.2), Inches(SLIDE_H_IN - 0.5), Inches(0.7), Inches(0.25),
         font=F_NUM, size=8, color=fg, align='right')


# ─── スライドテンプレート ────────────────────────────────

def slide_cover(eyebrow, title_main, title_sub, presenter):
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    # 中央寄りやや上のタイトル
    text(s, eyebrow, Inches(0.7), Inches(0.9), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
         font=F_DISPLAY, size=11, color=INK_MUTED, bold=True, letter_spacing=600, align='left')
    text(s, title_main, Inches(0.7), Inches(2.4), Inches(SLIDE_W_IN - 1.4), Inches(2.5),
         font=F_JP, size=92, color=INK, bold=True, align='left')
    text(s, title_sub, Inches(0.7), Inches(4.9), Inches(SLIDE_W_IN - 1.4), Inches(0.9),
         font=F_JP_MIN, size=24, color=INK_MUTED, italic=True, align='left')
    # アクセントライン (横棒、ゴールド)
    rect(s, Inches(0.7), Inches(6.0), Inches(2.6), Inches(0.04), fill=CORE_GOLD)
    text(s, presenter, Inches(0.7), Inches(6.25), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
         font=F_TEXT, size=12, color=INK_MUTED, align='left', letter_spacing=200)
    return s


def slide_chapter(num, title_jp, title_en, accent=CORE_GOLD):
    """ダーク章扉"""
    s = prs.slides.add_slide(BLANK)
    bg(s, BG_DARK)
    # 巨大番号 (右下に薄く)
    text(s, f'{num:02d}', Inches(SLIDE_W_IN - 4), Inches(SLIDE_H_IN - 3.5), Inches(3.5), Inches(3.0),
         font=F_DISPLAY, size=260, color=RGBColor(0x2A, 0x2A, 0x30), bold=True, align='right')
    # eyebrow
    text(s, title_en, Inches(0.7), Inches(2.5), Inches(SLIDE_W_IN - 1.4), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=700, align='left')
    # 日本語タイトル
    text(s, title_jp, Inches(0.7), Inches(3.1), Inches(SLIDE_W_IN - 1.4), Inches(2.0),
         font=F_JP, size=80, color=DARK_INK, bold=True, align='left')
    # アクセント
    rect(s, Inches(0.7), Inches(5.4), Inches(1.6), Inches(0.04), fill=accent)
    return s


def slide_full_statement(text_main, text_sub=None, light=True, big_size=64):
    """大きな宣言文 (1 ステートメント 1 スライド)"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE if light else BG_DARK)
    fg = INK if light else DARK_INK
    fg_dim = INK_MUTED if light else DARK_MUTED
    text(s, text_main, Inches(1.0), Inches(2.2), Inches(SLIDE_W_IN - 2.0), Inches(3.5),
         font=F_JP, size=big_size, color=fg, bold=True, align='left')
    if text_sub:
        text(s, text_sub, Inches(1.0), Inches(5.7), Inches(SLIDE_W_IN - 2.0), Inches(1.0),
             font=F_JP_MIN, size=20, color=fg_dim, italic=True, align='left')
    return s


def slide_quote_big(quote, attribution=None, light=True):
    """超大きな引用"""
    s = prs.slides.add_slide(BLANK)
    bg(s, OFF_WHITE if light else BG_DARK)
    fg = INK if light else DARK_INK
    text(s, '"', Inches(0.8), Inches(0.9), Inches(2), Inches(2),
         font='Georgia', size=180, color=CORE_GOLD, bold=True, align='left')
    text(s, quote, Inches(1.2), Inches(2.8), Inches(SLIDE_W_IN - 2.4), Inches(3.5),
         font=F_JP_MIN, size=44, color=fg, italic=True, bold=True, align='left')
    if attribution:
        text(s, '— ' + attribution, Inches(1.2), Inches(6.3), Inches(SLIDE_W_IN - 2.4), Inches(0.4),
             font=F_TEXT, size=14, color=INK_MUTED, align='left')
    return s


def slide_number_hero(number, label, sub, color=PRISM_BLUE, light=True):
    """巨大数字"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE if light else BG_DARK)
    fg_label = INK_MUTED if light else DARK_MUTED
    fg_sub = INK if light else DARK_INK
    text(s, label, Inches(1.0), Inches(1.5), Inches(SLIDE_W_IN - 2.0), Inches(0.5),
         font=F_DISPLAY, size=12, color=fg_label, bold=True, letter_spacing=600, align='left')
    text(s, number, Inches(1.0), Inches(2.2), Inches(SLIDE_W_IN - 2.0), Inches(3.5),
         font=F_DISPLAY, size=260, color=color, bold=True, align='left')
    if sub:
        text(s, sub, Inches(1.0), Inches(5.9), Inches(SLIDE_W_IN - 2.0), Inches(1.0),
             font=F_JP, size=22, color=fg_sub, align='left')
    return s


def slide_brand_hero(brand_name, tagline_jp, tagline_en, accent, light=True):
    """ブランドの大ヒーロー"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE if light else BG_DARK)
    fg = INK if light else DARK_INK
    fg_dim = INK_MUTED if light else DARK_MUTED
    # eyebrow EN
    text(s, tagline_en, Inches(1.0), Inches(1.4), Inches(SLIDE_W_IN - 2.0), Inches(0.5),
         font=F_DISPLAY, size=12, color=accent, bold=True, letter_spacing=700, align='left')
    # ブランド名 (超大)
    text(s, brand_name, Inches(1.0), Inches(2.1), Inches(SLIDE_W_IN - 2.0), Inches(2.5),
         font=F_DISPLAY, size=160, color=fg, bold=True, align='left')
    # タグライン
    text(s, tagline_jp, Inches(1.0), Inches(4.9), Inches(SLIDE_W_IN - 2.0), Inches(1.2),
         font=F_JP, size=32, color=fg, align='left')
    # アクセント
    rect(s, Inches(1.0), Inches(6.2), Inches(2.6), Inches(0.05), fill=accent)
    return s


def slide_section_title(title, sub=None, accent=CORE_GOLD, light=True):
    """中規模セクションタイトル"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE if light else BG_DARK)
    fg = INK if light else DARK_INK
    fg_dim = INK_MUTED if light else DARK_MUTED
    text(s, title, Inches(1.0), Inches(2.8), Inches(SLIDE_W_IN - 2.0), Inches(2.0),
         font=F_JP, size=64, color=fg, bold=True, align='left')
    rect(s, Inches(1.0), Inches(5.0), Inches(1.6), Inches(0.05), fill=accent)
    if sub:
        text(s, sub, Inches(1.0), Inches(5.2), Inches(SLIDE_W_IN - 2.0), Inches(1.0),
             font=F_JP, size=20, color=fg_dim, align='left')
    return s


def slide_feature(eyebrow, title, body, color=PRISM_BLUE, side='right'):
    """機能 1 つを詳しく (左テキスト + 右ビジュアル円)"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    # 左テキスト
    text(s, eyebrow, Inches(0.9), Inches(1.4), Inches(6.5), Inches(0.4),
         font=F_DISPLAY, size=11, color=color, bold=True, letter_spacing=600, align='left')
    text(s, title, Inches(0.9), Inches(2.0), Inches(6.5), Inches(2.2),
         font=F_JP, size=52, color=INK, bold=True, align='left')
    text(s, body, Inches(0.9), Inches(4.4), Inches(6.5), Inches(2.4),
         font=F_JP, size=18, color=INK_MUTED, align='left', line_spacing=1.55)
    # 右側ビジュアル (色付き円 + 内部の白円)
    cx = Inches(SLIDE_W_IN - 3.0)
    cy = Inches(SLIDE_H_IN / 2)
    circle(s, cx, cy, Inches(2.2), fill=color, alpha=0.08)
    circle(s, cx, cy, Inches(1.5), fill=color, alpha=0.18)
    circle(s, cx, cy, Inches(0.8), fill=color)
    return s


def slide_agent_card(num, role_jp, role_en, body, color):
    """エージェント 1 つを詳しく"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    # 大きな円アクセント (左、淡色)
    circle(s, Inches(2.2), Inches(SLIDE_H_IN / 2), Inches(1.8), fill=color, alpha=0.12)
    circle(s, Inches(2.2), Inches(SLIDE_H_IN / 2), Inches(0.7), fill=color)
    # 番号
    text(s, f'AGENT {num:02d}', Inches(0.9), Inches(0.9), Inches(5), Inches(0.4),
         font=F_DISPLAY, size=10, color=INK_MUTED, bold=True, letter_spacing=600, align='left')
    # role en
    text(s, role_en, Inches(4.6), Inches(2.0), Inches(SLIDE_W_IN - 5.3), Inches(0.4),
         font=F_DISPLAY, size=11, color=color, bold=True, letter_spacing=500, align='left')
    # role jp
    text(s, role_jp, Inches(4.6), Inches(2.5), Inches(SLIDE_W_IN - 5.3), Inches(1.5),
         font=F_JP, size=64, color=INK, bold=True, align='left')
    # body
    text(s, body, Inches(4.6), Inches(4.5), Inches(SLIDE_W_IN - 5.3), Inches(2.2),
         font=F_JP, size=18, color=INK_MUTED, align='left')
    return s


def slide_three_columns(title, items, accent=PRISM_BLUE, eyebrow=None):
    """3 列 (各列: タイトル + 本文)"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    if eyebrow:
        text(s, eyebrow, Inches(0.9), Inches(0.9), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
             font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, title, Inches(0.9), Inches(1.5), Inches(SLIDE_W_IN - 1.8), Inches(1.6),
         font=F_JP, size=44, color=INK, bold=True, align='left')
    # 3 列
    col_w = Inches((SLIDE_W_IN - 2.4) / 3)
    col_h = Inches(3.0)
    y = Inches(3.6)
    start_x = Inches(0.9)
    gap = Inches(0.3)
    for i, (t_, b_) in enumerate(items[:3]):
        x = start_x + i * (col_w + gap)
        # 番号 (薄く)
        text(s, f'{i+1:02}', x, y, col_w, Inches(0.5),
             font=F_NUM, size=14, color=accent, bold=True, align='left')
        text(s, t_, x, y + Inches(0.5), col_w, Inches(1.0),
             font=F_JP, size=22, color=INK, bold=True, align='left')
        text(s, b_, x, y + Inches(1.5), col_w, col_h - Inches(1.5),
             font=F_JP, size=14, color=INK_MUTED, align='left')
    return s


def slide_six_grid(title, items, accent=IRIS_PINK, eyebrow=None):
    """6 個グリッド (Iris 6 ファセット用)"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    if eyebrow:
        text(s, eyebrow, Inches(0.9), Inches(0.7), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
             font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, title, Inches(0.9), Inches(1.2), Inches(SLIDE_W_IN - 1.8), Inches(1.3),
         font=F_JP, size=40, color=INK, bold=True, align='left')
    # 3x2
    card_w = Inches(3.85)
    card_h = Inches(2.0)
    gap = Inches(0.18)
    start_x = Inches(0.9)
    start_y = Inches(3.0)
    for i, (t_, b_) in enumerate(items[:6]):
        row, col = i // 3, i % 3
        x = start_x + col * (card_w + gap)
        y = start_y + row * (card_h + gap)
        rect(s, x, y, card_w, card_h, fill=BG_GRAY, radius=True)
        # 小ドット
        circle(s, x + Inches(0.3), y + Inches(0.45), Inches(0.06), fill=accent)
        text(s, t_, x + Inches(0.5), y + Inches(0.3), card_w - Inches(0.7), Inches(0.5),
             font=F_JP, size=20, color=INK, bold=True, align='left', anchor='middle')
        text(s, b_, x + Inches(0.35), y + Inches(1.0), card_w - Inches(0.6), card_h - Inches(1.2),
             font=F_JP, size=13, color=INK_MUTED, align='left')
    return s


def slide_seven_agents_overview(accent=PRISM_BLUE):
    """7 エージェントを 1 画面に"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    text(s, 'SEVEN AGENTS · ONE OS', Inches(0.9), Inches(0.7), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, '役割ごとに、専属の AI を。', Inches(0.9), Inches(1.2), Inches(SLIDE_W_IN - 1.8), Inches(1.3),
         font=F_JP, size=40, color=INK, bold=True, align='left')
    agents = [
        ('経営', 'CEO Agent', '戦略立案・KPI 監視'),
        ('営業', 'Sales Agent', '商談スクリプト・提案書'),
        ('財務', 'CFO Agent', 'P&L・キャッシュ予測'),
        ('創造', 'Creative Agent', '画像・スライド・ブランド'),
        ('学び', 'Knowledge Agent', '要約・読書ノート・検索'),
        ('人材', 'People Agent', '1on1・センチメント'),
        ('生活', 'Life Agent', '健康・家族・心の整え'),
    ]
    card_w = Inches(1.62)
    card_h = Inches(3.0)
    gap = Inches(0.06)
    total_w = card_w * 7 + gap * 6
    start_x = Emu((prs.slide_width - total_w) / 2)
    start_y = Inches(3.0)
    for i, (jp, en, desc) in enumerate(agents):
        c = AGENT_COLORS[jp]
        x = start_x + i * (card_w + gap)
        rect(s, x, start_y, card_w, card_h, fill=BG_GRAY, radius=True)
        # 上部ドット
        circle(s, x + card_w / 2, start_y + Inches(0.55), Inches(0.13), fill=c)
        # 番号
        text(s, f'{i+1:02}', x, start_y + Inches(0.85), card_w, Inches(0.3),
             font=F_NUM, size=10, color=INK_FAINT, align='center')
        # JP
        text(s, jp, x, start_y + Inches(1.15), card_w, Inches(0.6),
             font=F_JP, size=22, color=INK, bold=True, align='center')
        # EN
        text(s, en, x, start_y + Inches(1.75), card_w, Inches(0.3),
             font=F_DISPLAY, size=9, color=c, bold=True, letter_spacing=300, align='center')
        # desc
        text(s, desc, x + Inches(0.15), start_y + Inches(2.15), card_w - Inches(0.3), Inches(0.75),
             font=F_JP, size=10, color=INK_MUTED, align='center')
    return s


def slide_six_facets_overview(accent=IRIS_PINK):
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    text(s, 'SIX FACETS OF LIGHT', Inches(0.9), Inches(0.7), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, '光は、6 つの色を持つ。', Inches(0.9), Inches(1.2), Inches(SLIDE_W_IN - 1.8), Inches(1.3),
         font=F_JP, size=40, color=INK, bold=True, align='left')
    facets = [
        ('案件', 'Briefs', '受注→投稿→レポート', IRIS_PINK),
        ('分析', 'Analytics', 'Instagram 解析', IRIS_PURPLE),
        ('創作', 'Creation', 'キャプション・サムネ', IRIS_GOLD),
        ('交渉', 'Negotiation', '料金・媒体資料・提案', RGBColor(0xFF,0xA9,0x4D)),
        ('ブランド', 'Brand', 'フォント・色・トーン', RGBColor(0xC7,0x7D,0xFF)),
        ('仲間', 'Community', '招待制コミュニティ', RGBColor(0xFD,0x7C,0xB8)),
    ]
    card_w = Inches(1.92)
    card_h = Inches(3.0)
    gap = Inches(0.1)
    total_w = card_w * 6 + gap * 5
    start_x = Emu((prs.slide_width - total_w) / 2)
    start_y = Inches(3.0)
    for i, (jp, en, desc, c) in enumerate(facets):
        x = start_x + i * (card_w + gap)
        rect(s, x, start_y, card_w, card_h, fill=BG_GRAY, radius=True)
        circle(s, x + card_w / 2, start_y + Inches(0.55), Inches(0.14), fill=c)
        text(s, f'{i+1:02}', x, start_y + Inches(0.9), card_w, Inches(0.3),
             font=F_NUM, size=10, color=INK_FAINT, align='center')
        text(s, jp, x, start_y + Inches(1.25), card_w, Inches(0.6),
             font=F_JP, size=24, color=INK, bold=True, align='center')
        text(s, en, x, start_y + Inches(1.9), card_w, Inches(0.3),
             font=F_DISPLAY, size=10, color=c, bold=True, letter_spacing=300, align='center')
        text(s, desc, x + Inches(0.15), start_y + Inches(2.3), card_w - Inches(0.3), Inches(0.6),
             font=F_JP, size=11, color=INK_MUTED, align='center')
    return s


def slide_compare_two(title, eyebrow,
                      left_label, left_body, right_label, right_body,
                      accent=PRISM_BLUE, left_color=INK_FAINT, right_color=None):
    if right_color is None:
        right_color = accent
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    text(s, eyebrow, Inches(0.9), Inches(0.8), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, title, Inches(0.9), Inches(1.3), Inches(SLIDE_W_IN - 1.8), Inches(1.4),
         font=F_JP, size=44, color=INK, bold=True, align='left')
    col_w = Inches(5.6)
    col_h = Inches(3.5)
    col_y = Inches(3.2)
    # Left (Before / 他社)
    rect(s, Inches(0.9), col_y, col_w, col_h, fill=BG_GRAY, radius=True)
    text(s, left_label, Inches(1.2), col_y + Inches(0.3), col_w, Inches(0.4),
         font=F_DISPLAY, size=10, color=left_color, bold=True, letter_spacing=500, align='left')
    text(s, left_body, Inches(1.2), col_y + Inches(1.0), col_w - Inches(0.6), col_h - Inches(1.3),
         font=F_JP, size=18, color=INK_MUTED, align='left')
    # Right (After / CORE)
    rect(s, Inches(SLIDE_W_IN - 6.5), col_y, col_w, col_h, fill=right_color, alpha=0.06, radius=True)
    text(s, right_label, Inches(SLIDE_W_IN - 6.2), col_y + Inches(0.3), col_w, Inches(0.4),
         font=F_DISPLAY, size=10, color=right_color, bold=True, letter_spacing=500, align='left')
    text(s, right_body, Inches(SLIDE_W_IN - 6.2), col_y + Inches(1.0), col_w - Inches(0.6), col_h - Inches(1.3),
         font=F_JP, size=18, color=INK, align='left')
    return s


def slide_table_plans(brand, plans_data, accent):
    """プラン比較表"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    text(s, brand.upper() + ' · PLANS', Inches(0.9), Inches(0.7), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, '料金プラン', Inches(0.9), Inches(1.2), Inches(SLIDE_W_IN - 1.8), Inches(1.0),
         font=F_JP, size=40, color=INK, bold=True, align='left')
    # 列
    n = len(plans_data)
    card_w = Inches((SLIDE_W_IN - 1.8) / n - 0.2)
    gap = Inches(0.2)
    start_x = Inches(0.9)
    start_y = Inches(2.8)
    card_h = Inches(4.0)
    for i, plan in enumerate(plans_data):
        x = start_x + i * (card_w + gap)
        is_popular = plan.get('popular', False)
        bg_fill = accent if is_popular else BG_GRAY
        alpha = 0.08 if is_popular else None
        rect(s, x, start_y, card_w, card_h, fill=bg_fill, alpha=alpha, radius=True,
             line=accent if is_popular else None)
        # popular badge
        if is_popular:
            text(s, '人気 No.1', x + Inches(0.25), start_y + Inches(0.15), card_w - Inches(0.5), Inches(0.3),
                 font=F_DISPLAY, size=9, color=accent, bold=True, letter_spacing=300, align='center')
        # plan name
        text(s, plan['name'], x + Inches(0.25), start_y + Inches(0.55), card_w - Inches(0.5), Inches(0.6),
             font=F_DISPLAY, size=22, color=INK, bold=True, align='center')
        # price
        text(s, plan['price'], x + Inches(0.25), start_y + Inches(1.2), card_w - Inches(0.5), Inches(0.8),
             font=F_DISPLAY, size=36, color=INK, bold=True, align='center')
        text(s, plan.get('priceLabel', '/月'), x + Inches(0.25), start_y + Inches(2.0), card_w - Inches(0.5), Inches(0.3),
             font=F_JP, size=11, color=INK_MUTED, align='center')
        # tag
        text(s, plan.get('tag', ''), x + Inches(0.25), start_y + Inches(2.35), card_w - Inches(0.5), Inches(0.5),
             font=F_JP, size=12, color=INK_MUTED, italic=True, align='center')
        # divider
        rect(s, x + Inches(0.5), start_y + Inches(2.85), card_w - Inches(1.0), Inches(0.02), fill=INK_FAINT, alpha=0.4)
        # features (max 5)
        feat_y = start_y + Inches(3.0)
        for j, f in enumerate(plan.get('features', [])[:5]):
            text(s, '· ' + f, x + Inches(0.25), feat_y + Inches(j * 0.18), card_w - Inches(0.5), Inches(0.2),
                 font=F_JP, size=9, color=INK_MUTED, align='left')
    return s


def slide_day_in_life(title, eyebrow, slots, accent):
    """1 日のタイムライン (時刻 + 行動 + AI 補助)"""
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    text(s, eyebrow, Inches(0.9), Inches(0.7), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, title, Inches(0.9), Inches(1.2), Inches(SLIDE_W_IN - 1.8), Inches(1.3),
         font=F_JP, size=40, color=INK, bold=True, align='left')
    y = Inches(2.8)
    for time_label, action, ai_help in slots:
        rect(s, Inches(0.9), y, Inches(SLIDE_W_IN - 1.8), Inches(0.65), fill=BG_GRAY, radius=True)
        text(s, time_label, Inches(1.1), y, Inches(1.4), Inches(0.65),
             font=F_NUM, size=14, color=accent, bold=True, align='left', anchor='middle')
        text(s, action, Inches(2.7), y, Inches(4.2), Inches(0.65),
             font=F_JP, size=13, color=INK, bold=True, align='left', anchor='middle')
        text(s, '→ ' + ai_help, Inches(7.0), y, Inches(SLIDE_W_IN - 7.7), Inches(0.65),
             font=F_JP, size=12, color=INK_MUTED, align='left', anchor='middle')
        y += Inches(0.75)
    return s


def slide_qr_final(url, code='KEYNOTE30'):
    s = prs.slides.add_slide(BLANK)
    bg(s, BG_DARK)
    # QR
    qr = qrcode.QRCode(version=4, error_correction=qrcode.constants.ERROR_CORRECT_M, box_size=20, border=2)
    qr.add_data(url); qr.make(fit=True)
    img = qr.make_image(fill_color='white', back_color='#0B0B10')
    buf = BytesIO(); img.save(buf, format='PNG'); buf.seek(0)
    qr_size = Inches(4.0)
    qr_x = Inches(SLIDE_W_IN - 5.0)
    qr_y = Inches((SLIDE_H_IN - 4.0) / 2)
    s.shapes.add_picture(buf, qr_x, qr_y, qr_size, qr_size)
    # 左テキスト
    text(s, 'START NOW', Inches(0.9), Inches(1.6), Inches(6), Inches(0.4),
         font=F_DISPLAY, size=11, color=CORE_GOLD, bold=True, letter_spacing=600, align='left')
    text(s, 'はじめる。', Inches(0.9), Inches(2.2), Inches(6), Inches(2.0),
         font=F_JP, size=80, color=DARK_INK, bold=True, align='left')
    text(s, '14 日間無料 · クレカ不要', Inches(0.9), Inches(4.5), Inches(6), Inches(0.5),
         font=F_JP, size=18, color=DARK_MUTED, align='left')
    text(s, f'限定コード: {code}', Inches(0.9), Inches(5.1), Inches(6), Inches(0.4),
         font=F_DISPLAY, size=14, color=CORE_GOLD, bold=True, align='left')
    text(s, url, Inches(0.9), Inches(5.6), Inches(6), Inches(0.4),
         font=F_NUM, size=12, color=DARK_MUTED, align='left')
    return s


def slide_simple_bullets(eyebrow, title, bullets, accent=PRISM_BLUE):
    s = prs.slides.add_slide(BLANK)
    bg(s, WHITE)
    text(s, eyebrow, Inches(0.9), Inches(0.9), Inches(SLIDE_W_IN - 1.8), Inches(0.4),
         font=F_DISPLAY, size=11, color=accent, bold=True, letter_spacing=600, align='left')
    text(s, title, Inches(0.9), Inches(1.5), Inches(SLIDE_W_IN - 1.8), Inches(1.5),
         font=F_JP, size=44, color=INK, bold=True, align='left')
    y = Inches(3.4)
    for b in bullets[:6]:
        circle(s, Inches(1.05), y + Inches(0.18), Inches(0.06), fill=accent)
        text(s, b, Inches(1.4), y, Inches(SLIDE_W_IN - 2.4), Inches(0.5),
             font=F_JP, size=18, color=INK, align='left')
        y += Inches(0.6)
    return s


# ========================================================================
# === スライド構成 ========================================================
# ========================================================================

slides = []  # 各 entry: (callable, *args, **kwargs)

def add(fn, *args, **kwargs):
    slides.append((fn, args, kwargs))


# ───────────────────────────────────────────────────────
# Section 0 — COVER & WELCOME (8 枚)
# ───────────────────────────────────────────────────────
add(slide_cover,
    eyebrow='CORE IDENTITY OS · PRODUCT CATALOG 2026',
    title_main='ひとつの肉体に、\n無限の自分を。',
    title_sub='— 役割の数だけ、AI を。',
    presenter='株式会社コアプリズム  ·  Naoki Ide  ·  May 2026')

add(slide_full_statement,
    'これは、\n単なる業務効率化ツールではない。',
    'これは、人間の生き方を変える、自己拡張の OS。', big_size=60)

add(slide_full_statement,
    '中世の時計が作った\n「時間の奴隷」を、終わらせる。',
    'CORE Identity OS — for the post-clock era.', big_size=58)

add(slide_quote_big,
    '左脳の作業を AI に渡し、\n右脳の創造を、取り戻す。',
    'CORE が提案する、新しい働き方。')

add(slide_number_hero, '23', 'MINUTES', '一つの通知で失う、深い集中時間。CORE はこれを取り戻す。', color=RED)

add(slide_number_hero, '7 + 6', 'AGENTS', '事業家のための 7 つ。クリエイターのための 6 つ。', color=PRISM_BLUE)

add(slide_number_hero, '24 / 7', 'AUTONOMOUS', 'あなたが眠っている間も、分身は働き続ける。', color=GREEN)

add(slide_full_statement,
    'CORE は、\n働き方ではなく、\n生き方を変える。', None, big_size=68)


# ───────────────────────────────────────────────────────
# Section 1 — PHILOSOPHY (12 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 1, '思想', 'PHILOSOPHY', CORE_GOLD)

add(slide_full_statement,
    '人間は、\nひとつではない。',
    '経営者・営業・財務・創造者・先生・親 — 役割の数だけ、思考が要る。')

add(slide_quote_big,
    '本来の人間は、\nダ・ヴィンチのように\n複数の顔を持つ。',
    'ルネサンスの記憶。')

add(slide_full_statement,
    '産業革命が、\n人間を「一つの役割」に\n閉じ込めた。',
    '工場のサイレンと、肉の歯車。')

add(slide_full_statement,
    'いま、AI が、\n人間をルネサンスへ\n戻す。',
    '第二のルネサンス、開幕。')

add(slide_three_columns, '3 つの哲学', [
    ('Multi-Identity', '一つに絞らない。複数の役割を、並行で生きる。'),
    ('Brain Preservation', '左脳の摩耗を防ぐ。右脳の創造を守る。'),
    ('Autonomous Agents', '道具ではなく、分身。命令ではなく、任せる。'),
], accent=CORE_GOLD, eyebrow='OUR PRINCIPLES')

add(slide_full_statement,
    '光は、分かれる。\n光は、受け止められる。',
    'そのどちらもが、CORE。', big_size=56)

add(slide_quote_big,
    'AI が仕事を奪うのではない。\n人間が、機械の真似事を辞めるためにある。',
    'CORE のスタンス。')

add(slide_full_statement,
    'あなたという人間の\nCore (核) だけを、\n世界へ抽出する。', None, big_size=64)

add(slide_section_title, 'CORE は、自分を失う装置ではない。',
    '不純物を取り除き、最も創造的な部分だけを抽出する装置。', accent=CORE_GOLD)

add(slide_full_statement,
    'あなたは、\n指揮者になる。', '楽器を演奏するのは、もう辞める。')

add(slide_quote_big, 'マルチ・アイデンティティの、\n交響曲。', 'これが CORE の世界観。')


# ───────────────────────────────────────────────────────
# Section 2 — ARCHITECTURE (10 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 2, 'アーキテクチャ', 'ARCHITECTURE', CORE_GOLD)

add(slide_full_statement,
    'CORE は、\n3 層構造で\n構築されている。', None, big_size=64)

add(slide_three_columns, '3 層', [
    ('Layer 1 — Identity', 'あなたという「核」を中心に置く、上位 OS。'),
    ('Layer 2 — Brand', '事業家のための Prism / 創作者のための Iris。'),
    ('Layer 3 — Agents', '役割ごとの専属 AI エージェント (7 + 6)。'),
], accent=CORE_GOLD, eyebrow='THE STACK')

add(slide_simple_bullets, 'BUILT ON',
    '世界最高峰の技術スタック。',
    ['Anthropic Claude (Opus / Sonnet / Haiku 切替)',
     'Google Gemini (Flash + Pro マルチプロバイダ)',
     'OpenAI TTS — 6 ボイス、自然な音声秘書',
     'Stripe — 課金、サブスク、Webhook',
     'Resend — トランザクションメール',
     'Vercel — Edge デプロイ、全世界配信'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'CROSS-PLATFORM',
    'すべての画面で、同じあなた。',
    ['Web (Chrome / Safari / Firefox / Edge)',
     'iOS PWA (iPhone / iPad、ホーム追加)',
     'macOS PWA (Mac、ホーム追加)',
     'Android PWA',
     'マルチデバイス自動同期 (5/13 〜)',
     'オフライン継続 (localStorage キャッシュ)'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'MULTI-TENANT',
    'チームでも、個人でも。',
    ['Supabase Magic Link 認証 (50,000 MAU まで無料)',
     'Row Level Security でテナント完全分離',
     'Organization 招待 + ロール (Owner / Admin / Member)',
     'Stripe Customer 1:1 マッピング',
     '個人 → チームへの移行はワンクリック',
     'SSO / SAML (Enterprise 対応予定)'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'I18N',
    '言語の壁を、超えていく。',
    ['日本語 (主)',
     'English',
     '中文 (簡体)',
     '全 LP / Dashboard / 通知に対応',
     'タイムゾーン自動検出',
     'ローカライズされた金額表記'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'PRIVACY',
    'あなたのデータは、あなたのもの。',
    ['東京リージョン (asia-northeast1) 保管',
     'パスワード不要 (Magic Link)',
     '個人情報保護法・GDPR 準拠',
     '30 日 auto-delete (古いブリーフ・履歴)',
     'ユーザー削除で全データ cascade 削除',
     '医療情報は端末内 (localStorage) のみ'],
    accent=PRISM_BLUE)

add(slide_full_statement,
    'すべては、\nあなたの「Core」を\n守るために。', None, big_size=64)

add(slide_section_title, 'では、各ブランドを見ていこう。',
    'Prism for entrepreneurs. Iris for creators.', accent=CORE_GOLD)


# ───────────────────────────────────────────────────────
# Section 3 — CORE PRISM (50 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 3, 'CORE Prism', 'FOR ENTREPRENEURS', PRISM_BLUE)

add(slide_brand_hero, 'CORE Prism', 'すべての事業家に、\n専属のエージェント AI を。', 'FOR EVERY ENTREPRENEUR', PRISM_BLUE)

add(slide_full_statement,
    'ひとつの白い光を、\n7 つの人格に分散させる。', 'それが、Prism の哲学。', big_size=62)

# 7 Agents overview
add(slide_seven_agents_overview, accent=PRISM_BLUE)

# Each agent in detail
agents_detail = [
    ('経営', 'CEO Agent',
     '事業戦略の立案、KPI モニタリング、意思決定メモの自動生成。\n\n"明日決断すべき、たった一つ" を毎朝あなたの机に置く。'),
    ('営業', 'Sales Agent',
     'リード探索から提案書ドラフト、反論対応、商談スクリプトまで。\n\n人を増やさず、売上を伸ばす。'),
    ('財務', 'CFO Agent',
     'P&L 自動生成、経費 OCR、予算配分、キャッシュ予測。\n\n会計士の役割を、AI が日次で。'),
    ('創造', 'Creative Agent',
     '画像生成、キャプション、ブランド設計、スライド自動化。\n\n表現の準備時間を、ゼロに。'),
    ('学び', 'Knowledge Agent',
     'YouTube 要約、読書ノート、知識グラフ、横断検索。\n\n蓄積したナレッジが、勝手に繋がる。'),
    ('人材', 'People Agent',
     '1on1 履歴、センチメント分析、採用面接準備、チームケア。\n\nチームの「健康診断」を、AI が。'),
    ('生活', 'Life Agent',
     '健康データ統合、スケジュール、家族の予定、心の整え。\n\n事業の決断は、身体性込みで。'),
]
for i, (jp, en, body) in enumerate(agents_detail, 1):
    add(slide_agent_card, i, jp, en, body, AGENT_COLORS[jp])

# Features
add(slide_section_title, 'Prism の主要機能。', '20 の機能を、すべて 1 つのアプリで。', accent=PRISM_BLUE)

prism_features = [
    ('PROACTIVE COACH', 'AI 戦略コーチ', '朝・昼・夜の 3 スロットで、AI が「今日の最適戦略」を自動で配信。\n\nあなたの体調・タスク・カレンダー・直近のメッセージから、最も重要な判断を 3 行に圧縮。'),
    ('MEETING INTELLIGENCE', '議事録 AI', '音声録音 → 文字起こし → 要約 → タスク抽出 → カレンダー登録、すべて自動。\n\n会議が終わった瞬間、明日の行動が決まっている。'),
    ('SLIDE STUDIO', 'スライド自動生成', '議事録や案件メモから、Beautiful.ai レベルの営業資料を 3 分で。\n\nテンプレート選択、構成、配色、すべて AI が決める。'),
    ('DOCUMENT SERIES', '書類スタジオ', '見積→発注→納品→請求 を一気通貫で生成。\n\n各書類が CRM 案件と自動連動、ステータス遷移も AI が提案。'),
    ('PEOPLE CARE', '人物ケア (1on1)', '1on1 履歴 + センチメント推移 + 可視化 + AI 改善提案。\n\nチームメンバー一人ひとりの「今」を、データで掴む。'),
    ('CRM', '案件管理', '商談ステージ・履歴・関連書類すべてを 1 つに統合。\n\nHubSpot / Salesforce との双方向同期 (Pro 以上)。'),
    ('VOICE MEMO', '音声メモ AI 振分', 'スマホで話すだけで、AI が 5 種類のフォルダに自動振り分け。\n\nアイデア / タスク / 学び / 課題 / 家族。'),
    ('YOUTUBE TO KNOWLEDGE', 'YouTube → ナレッジ', '動画 URL を貼るだけで、要約 + 重要発言 + 知識ベース登録まで完結。\n\n通勤中の視聴が、そのまま社内資産に。'),
    ('SHADOW SECRETARY', 'シャドー秘書', 'Gmail を 30 分ごとにポーリング、重要メールの返信下書きを事前に用意。\n\nあなたは "確認 → 送信" の 2 タップで返信完了。'),
    ('HEALTH DASHBOARD', 'ヘルスダッシュボード', 'Apple Health Export 取込で、HRV / 心拍 / 睡眠 / 歩数の 60 日トレンド。\n\n事業の重要決定前、健康指標で「今日の決断力」を可視化。'),
    ('SALES LEDGER', '売上台帳 + 会計連携', 'freee / マネーフォワード / 弥生と双方向 OAuth で同期。\n\n請求書発行から会計入力まで自動、確定申告も用意。'),
    ('PERSONA SWITCHING', 'ペルソナ切替', '1 タップで「経営者モード」「父親モード」「投資家モード」と文脈分離。\n\n切替後の AI は、その役割の記憶・スタイルで動く。'),
    ('VOICE SECRETARY', '音声秘書 (OpenAI TTS)', '6 種類のボイスから選択、AI の応答を自然な音声で読み上げ。\n\n運転中・移動中も、AI が "聞こえる秘書" に。'),
    ('MASTER MODE', 'マスターモード', 'API キー直叩き、Claude Opus を無制限で。\n\n品質を最優先したい時の最終モード。'),
    ('PROACTIVE PANEL', '能動提案パネル', 'あなたが何もしなくても、AI が「今あなたがすべきこと」を 3 つ提案。\n\n指示待ちから、提案先取りへ。'),
    ('IDENTITY DASHBOARD', 'アイデンティティ統合', '7 つのエージェントの状態を、1 画面でリアルタイム可視化。\n\n指揮台に立ち、オーケストラを俯瞰する感覚。'),
    ('CACHED MEMORY', 'コンテキスト記憶', 'ペルソナごとに別の長期記憶を持つ AI。\n\n「経営者モード」の文脈が「父親モード」に漏れない、完全分離。'),
    ('PWA NATIVE', 'PWA インストール', 'iPhone / Mac / Android にホーム追加可、ネイティブアプリ並みの起動。\n\nオフラインでも履歴・ナレッジは継続表示。'),
    ('INTEGRATIONS', 'Webhook 統合', 'Slack / Discord にデイリーブリーフを Webhook 配信。\n\n複数 Webhook で、複数ブランドへ同時配信。'),
    ('SECURITY', 'セキュリティ', 'localStorage + Supabase RLS でデータ完全分離。\n\n金融情報・健康情報は端末内、AI への送信時も暗号化。'),
]
for eyebrow, title, body in prism_features:
    add(slide_feature, eyebrow, title, body, PRISM_BLUE)

# Day in the life
add(slide_day_in_life, '事業家の 1 日 — with Prism.', 'A DAY IN THE LIFE', [
    ('06:30', '起床 — Apple Watch で HRV 確認', 'Prism: "HRV 良好。今日は新規提案 OK"'),
    ('07:00', '朝のブリーフ', 'CEO Agent: 今日決断すべき 1 件 + 注意 3 件'),
    ('09:30', '新規顧客との Zoom', 'Meeting AI: 議事録を自動録音 + 要約'),
    ('11:00', '見積書作成', 'Doc Studio: 過去案件から自動生成、3 分で完成'),
    ('13:00', '昼休み中に LinkedIn DM 確認', 'Shadow Secretary: 返信下書き 5 件待機'),
    ('15:00', 'スタッフの 1on1', 'People Agent: 過去 3 ヶ月のセンチメント推移を提示'),
    ('19:00', '夕方の振り返り', 'Prism: 今日の達成と、明日の優先 1 件'),
    ('22:00', '読書 — YouTube で経営者対談', 'Knowledge Agent: 自動要約 + ナレッジ化'),
], accent=PRISM_BLUE)

# Plans
add(slide_table_plans, 'CORE Prism', [
    {'name': 'Free', 'price': '¥0', 'priceLabel': '14 日間トライアル', 'tag': 'まずは試す',
     'features': ['全機能 14 日間', 'クレカ不要', '自動課金なし', '— —', '— —']},
    {'name': 'Starter', 'price': '¥4,980', 'priceLabel': '/ 月', 'tag': '個人・スタートアップ',
     'features': ['基本 AI 機能', '1 人格 / 1 ユーザー', 'ナレッジ 100 件', 'コミュニティサポート', '— —']},
    {'name': 'Standard', 'price': '¥9,800', 'priceLabel': '/ 月', 'tag': 'チームで本格活用', 'popular': True,
     'features': ['全 AI 機能 (商談 AI)', '無制限人格・ユーザー', 'ナレッジ無制限', 'OpenAI TTS 音声秘書', 'メール / Chat サポート']},
    {'name': 'Exclusive', 'price': '¥29,800', 'priceLabel': '/ 月', 'tag': '経営者・大型法人',
     'features': ['Standard 全機能', '専任 CS (1 営業日)', 'カスタム連携', '社内研修', '導入伴走']},
], accent=PRISM_BLUE)

add(slide_simple_bullets, 'WHY STANDARD',
    '⭐ なぜ Standard が選ばれるのか。',
    ['ChatGPT Plus + Notion AI + Granola + Calendly を 1 本化',
     '他社合計 ¥12,500 → Standard ¥9,800 (22% コスト削減)',
     '商談 AI は Standard 以上のみ、売上に直結',
     '音声秘書で運転中も使える',
     '無制限ペルソナで「複数事業」を 1 アプリに',
     '14 日間無料、いつでも解約'],
    accent=PRISM_BLUE)

add(slide_quote_big,
    '事業家の頭から、\n「作業」を消す。',
    'CORE Prism が約束すること。')


# ───────────────────────────────────────────────────────
# Section 4 — CORE IRIS (45 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 4, 'CORE Iris', 'FOR CREATORS', IRIS_PINK)

add(slide_brand_hero, 'CORE Iris', '数万のファンに、\nあなたの体温で届ける。', 'FOR EVERY CREATOR', IRIS_PINK)

add(slide_full_statement,
    'あなたの光が、\n世界をつくる。', 'Iris — 虹彩、瞳に宿る光のかけら。', big_size=62)

# 6 Facets overview
add(slide_six_facets_overview, accent=IRIS_PINK)

# Each facet detail
facets_detail = [
    ('案件', 'BRIEFS',
     'PR 案件のスクショを送るだけで、AI が:\n\n• ブランド名・報酬・納期を 30 秒で抽出\n• 媒体資料 + 希望ギャラを自動生成\n• 断り文 / カウンターオファー 3 案を提示', IRIS_PINK),
    ('分析', 'ANALYTICS',
     'Instagram アカウントを AI が解析:\n\n• 投稿時間の最適化\n• 反応率の高いテーマ抽出\n• フォロワー属性のクラスタリング\n• 次の 30 日で伸びる仮説 3 つ', IRIS_PURPLE),
    ('創作', 'CREATION',
     '撮影帰り、一言吹き込むだけで:\n\n• 投稿キャプション (3 案)\n• ストーリー台本 (1 分)\n• サムネ案 (画像 4 枚)\n• ハッシュタグ最適 30 個', IRIS_GOLD),
    ('交渉', 'NEGOTIATION',
     '料金交渉のロープレ + ドラフト:\n\n• 媒体資料 PDF を自動生成\n• 競合価格を AI がベンチマーク\n• 強気 / 標準 / 譲歩 の 3 トーン用意\n• 弁護士監修の利用許諾フォーマット', RGBColor(0xFF,0xA9,0x4D)),
    ('ブランド', 'BRAND',
     'あなたの世界観を、AI が学習:\n\n• フォント・色・トーンを統一\n• OG 画像・名刺・LP を一括生成\n• 撮影スタイルガイド出力\n• ブランドガイドラインの自動作成', RGBColor(0xC7,0x7D,0xFF)),
    ('仲間', 'COMMUNITY',
     '招待制クリエイターコミュニティ:\n\n• 案件シェア・コラボ募集\n• 成功事例のタイムライン\n• 月次オンライン研修 (Pro 以上)\n• オフ会 / イベント情報', RGBColor(0xFD,0x7C,0xB8)),
]
for jp, en, body, c in facets_detail:
    add(slide_feature, en, jp, body, c)

# Features
add(slide_section_title, 'Iris の主要機能。', 'クリエイターの "残酷なジレンマ" を、解く。', accent=IRIS_PINK)

iris_features = [
    ('SCREENSHOT TO DEAL', 'スクショから案件登録', 'DM のスクショを 1 枚投げるだけで、案件カードが自動完成。\n\nブランド名 / 報酬 / 納期 / 必須項目すべて AI が読み取り、入力時間 30 秒。'),
    ('VOICE TO POST', '声から投稿生成', '撮影帰りに 30 秒、一言吹き込むだけ。\n\n投稿文・ストーリー・ハッシュタグ・サムネ案を一括ドラフト。'),
    ('TRIAGE', '案件精査', '受信した依頼メールを AI が読み、優先度 / リスク / 返信案を 3 秒で提示。\n\n大量の DM の中から、真に良い案件だけを浮上させる。'),
    ('DIRECTOR', '丸投げ編集', '構成 → テロップ → キャプション → ハッシュタグ → サムネ。\n\n動画編集者に依頼するレベルの「全部」を、AI が代行。'),
    ('BEAUTY ADVISOR', '美容相談', 'スキンケア・PMS・ヘア・コスメ。\n\n肌の写真を撮って一言 → 原因仮説 + スキンケア順序 + 受診目安まで。'),
    ('INSTAGRAM ANALYZER', 'Instagram 解析', '伸びる時間帯・テーマ・タグを AI が分析。\n\n月 10 回 (Standard) 〜 無制限 (Pro 以上) で最適化サイクル。'),
    ('MEDIA KIT', 'メディアキット PDF', 'チームメンバー / 過去案件 / フォロワー属性を 1 枚にまとめた媒体資料 PDF。\n\nブランド側に送るだけで、案件単価が上がる。'),
    ('BRAND MATCH', 'ブランドマッチ', 'Prism の企業リストと連動。\n\n企業が直接 Iris クリエイターを発注、双方向マッチング。'),
    ('MULTI-ACCOUNT', 'マルチアカウント', '個人 + 副ブランド + チーム、複数 Instagram アカウントをワンタップ切替。\n\n副業ブランドの管理も 1 アプリで完結。'),
    ('HEALTH INTEGRATION', 'ヘルス連携', 'Apple Watch のデータと美容相談を統合。\n\n肌の調子・睡眠・PMS 周期と、コンテンツ生成のリズムを連動。'),
    ('COMMUNITY', 'コミュニティ', '招待制クリエイターコミュニティで、案件・コラボ・成功事例を共有。\n\n孤独になりがちなクリエイター業を、横の繋がりで支える。'),
    ('INSTAGRAM SHARE', 'Instagram 直接投稿', '生成キャプション + 画像をワンタップで Instagram アプリへ渡す。\n\nWeb Share Level 2 + URL Scheme + クリップボードの三段戦略。'),
    ('VOICE INPUT', '音声入力', '声で全機能を操作可能。\n\n手が空かない撮影現場でも、Iris は動く。'),
    ('AI MANAGER', 'AI マネージャー', '"おかえり。なんでも、話して。"\n\nIris の AI ホームは、24h あなたの専属マネージャー。'),
    ('TEMPERATURE LEARNING', '体温学習', 'あなたの過去発信・言い回し・ニュアンスを深く学習。\n\n数万のファンに "あなただけ" のメッセージが届く。'),
    ('AUTO-DM', 'DM 自動応答', 'ファンへの DM 返信を Iris が代行 (Pro 以上)。\n\n本物のあなたと区別がつかないレベルの応答品質。'),
    ('STRATEGIC ARC', '30 日プラン', '次の 30 日で「何を投稿するか」を AI が戦略立案。\n\n投稿カレンダー、テーマ、目標 KPI まで一気通貫。'),
    ('WHITE LABEL', 'ホワイトラベル (Studio)', '自社ブランドで Iris の機能を提供可能。\n\n事務所・代理店のための完全カスタマイズ。'),
    ('API ACCESS', 'API 連携', 'Salesforce / Slack / 任意の外部システムと API 連携 (Studio)。\n\n大規模事務所の業務システム統合に。'),
    ('ICON CUSTOM', '世界観カスタム', '背景・カラー・フォントすべてを、あなたの世界観に。\n\nアプリ画面そのものがブランドの延長になる。'),
]
for eyebrow, title, body in iris_features:
    add(slide_feature, eyebrow, title, body, IRIS_PINK)

# Day in the life
add(slide_day_in_life, 'クリエイターの 1 日 — with Iris.', 'A DAY IN THE LIFE', [
    ('07:00', '起床 — 肌の調子確認', 'Iris: "PMS 周期 + 睡眠不足。今日は軽めの撮影で"'),
    ('09:00', 'DM チェック — 50 件', 'Triage: 案件 3 件浮上、返信下書き完了'),
    ('11:00', '撮影現場', 'Voice: "今日は新作リップ。秋色"'),
    ('15:00', '帰宅 — 編集', 'Director: 構成 + キャプション + ハッシュタグ 一括生成'),
    ('17:00', 'PR 案件交渉', 'Negotiation: 媒体資料 + 強気価格案を 3 分で'),
    ('19:00', '美容相談', 'Beauty: "今夜は保湿強化、明日朝はビタミン C"'),
    ('21:00', '投稿 — Instagram', 'Direct Share: ワンタップで Insta アプリへ'),
    ('22:00', '振り返り', 'Iris: 今日のエンゲージメント分析 + 明日提案'),
], accent=IRIS_PINK)

# Plans
add(slide_table_plans, 'CORE Iris', [
    {'name': 'Free', 'price': '¥0', 'priceLabel': '14 日間', 'tag': 'まずは試す',
     'features': ['全機能', 'クレカ不要', '自動課金なし', '— —', '— —']},
    {'name': 'Lite', 'price': '¥1,980', 'priceLabel': '/ 月', 'tag': '入門・副業',
     'features': ['AI 戦略 30回/月', '案件管理 無制限', 'スクショ AI 10回/月', '美容 50回/月', '履歴 90 日']},
    {'name': 'Standard', 'price': '¥4,980', 'priceLabel': '/ 月', 'tag': '本気のクリエイター', 'popular': True,
     'features': ['AI 戦略 無制限', 'スクショ 無制限', '交渉文 無制限', '30 日プラン 5回', 'Instagram 解析 10回']},
    {'name': 'Pro', 'price': '¥9,800', 'priceLabel': '/ 月', 'tag': 'チーム・マネージャー',
     'features': ['Standard 全部', '連携アカ 5', 'ブランドマッチ', '自動 30 日プラン', '専任 CS']},
    {'name': 'Studio', 'price': '¥29,800', 'priceLabel': '/ 月', 'tag': '事務所・代理店',
     'features': ['Pro 全部', '連携アカ 無制限', 'ホワイトラベル', 'API 連携', '月次研修']},
], accent=IRIS_PINK)

add(slide_simple_bullets, 'WHY STANDARD',
    '⭐ なぜ Standard が一番選ばれるのか。',
    ['1 万フォロワー以上のクリエイターのスイートスポット',
     'AI 戦略 + スクショ + 交渉 が全部無制限',
     '美容相談まで含めた "横断パートナー"',
     'コミュニティ参加で横のつながりも',
     '他社分析ツールが ¥10,000 〜、Iris は ¥4,980',
     '14 日間無料、いつでも解約'],
    accent=IRIS_PINK)

add(slide_quote_big,
    '影響力を持つほど、\nファンとの距離が遠くなる。\nIris は、それを終わらせる。',
    '体温で届く、もう一人のあなた。')


# ───────────────────────────────────────────────────────
# Section 5 — CROSS-PLATFORM FEATURES (20 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 5, '共通機能', 'CROSS-PLATFORM', CORE_GOLD)

add(slide_section_title, '両ブランドで使える、共通機能。', 'どちらのプランでも、これらは含まれます。', accent=CORE_GOLD)

common = [
    ('PWA', 'ホーム画面追加', 'iPhone / Mac / Android にネイティブアプリと同じ感覚でインストール。\n\nオフラインでも履歴閲覧可、自動アップデート。'),
    ('REFERRAL', '紹介プログラム', '友人を紹介すると、お互いに 30 日トライアル延長。\n\nリンクを送るだけ、Stripe で自動精算。'),
    ('INVITE', '招待コード', 'チームメンバーや顧問を、メール 1 通で招待。\n\n受信者は Magic Link をクリックするだけで参加。'),
    ('DEMO DATA', 'デモデータ自動投入', '初回ログイン時、サンプル案件 3 件 + Welcome メッセージで Wow 体験。\n\n「使ったらすぐにできる」を実感できる導線。'),
    ('ERROR HANDLING', 'AI エラーガイド', 'quota / auth / network エラーを自動分類、3 ステップの解消手順を Pop に表示。\n\n60 秒 dismiss で同じエラーは再表示なし。'),
    ('CONTRAST', '自動コントラスト', '背景色から最適な文字色を AI が選択。\n\nどんな色を選んでも、視認性が落ちない。'),
    ('MASCOT', '精霊マスコット', 'Prism = 青いオーブ + 星、Iris = ピンクオーブ + 花。\n\nふんわり浮く、まばたきする、しゃべるとパクパク動く。'),
    ('VOICE INPUT', '音声入力', 'Web Speech API でほぼ全機能を声で操作。\n\nハンドフリーで AI と対話する未来体験。'),
    ('PERSONA', 'ペルソナ切替', 'ワンタップで「経営者モード」「父親モード」と文脈分離。\n\nAI は切替後のロールで動く、記憶は完全に独立。'),
    ('I18N', '多言語対応', '日本語 / English / 中文の 3 言語、全画面で切替。\n\nグローバル展開時もコード変更不要。'),
    ('OG IMAGES', 'SNS 共有最適化', 'X / Facebook / LINE で美しく表示される OG 画像、白背景ロゴ統一。\n\nリンクシェアが、そのままブランド露出に。'),
    ('SAFE AREA', 'iPhone 最適化', 'safe-area-inset、44px タッチ対象、自動ズーム防止、すべて対応。\n\niPhone Safari でストレスゼロ。'),
    ('NOTIFICATIONS', '通知設定', 'デイリーブリーフを Slack / Discord に Webhook 配信。\n\n複数チャネルへ同時配信、メンバー全員と同期。'),
    ('ACCOUNTING', '会計連携', 'freee / マネーフォワード / 弥生 と双方向 OAuth で同期。\n\n売上自動取込、仕訳逆送り、確定申告まで。'),
    ('HEALTH', 'Apple Health 取込', 'iPhone のヘルスケア書き出しを zip でドロップ、60 日分の指標を即可視化。\n\nデータは端末内のみ、サーバーには送信しない。'),
]
for eyebrow, title, body in common:
    add(slide_feature, eyebrow, title, body, CORE_GOLD)


# ───────────────────────────────────────────────────────
# Section 6 — DIFFERENTIATION (15 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 6, '差別化', 'WHY CORE', CORE_GOLD)

add(slide_compare_two,
    'なぜ ChatGPT ではないのか。', 'VS CHATGPT',
    'ChatGPT Plus',
    '単一文脈、能動提案なし、外部ツール統合弱い。\n\nプロンプトを毎回打つ必要があり、結局あなたが指揮を取り続ける。',
    'CORE',
    'ペルソナ毎の継続文脈、能動提案 (デイリーコーチ)、Stripe / HubSpot / Gmail 等と深い統合。\n\nあなたは "判断" だけ、それ以外は AI が走らせる。',
    accent=PRISM_BLUE)

add(slide_compare_two,
    'なぜ Notion AI ではないのか。', 'VS NOTION AI',
    'Notion AI',
    'メモ補助止まり、実行不可、健康統合なし。\n\n書く支援は得意だが、"次の一手" は教えてくれない。',
    'CORE',
    'メモ → タスク → カレンダー → 議事録 → 提案 → 実行 まで一気通貫。\n\n書いた瞬間に、AI が次の行動を提案する。',
    accent=PRISM_BLUE)

add(slide_compare_two,
    'なぜ Apple Intelligence ではないのか。', 'VS APPLE INTELLIGENCE',
    'Apple Intelligence',
    'OS 統合は強いが浅い、ペルソナ概念なし、ビジネス機能弱。\n\niPhone でしか使えない、Android ユーザーは置き去り。',
    'CORE',
    'クロスプラットフォーム (Web)、深いビジネス機能、ペルソナ毎の文脈分離。\n\nどのデバイスでも、同じあなたが続く。',
    accent=PRISM_BLUE)

add(slide_compare_two,
    'なぜ Granola ではないのか。', 'VS GRANOLA',
    'Granola',
    '議事録に特化、横展開なし。\n\n会議は記録できるが、その後の "実行" はあなた次第。',
    'CORE',
    '議事録 → タスク化 → カレンダー登録 → 振り返り の連鎖。\n\n会議が終わった瞬間、明日の行動が決まっている。',
    accent=PRISM_BLUE)

add(slide_compare_two,
    'なぜトレミルではないのか。', 'VS TOREMIL (FOR IRIS)',
    'トレミル / SINIS',
    '分析特化、"次の一手"・交渉・投稿生成できない。\n\n数字は見られるが、どう動くべきかは教えてくれない。',
    'CORE Iris',
    '分析 + 戦略 + 創作 + 交渉 + コミュニティを 1 アプリで。\n\n数字を見るのではなく、"次に何をすべきか" を提案する。',
    accent=IRIS_PINK)

add(slide_simple_bullets, 'COMPETITIVE MOAT',
    '6 つの構造的優位。',
    ['Founder 自身がマルチアイデンティティの典型ユーザー (Founder-Market Fit)',
     '絵本作家としての美術ディレクションが UI/UX に',
     'ペルソナ毎の継続文脈 (競合が真似しにくい設計思想)',
     '健康データ統合 (Apple Health/Watch/Withings)',
     'データ蓄積効果で乗り換えコストが指数関数的に増大',
     'Echo (姉妹サービス) からのクロスセル基盤'],
    accent=CORE_GOLD)

add(slide_simple_bullets, 'TECHNOLOGY EDGE',
    '技術スタックの優位。',
    ['Anthropic Claude + Google Gemini のマルチプロバイダ',
     'OpenAI TTS の 6 ボイス対応',
     'Edge Function で全世界 100ms 以下のレスポンス',
     '50,000 MAU まで完全無料の Supabase 基盤',
     'localStorage キャッシュでオフライン継続',
     'ベータ初日から PWA 完全対応'],
    accent=CORE_GOLD)

add(slide_full_statement,
    '他社が "個別ツール" なら、\nCORE は "オペレーティングシステム"。',
    '機能の数ではなく、構造で勝つ。', big_size=46)

add(slide_full_statement,
    '他社は、\nあなたが命令する道具。\n\nCORE は、\nあなたが任せる分身。',
    '道具と分身、その圧倒的な違い。', big_size=42)


# ───────────────────────────────────────────────────────
# Section 7 — PRICING & PLANS SUMMARY (10 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 7, '料金', 'PRICING', CORE_GOLD)

add(slide_section_title, 'シンプルな、5 つのプラン。', '14 日間無料、クレカ不要、いつでも解約。', accent=CORE_GOLD)

add(slide_table_plans, 'CORE Prism', [
    {'name': 'Free', 'price': '¥0', 'priceLabel': '14 日', 'tag': '試す', 'features': ['全機能', 'クレカ不要', '— —', '— —', '— —']},
    {'name': 'Starter', 'price': '¥4,980', 'priceLabel': '/ 月', 'tag': '個人',
     'features': ['基本 AI', '1 人格', 'ナレッジ 100', 'コミュニティ', '— —']},
    {'name': 'Standard', 'price': '¥9,800', 'priceLabel': '/ 月', 'tag': 'チーム本格', 'popular': True,
     'features': ['全 AI 機能', '無制限人格', '無制限ナレッジ', '音声秘書', 'メール / Chat']},
    {'name': 'Exclusive', 'price': '¥29,800', 'priceLabel': '/ 月', 'tag': '経営者',
     'features': ['Standard 全部', '専任 CS', 'カスタム連携', '導入伴走', '— —']},
], accent=PRISM_BLUE)

add(slide_table_plans, 'CORE Iris', [
    {'name': 'Free', 'price': '¥0', 'priceLabel': '14 日', 'tag': '試す', 'features': ['全機能', 'クレカ不要', '— —', '— —', '— —']},
    {'name': 'Lite', 'price': '¥1,980', 'priceLabel': '/ 月', 'tag': '入門', 'features': ['AI 30 回', 'スクショ 10 回', '美容 50 回', '履歴 90 日', '— —']},
    {'name': 'Standard', 'price': '¥4,980', 'priceLabel': '/ 月', 'tag': '本気', 'popular': True,
     'features': ['AI 無制限', 'スクショ 無制限', '交渉 無制限', '30 日プラン', 'Instagram 解析']},
    {'name': 'Pro', 'price': '¥9,800', 'priceLabel': '/ 月', 'tag': 'チーム',
     'features': ['Standard 全部', '連携 5', 'ブランドマッチ', '自動 30 日', '専任 CS']},
    {'name': 'Studio', 'price': '¥29,800', 'priceLabel': '/ 月', 'tag': '事務所',
     'features': ['Pro 全部', 'ホワイトラベル', 'API 連携', '研修', '無制限連携']},
], accent=IRIS_PINK)

add(slide_simple_bullets, 'GUIDE',
    'プラン選びのガイド。',
    ['"試してから決めたい" → Free 14 日',
     '"個人で完結したい" → Prism Starter / Iris Lite',
     '"本気で稼ぎたい" → Prism / Iris Standard ⭐',
     '"チームに導入したい" → Pro 以上',
     '"自社ブランドとして出したい" → Studio',
     '— 迷ったら Standard、9 割の方がここに着地'],
    accent=CORE_GOLD)

add(slide_simple_bullets, 'GUARANTEES',
    'CORE の 5 つの約束。',
    ['14 日間は完全無料、クレカ不要',
     '24 時間以内に体感できなければ即解約 OK',
     '解約はマイページから 1 タップ',
     '年払いで 2 ヶ月分無料',
     'データはいつでも CSV エクスポート可'],
    accent=CORE_GOLD)

add(slide_quote_big,
    '料金は、\n他に支払っていたコストを\n置き換えるためのもの。',
    '足し算ではなく、引き算で考える。')


# ───────────────────────────────────────────────────────
# Section 8 — ROADMAP (8 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 8, 'ロードマップ', 'ROADMAP', CORE_GOLD)

add(slide_section_title, '2026 年、3 つの大きな波。', 'これから 6 ヶ月で起こること。', accent=CORE_GOLD)

add(slide_simple_bullets, 'Q2 (5-7月)',
    'ベータ → 一般公開 → MRR ¥1M',
    ['5/12 ベータ同時公開 (Prism + Iris)',
     '5/13 Supabase Magic Link マルチテナント',
     '5/15 リファラル + 招待プログラム',
     '6 月 法人化 + 商標出願 (4 件)',
     '7 月 一般公開 + シードラウンド開始',
     '7 月末目標: MRR ¥1M / 有料 200 人'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'Q3 (8-10月)',
    '広告本格投下 + シード ¥30M クロージング',
    ['8 月 Google / Meta 広告本格投下',
     '8 月 法人セールス組織立ち上げ',
     '9 月 シードラウンド ¥30M クロージング',
     '9 月 デザイナー + エンジニア 2 名採用',
     '10 月 UI/UX 3.0 大幅刷新',
     '10 月末目標: MRR ¥10M / 有料 2,000 人'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'Q4 (11-1月)',
    '海外展開 + シリーズ A 準備',
    ['11 月 英語版 / 中国語版 本格展開',
     '11 月 Product Hunt #1 狙い',
     '12 月 1 周年キャンペーン',
     '1 月 US ベータ公開',
     '1 月 シリーズ A 調達準備開始',
     '1 月末目標: MRR ¥22M / 有料 5,000 人'],
    accent=PRISM_BLUE)

add(slide_simple_bullets, 'Year 2-3 (2027-2028)',
    '海外シフト + 大型調達',
    ['Q2 シリーズ A ¥3-5 億調達',
     'Q3 米国 / 中華圏セールス組織',
     'Q4 Core Family OS / Core Team OS リリース',
     '2028 末目標: ARR ¥31 億 / 有料 5 万 / 黒字化',
     '2030 末目標: ARR ¥220 億',
     'IPO 準備 (東証グロース)'],
    accent=CORE_GOLD)

add(slide_full_statement,
    'これは、\n単発のサービスではない。\n\n継続的に進化する、\nOS。', None, big_size=54)

add(slide_quote_big,
    '500 年ぶりの、\nダ・ヴィンチを\nあなたに。',
    'これが、CORE のミッション。')


# ───────────────────────────────────────────────────────
# Section 9 — CLOSING & CTA (10 枚)
# ───────────────────────────────────────────────────────
add(slide_chapter, 9, '始める', 'GET STARTED', CORE_GOLD)

add(slide_full_statement,
    '今日、\n知ってしまった以上、\nもう、戻れない。',
    'これが、CORE が一度触れた人に与える感覚。', big_size=58)

add(slide_full_statement,
    '時間を取り戻すか、\nそれとも、\nこれまで通り続けるか。',
    '選ぶのは、あなた。', big_size=54)

add(slide_simple_bullets, 'NEXT STEP',
    '今すぐ始める、3 ステップ。',
    ['1. core-prism-app.vercel.app/keynote にアクセス',
     '2. Prism / Iris、合うブランドを選ぶ',
     '3. メールアドレスを入れる → Magic Link で開始',
     '— 14 日間、すべての機能を試せます',
     '— カード登録なし、自動課金なし',
     '— 講演会限定 KEYNOTE30 で 30 日に延長'],
    accent=CORE_GOLD)

add(slide_qr_final, 'https://core-prism-app.vercel.app/keynote', code='KEYNOTE30')

add(slide_full_statement,
    'あちら側で、\nお会いしましょう。',
    '限界のない世界で。', big_size=66)

add(slide_quote_big,
    '未来のビジネスの価値は、\nいかに感性を研ぎ澄まし、\nAI と共に\n誰も見たことがない美しい世界を描けたか\nで決まる。',
    'Naoki Ide, Founder')

# 最後
s = prs.slides.add_slide(BLANK)
bg(s, WHITE)
text(s, 'CORE  Identity  OS', Inches(1.0), Inches(2.8), Inches(SLIDE_W_IN - 2.0), Inches(1.5),
     font=F_DISPLAY, size=72, color=INK, bold=True, letter_spacing=500, align='center')
rect(s, Inches(SLIDE_W_IN/2 - 1.3), Inches(4.3), Inches(2.6), Inches(0.05), fill=CORE_GOLD)
text(s, 'Thank you.', Inches(1.0), Inches(4.6), Inches(SLIDE_W_IN - 2.0), Inches(1.0),
     font=F_JP_MIN, size=28, color=INK_MUTED, italic=True, align='center')
text(s, '株式会社コアプリズム (仮)  ·  hello@core-inc.jp  ·  core-prism-app.vercel.app', Inches(1.0), Inches(5.8), Inches(SLIDE_W_IN - 2.0), Inches(0.4),
     font=F_TEXT, size=12, color=INK_FAINT, align='center')


# ========================================================================
# === 実行 ================================================================
# ========================================================================
print(f'Generating {len(slides) + 1} slides...')
TOTAL = len(slides) + 1  # +1 for the manual last slide

for i, (fn, args, kwargs) in enumerate(slides, 1):
    s = fn(*args, **kwargs)
    # Footer (skip if slide is the chapter cover - which has its own design)
    footer(s, i, TOTAL, dark=False)

# Footer for the manual last slide (Thank you)
last_slide = prs.slides[-1]
footer(last_slide, TOTAL, TOTAL, dark=False)

out = os.path.expanduser('~/Desktop/CORE_製品カタログ_2026-05-11.pptx')
prs.save(out)
print(f'✓ Saved: {out}')
print(f'  Total slides: {TOTAL}')
