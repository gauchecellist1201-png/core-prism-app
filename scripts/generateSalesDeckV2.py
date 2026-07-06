#!/usr/bin/env python3
# ============================================================
# CORE Continuum — 6サービス 営業資料 v2 (premium)
# 出力: ~/Desktop/CORE_Continuum_6サービス_営業資料.pptx
# デザイン: 深黒 × ゴールド × 製品別アクセント / Didot(英) × ヒラギノ明朝・角ゴ(和)
# 内容: 各サービス 課題→具体機能(詳細)→変化→料金 の5枚構成 + 導入編
# ============================================================
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── カラーシステム ───────────────────────────
BG        = RGBColor(0x08, 0x08, 0x0C)
BG_PANEL  = RGBColor(0x12, 0x12, 0x18)
BG_PANEL2 = RGBColor(0x17, 0x16, 0x1D)
INK       = RGBColor(0xF7, 0xF3, 0xEA)
INK_DIM   = RGBColor(0xAD, 0xA5, 0x97)
INK_FAINT = RGBColor(0x5C, 0x56, 0x4E)
GOLD      = RGBColor(0xC9, 0xA2, 0x4B)
GOLD_HI   = RGBColor(0xE8, 0xCE, 0x93)
LINE      = RGBColor(0x2A, 0x27, 0x22)
RED_SOFT  = RGBColor(0xC2, 0x5B, 0x54)

# 製品別アクセント
C_RES  = RGBColor(0x3F, 0xB6, 0xA0)   # teal — Resonance
C_GLD  = RGBColor(0x6E, 0x8B, 0xFF)   # indigo — Guild
C_IRI  = RGBColor(0xE8, 0x5B, 0x9E)   # magenta — Iris
C_LUM  = RGBColor(0xF2, 0xA0, 0x3D)   # amber — Lume
C_PRI  = RGBColor(0xB5, 0x83, 0xE6)   # violet — Prism
C_CRY  = RGBColor(0x7F, 0xB0, 0xE8)   # steel blue — Crystal

# ── フォント (macOS 標準の高級書体) ───────────────
F_DIDOT  = 'Didot'                       # 英字ラグジュアリー見出し
F_OPTIMA = 'Optima'                      # 英字ラベル
F_MINCHO = 'Hiragino Mincho ProN'        # 和文・感情的な大見出し
F_GOTHIC = 'Hiragino Kaku Gothic StdN'   # 和文・力強い見出し
F_SANS   = 'Hiragino Sans'               # 和文本文
F_NUM    = 'Avenir Next'                 # 数字

SLIDE_W = 13.333
SLIDE_H = 7.5

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]

_pages = {'n': 0, 'total': 0}


# ── 低レベルヘルパー ───────────────────────────
def _slide():
    return prs.slides.add_slide(BLANK)


def _alpha(sh, alpha):
    sp = sh.fill._xPr
    sf = sp.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            a = srgb.find(qn('a:alpha'))
            if a is None:
                a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))


def bg(slide, color=BG):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    slide.shapes._spTree.remove(sh._element)
    slide.shapes._spTree.insert(2, sh._element)
    return sh


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None, alpha=None, rounded=False):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, x, y, w, h)
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if rounded:
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    return sh


def circle(slide, cx, cy, r, fill=None, alpha=None, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, s, x, y, w, h, font=F_SANS, size=14, color=INK, bold=False,
         italic=False, align='left', anchor='top', spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    al = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT, 'center': PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    for i, ln in enumerate(s.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.italic = italic
        if spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(spacing))
    return tb


def glow(slide, cx, cy, r, color, alpha=0.10):
    circle(slide, cx, cy, r, fill=color, alpha=alpha)


def grad_bar(slide, x, y, w, h, c1=GOLD, c2=GOLD_HI):
    seg = w // 2
    rect(slide, x, y, seg, h, fill=c1)
    rect(slide, x + seg, y, w - seg, h, fill=c2)


def footer(slide, accent=GOLD, dark=False):
    _pages['n'] += 1
    col = INK_FAINT
    text(slide, 'CORE  CONTINUUM', Inches(0.62), Inches(SLIDE_H - 0.44), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=8, color=col, align='left', spacing=250)
    text(slide, f"{_pages['n']:02d}", Inches(SLIDE_W - 1.5), Inches(SLIDE_H - 0.44), Inches(0.9), Inches(0.3),
         font=F_NUM, size=9, color=col, align='right')


def eyebrow(slide, en, x=Inches(0.75), y=Inches(0.62), accent=GOLD, jp_no=None):
    """番号 + アクセント短線 + 英字ラベル"""
    cx = x
    if jp_no:
        text(slide, jp_no, cx, y - Inches(0.02), Inches(0.7), Inches(0.35),
             font=F_NUM, size=13, color=accent, bold=True, align='left')
        cx = x + Inches(0.62)
    rect(slide, cx + Inches(0.02), y + Inches(0.13), Inches(0.34), Emu(20000), fill=accent)
    text(slide, en, cx + Inches(0.46), y, Inches(8), Inches(0.35),
         font=F_OPTIMA, size=12, color=accent, bold=True, align='left', spacing=340)


def running_head(slide, name, accent):
    """製品スライド共通の右上ランニングヘッダ"""
    circle(slide, Inches(SLIDE_W - 1.62), Inches(0.78), Inches(0.05), fill=accent)
    text(slide, name.upper(), Inches(SLIDE_W - 3.4), Inches(0.62), Inches(1.6), Inches(0.32),
         font=F_DIDOT, size=13, color=INK_DIM, align='right', spacing=200)


# ══════════════════════════════════════════════
#  INTRO
# ══════════════════════════════════════════════
def slide_cover():
    s = _slide(); bg(s)
    glow(s, Inches(3.0), Inches(2.2), Inches(3.4), GOLD, 0.09)
    glow(s, Inches(11.4), Inches(6.2), Inches(2.8), GOLD_HI, 0.06)
    glow(s, Inches(10.6), Inches(1.4), Inches(1.7), C_PRI, 0.10)
    text(s, 'CORE  CONTINUUM', Inches(0.9), Inches(1.35), Inches(11), Inches(0.5),
         font=F_DIDOT, size=19, color=GOLD, bold=True, align='left', spacing=500)
    text(s, '事業のあらゆる接点に、\n途切れない知性を。', Inches(0.86), Inches(2.15), Inches(11.6), Inches(2.5),
         font=F_MINCHO, size=52, color=INK, bold=True, align='left', line_spacing=1.12)
    grad_bar(s, Inches(0.9), Inches(4.78), Inches(2.6), Inches(0.05))
    text(s, '接客・コミュニティ・SNS運用・収益化・秘書業務・顧客対応 ―― その一つひとつに、専属のAIを。',
         Inches(0.9), Inches(5.12), Inches(11.6), Inches(0.5), font=F_GOTHIC, size=16, color=INK_DIM, align='left')
    text(s, 'RESONANCE   ·   GUILD   ·   IRIS   ·   LUME   ·   PRISM   ·   CRYSTAL',
         Inches(0.9), Inches(5.95), Inches(11.6), Inches(0.4),
         font=F_OPTIMA, size=13, color=GOLD_HI, align='left', spacing=200)
    text(s, '6サービス 統合ご提案資料  ／  CORE（設立準備中）',
         Inches(0.9), Inches(6.7), Inches(11), Inches(0.35), font=F_SANS, size=11, color=INK_FAINT, align='left')
    footer(s)


def slide_problem_intro():
    s = _slide(); bg(s)
    eyebrow(s, 'THE PROBLEM')
    text(s, 'その事業、\n「人が足りない」だけで諦めていませんか。', Inches(0.75), Inches(1.15), Inches(12), Inches(1.9),
         font=F_MINCHO, size=34, color=INK, bold=True, align='left', line_spacing=1.15)
    text(s, '顧客とのやり取りは、24時間どこかで発生している。けれど人の手と時間には限りがある。その差が、静かに売上を削っていく。',
         Inches(0.75), Inches(3.0), Inches(11.6), Inches(0.7), font=F_SANS, size=14, color=INK_DIM, align='left', line_spacing=1.4)
    holes = [
        ('返信の遅れ', '営業時間外・作業中に来た問い合わせを取りこぼし、熱が冷めてから返す。'),
        ('刺さらない発信', '全員に同じ告知。開かれず、むしろブロックされる。'),
        ('活かせない情報', '会議・資料・数字が頭とフォルダに散らばり、次の一手に変わらない。'),
        ('離脱する訪問者', 'サイトやプロフィールに来ても、何も残さず去っていく。'),
        ('報われない貢献', 'コミュニティを手伝う人が可視化されず、熱量の高い人から抜ける。'),
        ('青天井の人件費', '対応量を増やすには、採用・教育・外注のコストが積み上がる。'),
    ]
    cw, ch, gx, gy = Inches(3.78), Inches(1.5), Inches(0.24), Inches(0.24)
    x0, y0 = Inches(0.75), Inches(3.85)
    for i, (t, d) in enumerate(holes):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y + Inches(0.22), Emu(34000), Inches(0.55), fill=RED_SOFT)
        text(s, t, x + Inches(0.3), y + Inches(0.22), cw - Inches(0.5), Inches(0.4),
             font=F_GOTHIC, size=16, color=INK, bold=True)
        text(s, d, x + Inches(0.3), y + Inches(0.68), cw - Inches(0.55), Inches(0.7),
             font=F_SANS, size=10.5, color=INK_DIM, line_spacing=1.25)
    footer(s)


def slide_thesis():
    s = _slide(); bg(s)
    glow(s, Inches(10.8), Inches(3.6), Inches(3.2), GOLD, 0.07)
    eyebrow(s, 'THE INSIGHT')
    text(s, '穴は6つ。ならば、\n専属のAIも6つでいい。', Inches(0.75), Inches(1.4), Inches(11.8), Inches(2.0),
         font=F_MINCHO, size=40, color=INK, bold=True, align='left', line_spacing=1.15)
    text(s, '一つの万能ツールで全部を薄くこなすのではなく、接点ごとに最適化した専門のAIを置く。\nそれぞれが深く、そして裏側でひとつの思想でつながる。それが CORE Continuum の設計です。',
         Inches(0.75), Inches(3.45), Inches(11.4), Inches(1.2), font=F_SANS, size=15, color=INK_DIM, align='left', line_spacing=1.5)
    text(s, '“ 全部入り ” ではなく、“ 全部それぞれに一流 ” を。', Inches(0.75), Inches(5.15), Inches(11.4), Inches(0.7),
         font=F_MINCHO, size=22, color=GOLD_HI, italic=True, align='left')
    grad_bar(s, Inches(0.78), Inches(6.0), Inches(2.2), Inches(0.045))
    footer(s)


PRODUCTS = [
    dict(
        key='RES', name='Resonance', accent=C_RES, no='01', en='LINE ENGAGEMENT AI',
        area='LINE公式アカウント運用',
        promise='公式LINEの返信を、AIがあなたの声で。\n月200通の無料枠を1通も使わずに。',
        metrics=[('0通', '無料枠の消費', '自動返信は返信トークン'), ('24h', '休まない応対', '深夜も作業中も'), ('1日20通', 'AI返信 (Pro)', '一人ひとりに個別最適')],
        problem_head='「返信、後回しにしていませんか。」',
        pains=[
            ('対応の遅れが失注に', '施術中・営業時間外の問い合わせを取りこぼし、熱が冷めてから返信。'),
            ('200通の壁', 'LINE公式の無料枠は月200通。使い切る不安で、告知を我慢する。'),
            ('刺さらない一斉配信', '全員に同じ文面。名前も文脈もなく、開かれずブロックされる。'),
            ('埋もれる予約の約束', '「来週の火曜で」がトークに流れ、カレンダー転記を忘れる。'),
        ],
        caps=[
            ('AI自動返信 ― 無料枠ゼロ消費', '過去8件の履歴・口調・人格コアを読み、本人になりきった返信を即送信。'),
            ('口調の自動追従', '敬語にもタメ口にも合わせ、迷えば敬語に倒して誤爆を防ぐ。'),
            ('人格コア（あなたの核）', '自己紹介文を貼るだけで人柄・目的・強み・客像を抽出し、その声で話す。'),
            ('NGワード誤送信ガード', '「完治・必ず・返金」等を含む返信は送信前に機械的に止める。'),
            ('アポ自動捕捉→カレンダー＆Meet', '会話から日時を抽出し予定を自動作成、オンラインならMeet URLも発行。'),
            ('熱量スコア', 'やり取りの頻度と鮮度から関心度を0〜100で数値化し、色分け表示。'),
            ('ファンカルテ', '会話からその人の人柄・興味・刺さる言葉・次の一手を自動でまとめる。'),
            ('返信キュー（優先度つき）', '未返信客を優先度順に並べ、上位はAIが下書き済み。3トーン即切替。'),
            ('パーソナライズ一斉配信', 'セグメント×興味タグで絞り、一人ずつ向けの案内文にAIが書き分け。'),
            ('はじめの7日間ステップ配信', '友だち追加からN日目に自動送信。これも無料枠を消費しない。'),
        ],
        before_after=[
            ('気づけず翌日返信、客が冷めて予約が流れる', '24時間、相手の口調と文脈を踏まえた返信が自動で即着信'),
            ('「全員こんにちは」の告知が開かれずブロック増', '興味タグで絞り、一人ずつ名前入りの案内をAIが用意'),
            ('「火曜で」の約束をうっかり転記し忘れる', '日時を自動抽出しGoogleカレンダー登録＋Meet発行まで'),
        ],
        diff=[
            'LINE無料枠(月200通)を1通も使わずAI自動返信を回す独自設計',
            '定型応答ではなく、あなたの口調まで憑依する人格コア',
            'AIの言い間違いを送信前に機械で止める安全ガード',
            '会話→アポ抽出→カレンダー＆Meetまで一気通貫',
        ],
        ideal='美容室・サロン・整体・教室・カフェなど、公式LINEで顧客とやり取りする店舗・個人事業主。',
        pricing=[('Free', '¥0', '名簿・トーク画面・熱量スコア'), ('Solo', '¥1,980', 'AI返信 5通/日・1アカウント'),
                 ('Pro', '¥6,980', 'AI返信 20通/日・3アカウント・配信解放'), ('Business', '¥14,800', 'AI返信 無制限・10アカウント')],
    ),
    dict(
        key='GLD', name='Guild', accent=C_GLD, no='02', en='COMMUNITY GOVERNANCE OS',
        area='DAO型コミュニティ運営OS',
        promise='決めて、残して、報いる。\n改ざんできない透明性で、コミュニティが自走する。',
        metrics=[('1/3', '可決の定足数', 'メンバーの1/3以上が投票'), ('+200', '声の採用で謝礼', '貢献が実額の価値に'), ('検証可', '改ざんゼロを', '各自のブラウザで確認')],
        problem_head='「また同じことで、揉めていませんか。」',
        pains=[
            ('決め事が流れて残らない', 'Discord/LINEで口頭で決まり、「なぜそう決めたか」が後から追えない。'),
            ('少数が仕切り、想いが埋もれる', '声の大きい人に発言力が偏り、多数決では強い想いが通らない。'),
            ('貢献が可視化されず離脱', '手伝っても報われず記録も残らず、熱量の高い人ほど冷めて抜ける。'),
            ('「言っても変わらない」不信', '要望が採用されたか握りつぶされたか分からず、運営が信頼されない。'),
        ],
        caps=[
            ('提案→クアドラティック投票', '票の強さを持ち票から選べる。単純多数決でなく「本気度」を数量で表現。'),
            ('締切で自動確定', '期限到達で自動集計し可決/不成立を確定。放置しても議題が腐らない。'),
            ('改ざん検知のハッシュ連鎖台帳', '全操作を暗号ハッシュで連結。1件でも書き換えると連鎖が壊れ検出。'),
            ('ブラウザで独立検証', 'サーバーを信用せず、各自の端末でハッシュを再計算し「改ざんゼロ」を確認。'),
            ('貢献ポイント（実データ）', '提案+5・投票+1・コメント+1・可決著者+10・完了担当+25で自動計算。'),
            ('称号レベルの解放', '持ち票残高で 仲間100→中核300→立役者700→殿堂1500 と段階解放。'),
            ('声の採用→謝礼トークン', '採用された改善提案の投稿者に、既定200の謝礼が届き台帳に記録。'),
            ('48時間の歓迎', '参加直後だけ歓迎ピンと「まず読む1件・反応できる1件」を最上部に表示。'),
            ('決定トラッキング', '可決した提案を 未着手→進行中→完了 で追跡、担当・期限を管理。'),
            ('公開ロードマップ', '届いた声を 次やる/検討中/実装済み/見送り の4列で公開、応援投票つき。'),
        ],
        before_after=[
            ('「来月の予算どうする？」が雑談に流れ誰も覚えていない', '提案化→投票→自動確定、AI書記が要約し台帳に永久記録'),
            ('手伝ってくれた人が報われず静かに離脱する', '貢献が点数・称号・順位で可視、採用された声に謝礼が届く'),
            ('「運営が裏で決めている」という不信', '全決定がハッシュ台帳に残り、各自が改ざんゼロを検証できる'),
        ],
        diff=[
            '改ざん検知の透明性 ― メンバー自身が端末で独立検証（競合になし）',
            'クアドラティック投票で「本気度」を票に反映',
            '声の採用が実額の謝礼トークンに直結',
            '不採用の判断すら台帳に残す＝握りつぶさない',
        ],
        ideal='オンラインサロン運営者・DAO・部活/サークル・習い事教室・社内の有志組織。',
        pricing=[('Free', '¥0', '15人まで・1ギルド・改ざん検知台帳'), ('Community', '¥980', '人数無制限・ギルド5つ・謝礼トークン'),
                 ('Pro', '¥3,980', 'ギルド無制限・決定トラッキング・監査出力'), ('Enterprise', '個別', 'SSO・組織改革導入・個別サポート')],
    ),
    dict(
        key='IRI', name='Iris', accent=C_IRI, no='03', en='INSTAGRAM GROWTH AI',
        area='Instagram運用支援AI',
        promise='分析から、企画・台本・案件受注まで。\nインフルエンサーの「仕事」を丸ごと引き受ける。',
        metrics=[('10 / 8種', '企画・台本テンプレ', '白紙ゼロで即生成'), ('30', '同時運用アカウント', 'Agencyプラン'), ('SRT', '字幕を自動出力', 'CapCut / Editsへ')],
        problem_head='「毎朝、白紙から悩んでいませんか。」',
        pains=[
            ('毎回ネタ切れ、白紙から企画', '何を投稿するか毎日ゼロから考える。運用者の最大の消耗ポイント。'),
            ('分析しても次の一手が不明', '数字は見えるが、それをどう行動に変えるかで手が止まる。'),
            ('代行会社は台本作成が重い', '企画者・台本ライターの人件費が利益を圧迫。属人的で量産できない。'),
            ('案件の真贋・相場が読めない', '詐欺DM・無償依頼の見極め、適正報酬の判断に時間と不安。'),
        ],
        caps=[
            ('企画テンプレ10種', '悩み解決・比較・ルーティン等をタップ→ジャンルと人物像に合わせ具体ネタへ展開。'),
            ('台本テンプレ8種', '1企画→冒頭フック3案・5〜8カットの画角/セリフ/テロップ・BGM・投稿本文まで自動。'),
            ('SRT字幕を自動出力', '台本からSRTを生成→CapCut / CapCut Web / Edits に貼るだけ。'),
            ('アカウント分析（5軸スコア）', 'バイオ/世界観/独自性/効率/商業性を0-100で採点＋推定報酬幅と30日プラン。'),
            ('最適投稿時間の提案', '曜日×5時間帯で提案。実績が無ければ「目安」と明記し、数字を捏造しない。'),
            ('9マス・グリッドプレビュー', '投稿予約を3×3で一望、ドラッグ／矢印で公開前に統一感を並べ替え。'),
            ('案件管理（受注→報告）', 'DMスクショから報酬・締切を自動抽出、詐欺リスクも自動検知。'),
            ('DM・交渉ドラフト', '4トーンで初回DM生成、相手の返信を3予測、ワンタップでDM起動。'),
            ('ブランドマッチ', '企業リストとフォロワー条件で相性を採点し、応募メールを自動生成。'),
            ('メディアキット生成', '企業送付用の美しいHTML1枚を出力。数字は実データのみ。'),
        ],
        before_after=[
            ('毎朝、白紙から30分ネタを悩む', 'テンプレ10種＋人格注入で、その人らしい企画が打鍵ゼロで即完成'),
            ('口頭・箇条書きの指示で撮影者が迷う', 'カット割り・画角・テロップ・SRT字幕まで揃い貼るだけ'),
            ('怪しいDMを手で見極め、勘で交渉', 'スクショ投入で自動抽出＋詐欺検知＋相場提示＋交渉文を即発射'),
        ],
        diff=[
            '分析→企画→台本→字幕→投稿→受注→交渉→請求を一本で繋ぐ縦統合',
            'クライアント文脈＋人格コアで「その人らしさ」を再現（代行特化）',
            '運用代行会社は台本ライターの工数をゼロに',
            '実績が無ければ「目安」と明記する誠実な数字設計',
        ],
        ideal='マイクロ〜ミドルのインフルエンサー、そして本命はSNS運用代行会社（Agency）。',
        pricing=[('Lite', '¥2,980', 'AI相談30回・案件管理無制限'), ('Standard', '¥6,980', 'AI相談無制限・IG解析月10回'),
                 ('Pro', '¥12,800', '企画/台本スタジオ・5アカウント'), ('Agency', '¥29,800', 'スタジオ無制限・30アカウント・専任')],
    ),
    dict(
        key='LUM', name='Lume', accent=C_LUM, no='04', en='LINK-IN-BIO COMMERCE',
        area='プロフィール収益化 × 継続支援',
        promise='リンクを「置くだけ」で終わらせない。\nそのプロフィールから、その場で売る。',
        metrics=[('30秒', 'で公開', '3問答えるだけ'), ('12種', 'テーマ', '切替は演出つき'), ('8週', '成長トレンド', '先週比・連続成長を表示')],
        problem_head='「フォロワーはいるのに、売上ゼロ。」',
        pains=[
            ('リンク集は売上につながらない', 'Linktreeはリンクを置くだけ。プロフィールから直接お金にできない。'),
            ('どのリンクが効いてるか不明', '無料ツールはクリック総数のみ。押されていないリンクが見えない。'),
            ('一度きりの訪問者を活かせない', '来た人のメールを集める仕組みがなく、二度と届けられない。'),
            ('数字を見ても次が分からない', '分析はあっても「改善の一手」も、続ける動機も湧かない。'),
        ],
        caps=[
            ('30秒AIオンボーディング', '表示名・活動・SNSを選ぶだけで、リンクとテーマ入りの完成ページを自動生成。'),
            ('ライブプレビュー編集', 'PCは左設定/右スマホ枠、保存も再読込もなくリアルタイムに反映。'),
            ('その場で収益化カード（4種）', '購入・予約・投げ銭・メンバー登録を、あなたの決済URLへ直結。'),
            ('メール登録フォーム→CSV', '訪問者のメールを最大5,000件集めて資産化、CSVで書き出し。'),
            ('期間限定リンク', '開始・終了日時を入れるだけで自動公開→自動非表示。貼り替えゼロ。'),
            ('クリックヒートマップ', '押された比率を熱（赤=熱い/青=冷たい）で可視化。'),
            ('流入元クロス分析', '「Instagramから来た人はEC、TikTokからはYouTube」まで分解。'),
            ('成長トレンド（8週）', '日別データを8週グラフに集計、先週比％と連続成長週を表示。'),
            ('週次AIレポート＋コーチ', '毎週の通信簿と「次の一手」を日本語で。並べ替えはワンタップ。'),
            ('選べるストリーク（罰なし）', '2/5/7/14日から継続日数を選択、途切れても責めず「また今日から」。'),
        ],
        before_after=[
            ('Linktreeにリンクを並べたが、1万人でも売上ゼロ', 'プロフィールに「購入 ¥1,500」カードを置き、その場で売れる'),
            ('どのリンクが効いてるか分からず勘で並べる', 'ヒートマップで一目、コーチが並べ替えをワンタップ提案'),
            ('告知リンクを毎回手で貼って消し忘れる', '期間限定リンクで時刻を入れるだけ、自動で公開→非表示'),
        ],
        diff=[
            'リンク集で終わらず、販売・予約・投げ銭でその場で収益化',
            'ヒートマップ＋流入元クロス分析（他社はクリック総数のみ）',
            '改善コーチ＋週次AIレポートで「次の一手」まで日本語で',
            '外部読込ゼロ・iOS標準フォントで爆速表示',
        ],
        ideal='ミュージシャン・クリエイター・フォトグラファー、サロン/教室、個人事業主・コーチ・士業。',
        pricing=[('無料', '7日間', 'リンク無制限・基本テーマ・高速表示'), ('Pro', '¥1,480', '全テーマ・ヒートマップ・独自ドメイン'),
                 ('Business', '¥3,480', 'Pro全機能＋チーム管理・EC連携・分析')],
    ),
    dict(
        key='PRI', name='Prism', accent=C_PRI, no='05', en='PERSONAL AGENT OS',
        area='パーソナルAIエージェント（専属秘書）',
        promise='あなたの資料・数字・タスクを覚えている専属秘書。\n知識を足すほど、提案が賢くなる。',
        metrics=[('14', 'CXO役員体制', '専門AIに委譲'), ('7業種', '業界別パック', 'KPI・悩み・施策を注入'), ('朝夜', '自律で先回り', '「今日の一手」を提案')],
        problem_head='「戦略を考える時間が、ありますか。」',
        pains=[
            ('雑務に追われ戦略に時間が割けない', '経営者は雑務に週12時間。時給換算で月20万円超の時間が消える。'),
            ('頭の中の情報が活かせない', '会議・資料・数字が散らばり、次の一手に変わらないまま眠る。'),
            ('AIツールは毎回文脈ゼロ', 'ChatGPTに毎回ゼロから説明。文脈が混ざり、自分事にならない。'),
            ('一人で全役割を抱える', '採用するほどではない雑務に、毎日一人で追われる。'),
        ],
        caps=[
            ('ナレッジベース自動分析', 'テキスト/URL/PDF/画像を投入→要約・洞察・戦略・アクション・リスクに構造化。'),
            ('財務データ自動抽出', '資料から売上・経費・MRR・バーンを検知し、円単位で月額換算。'),
            ('プロアクティブ提案「今日の一手」', '朝夜に巡回し、人格・タスク・知識・実データから具体アクションを提示。'),
            ('根拠チップ（嘘数字禁止）', '提案に使った知識の実件数だけを表示、0件なら演出ごと非表示。'),
            ('アクション分解', '提案文を「明日から動ける」命令形3-5個に分解、抽象論を禁止。'),
            ('会議要約（音声→ナレッジ）', 'Zoom/Meetの録音を文字起こし→決定事項/アクションに要約→知識化。'),
            ('クイックキャプチャ', '画面左下に常設、作業を止めず1-2タップで思いつきを知識に保存。'),
            ('コマンドパレット Cmd+K', '50以上のコマンドを1本に集約、最近使った操作を学習。'),
            ('人格コア＋CXO14役員', 'CEO/CFO/CMO等の専門AIに役割を割り当て、全AIに憑依。'),
            ('司令塔ループ', 'Lume→Iris→Prism→Resonanceが1本の信号を読み書きし連動。'),
        ],
        before_after=[
            ('会議録音を前に、要約する気力がない', '録音を投げるだけで決定事項とアクションに構造化（30分→3分）'),
            ('提案書を白紙から書き始める', '蓄積した知識を根拠に下書きが即完成（20分→2分）'),
            ('「あとで書こう」で気づきが消える', 'クイックキャプチャで摩擦ゼロ蓄積、知識が提案の根拠に育つ'),
        ],
        diff=[
            '知識を足すほど提案の根拠が増える自己強化ループ',
            '使用知識の件数を水増し不可能に可視化（嘘数字禁止）',
            '人格別にAIを分離し、24時間の自律実行まで進める',
            'バラバラの月額ツール群を一つに束ねる',
        ],
        ideal='一人で全部こなす個人事業主・一人社長・高単価フリーランス・1〜10名の法人代表。',
        pricing=[('Starter', '¥2,980', '3人格・商談/議事録/スライドAI'), ('Standard', '¥9,800', '7人格全員・提案書・Gmail秘書'),
                 ('Exclusive', '¥29,800', '全機能・API・チーム5名・戦略コーチ')],
    ),
    dict(
        key='CRY', name='Crystal', accent=C_CRY, no='06', en='AI CONCIERGE',
        area='AIコンシェルジュ・接客エージェント',
        promise='話しかけるだけの、24時間コンシェルジュ。\nブランドの言葉のまま、取りこぼしをゼロに。',
        metrics=[('$0.99→定額', '従量課金を回避', 'Intercom Fin比'), ('最短1分', 'で設置', 'リンクを貼るだけ'), ('多言語', '自動応対', '相手の言語で返す')],
        problem_head='「夜間の問い合わせ、逃していませんか。」',
        pains=[
            ('営業時間外の問い合わせを取りこぼす', '人が対応できない時間帯に来た訪問者が放置され、機会損失になる。'),
            ('有人チャットの人件費が重い', '世界の競合は従量課金で高額。有人対応は人件費が青天井。'),
            ('サイトに来ても離脱する', '訪問者が迷って何も残さず去っていく。'),
            ('会話が商談化・記録されない', 'リード情報が担当に届かず、有望顧客の見極めもできない。'),
        ],
        caps=[
            ('24時間・365日の応対', 'ブランドの言葉づかいのまま自動応対、深夜も休日も取りこぼしゼロ。'),
            ('声で話す・声で返す', 'マイクに話すと文字起こし→字幕表示＋音声で読み上げ。'),
            ('ナレッジ貼るだけ学習', '会社案内・料金表を貼るだけ（4,000字）。書いていないことは創作しない。'),
            ('FAQ自動生成', '貼った文章から最大6組のQ&Aを自動で起こす（原文にない情報は作らない）。'),
            ('商談・来店の日程獲得', '関心が高まると連絡先入力カードを自動で開く。'),
            ('見込み客の見極め（AI SDR）', '有望条件を設定すると、会話の流れで自然に1つずつ確認。'),
            ('会話まるごとメール通知', '名前・連絡先・会話全文（最大24往復）を要約し担当へ通知。'),
            ('多言語の自動応対', 'お客様の言語（日/英/中ほか）を自動判別し同じ言語で応対。'),
            ('先に話しかける接客', '設定秒数後にバブルの上からそっと一言、声かけから商談へ。'),
            ('設置は3通り・HTML不要', '専用リンク＋QR / タグ1行 / 設置代行。最短1分で稼働。'),
        ],
        before_after=[
            ('営業時間外は無人で放置、翌朝には熱が冷める', '深夜でもブランドの言葉で即応対し、その場で連絡先を確保'),
            ('サイト訪問者が何も残さず去っていく', '数秒後にそっと声かけ→会話→有望なら日程獲得まで自動'),
            ('チャットボット導入はHTML改修が必要', '専用リンクをInstagram/LINE/QRに貼るだけ、最短1分'),
        ],
        diff=[
            '月額定額・従量課金なし（Fin=$0.99/件、Qualified=年$68,000〜に対し）',
            '声で話す接客＋多言語（多くの競合はテキストのみ）',
            'サーバー設定ゼロ・専用リンク/QRで設置、エンジニア不要',
            '応対・見極め(SDR)・日程獲得を1つに統合',
        ],
        ideal='高級不動産・ホテル/旅館・サロン・クリニック・高単価サービス業。',
        pricing=[('Standard', '¥29,800/月', '初期¥98,000・12機能すべて'),
                 ('Luxury', '¥49,800/月', '初期¥298,000・専任チューニング＋複数サイト')],
    ),
]


# ── 製品スライド群 ───────────────────────────
def slide_divider(p):
    s = _slide(); bg(s)
    a = p['accent']
    glow(s, Inches(11.0), Inches(2.0), Inches(3.0), a, 0.12)
    glow(s, Inches(1.8), Inches(6.2), Inches(2.2), a, 0.07)
    text(s, p['no'], Inches(0.7), Inches(0.7), Inches(4), Inches(2.4),
         font=F_NUM, size=150, color=BG_PANEL2, bold=True, align='left')
    circle(s, Inches(0.98), Inches(3.55), Inches(0.07), fill=a)
    text(s, p['en'], Inches(1.2), Inches(3.38), Inches(9), Inches(0.4),
         font=F_OPTIMA, size=13, color=a, bold=True, align='left', spacing=360)
    text(s, p['name'], Inches(0.7), Inches(3.75), Inches(9), Inches(1.3),
         font=F_DIDOT, size=76, color=INK, bold=True, align='left')
    text(s, p['promise'], Inches(0.75), Inches(5.15), Inches(9.4), Inches(1.1),
         font=F_MINCHO, size=19, color=INK_DIM, align='left', line_spacing=1.3)
    # メトリクス3枚（右下）
    mx = Inches(0.75); my = Inches(6.45)
    for i, (num, lab, sub) in enumerate(p['metrics']):
        x = mx + i * Inches(2.95)
        text(s, num, x, my - Inches(0.05), Inches(2.8), Inches(0.5),
             font=F_NUM, size=26, color=a, bold=True, align='left')
        text(s, lab + ' ｜ ' + sub, x, my + Inches(0.48), Inches(2.85), Inches(0.35),
             font=F_SANS, size=9.5, color=INK_FAINT, align='left')
    footer(s, accent=a)


def slide_problem(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'THE CHALLENGE', accent=a, jp_no=p['no'])
    text(s, p['problem_head'], Inches(0.75), Inches(1.2), Inches(11.6), Inches(1.1),
         font=F_MINCHO, size=33, color=INK, bold=True, align='left')
    text(s, f"― {p['area']}が、今ぶつかっている現実。", Inches(0.78), Inches(2.35), Inches(11), Inches(0.4),
         font=F_SANS, size=13, color=INK_DIM, align='left')
    cw, ch, gx, gy = Inches(5.85), Inches(1.72), Inches(0.3), Inches(0.28)
    x0, y0 = Inches(0.75), Inches(3.0)
    for i, (t, d) in enumerate(p['pains']):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y + Inches(0.28), Emu(38000), Inches(1.16), fill=RED_SOFT)
        text(s, t, x + Inches(0.42), y + Inches(0.3), cw - Inches(0.7), Inches(0.5),
             font=F_GOTHIC, size=18, color=INK, bold=True)
        text(s, d, x + Inches(0.42), y + Inches(0.86), cw - Inches(0.8), Inches(0.7),
             font=F_SANS, size=12, color=INK_DIM, line_spacing=1.35)
    footer(s, accent=a)


def slide_caps(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'WHAT IT DOES', accent=a, jp_no=p['no'])
    text(s, '具体的に、できること。', Inches(0.75), Inches(1.16), Inches(11), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    caps = p['caps'][:10]
    n = len(caps)
    left = caps[:5]; right = caps[5:]
    col_x = [Inches(0.75), Inches(6.95)]
    col_w = Inches(5.85)
    y0 = Inches(2.05)
    row_h = Inches(0.93)
    for ci, col in enumerate([left, right]):
        for ri, (name, desc) in enumerate(col):
            x = col_x[ci]; y = y0 + ri * row_h
            circle(s, x + Inches(0.07), y + Inches(0.16), Inches(0.055), fill=a)
            text(s, name, x + Inches(0.32), y, col_w - Inches(0.35), Inches(0.36),
                 font=F_GOTHIC, size=14, color=INK, bold=True)
            text(s, desc, x + Inches(0.32), y + Inches(0.36), col_w - Inches(0.4), Inches(0.5),
                 font=F_SANS, size=10.5, color=INK_DIM, line_spacing=1.25)
    footer(s, accent=a)


def slide_transform(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'BEFORE  /  AFTER', accent=a, jp_no=p['no'])
    text(s, '導入すると、こう変わる。', Inches(0.75), Inches(1.16), Inches(11), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    # ヘッダラベル
    text(s, 'BEFORE', Inches(0.9), Inches(2.15), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=11, color=RED_SOFT, bold=True, spacing=300)
    text(s, 'AFTER', Inches(7.35), Inches(2.15), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=11, color=a, bold=True, spacing=300)
    y0 = Inches(2.6); rh = Inches(1.42)
    bw = Inches(5.75)
    for i, (bfr, aft) in enumerate(p['before_after']):
        y = y0 + i * (rh + Inches(0.12))
        # before
        rect(s, Inches(0.75), y, bw, rh, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, bfr, Inches(1.05), y, bw - Inches(0.6), rh, font=F_SANS, size=13, color=INK_DIM,
             anchor='middle', line_spacing=1.3)
        # arrow
        text(s, '→', Inches(6.55), y, Inches(0.75), rh, font=F_NUM, size=22, color=a, align='center', anchor='middle')
        # after
        rect(s, Inches(7.3), y, bw, rh, fill=a, alpha=0.12, rounded=True)
        rect(s, Inches(7.3), y, Emu(40000), rh, fill=a)
        text(s, aft, Inches(7.62), y, bw - Inches(0.6), rh, font=F_GOTHIC, size=13, color=INK, bold=True,
             anchor='middle', line_spacing=1.3)
    footer(s, accent=a)


def slide_offer(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'WHY & PRICING', accent=a, jp_no=p['no'])
    text(s, '選ばれる理由と、料金。', Inches(0.75), Inches(1.16), Inches(11), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    # 左: 差別化 + 理想顧客
    lx = Inches(0.75); lw = Inches(6.1)
    text(s, '他にはない強み', lx, Inches(2.1), lw, Inches(0.35),
         font=F_GOTHIC, size=14, color=a, bold=True)
    for i, d in enumerate(p['diff']):
        y = Inches(2.6) + i * Inches(0.72)
        circle(s, lx + Inches(0.08), y + Inches(0.15), Inches(0.05), fill=a)
        text(s, d, lx + Inches(0.32), y, lw - Inches(0.4), Inches(0.65),
             font=F_SANS, size=12.5, color=INK, line_spacing=1.3)
    iy = Inches(2.6) + len(p['diff']) * Inches(0.72) + Inches(0.15)
    rect(s, lx, iy, lw, Inches(1.0), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'こんな方に', lx + Inches(0.3), iy + Inches(0.16), lw - Inches(0.5), Inches(0.3),
         font=F_GOTHIC, size=11, color=a, bold=True)
    text(s, p['ideal'], lx + Inches(0.3), iy + Inches(0.46), lw - Inches(0.6), Inches(0.5),
         font=F_SANS, size=11.5, color=INK_DIM, line_spacing=1.3)
    # 右: 料金カード
    rx = Inches(7.2); rw = Inches(5.4)
    rect(s, rx, Inches(2.1), rw, Inches(4.55), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'PRICING', rx + Inches(0.35), Inches(2.35), rw - Inches(0.7), Inches(0.35),
         font=F_OPTIMA, size=11, color=a, bold=True, spacing=300)
    plans = p['pricing']
    py0 = Inches(2.85); prow = Inches(3.55) / max(len(plans), 1)
    for i, (plan, price, note) in enumerate(plans):
        y = py0 + prow * i
        if i > 0:
            rect(s, rx + Inches(0.35), y, rw - Inches(0.7), Emu(9000), fill=LINE)
        text(s, plan, rx + Inches(0.35), y + Inches(0.1), Inches(2.0), Inches(0.4),
             font=F_GOTHIC, size=14, color=INK, bold=True)
        text(s, price, rx + Inches(0.35), y + Inches(0.1), rw - Inches(0.7), Inches(0.4),
             font=F_NUM, size=19, color=GOLD_HI, bold=True, align='right')
        text(s, note, rx + Inches(0.35), y + Inches(0.5), rw - Inches(0.7), Inches(0.35),
             font=F_SANS, size=9.5, color=INK_DIM)
    footer(s, accent=a)


# ══════════════════════════════════════════════
#  CLOSING
# ══════════════════════════════════════════════
def slide_why_now():
    s = _slide(); bg(s)
    glow(s, Inches(10.6), Inches(3.4), Inches(3.0), GOLD, 0.07)
    eyebrow(s, 'WHY NOW')
    text(s, '世界はもう、AIエージェントに\n人件費以上を払っている。', Inches(0.75), Inches(1.2), Inches(11.8), Inches(1.7),
         font=F_MINCHO, size=32, color=INK, bold=True, align='left', line_spacing=1.15)
    cards = [
        ('Intercom  Fin', '$0.99 / 解決1件', '問い合わせを解決するたび従量課金。'),
        ('Qualified', '$68,000〜 / 年', 'AI営業担当の年間コスト。'),
        ('Sierra', 'Fortune 500 が採用', '大企業がAI接客に本気で投資。'),
    ]
    cw = Inches(3.85); ch = Inches(1.9); gx = Inches(0.28)
    x0 = Inches(0.75); y0 = Inches(3.3)
    for i, (t, big, d) in enumerate(cards):
        x = x0 + i * (cw + gx)
        rect(s, x, y0, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, t, x + Inches(0.3), y0 + Inches(0.25), cw - Inches(0.6), Inches(0.35),
             font=F_DIDOT, size=15, color=INK, bold=True)
        text(s, big, x + Inches(0.3), y0 + Inches(0.7), cw - Inches(0.6), Inches(0.5),
             font=F_NUM, size=20, color=GOLD_HI, bold=True)
        text(s, d, x + Inches(0.3), y0 + Inches(1.25), cw - Inches(0.6), Inches(0.5),
             font=F_SANS, size=10.5, color=INK_DIM, line_spacing=1.25)
    text(s, 'CORE Continuum は、同じ水準の知性を ―― 日本語ネイティブの品質で、月額定額で、6つの接点すべてに。',
         Inches(0.75), Inches(5.7), Inches(11.8), Inches(0.7), font=F_MINCHO, size=18, color=GOLD_HI, italic=True,
         align='left', line_spacing=1.3)
    footer(s)


def slide_journey():
    s = _slide(); bg(s)
    eyebrow(s, 'ONE CUSTOMER JOURNEY')
    text(s, '顧客の一生を、6つで途切れなく。', Inches(0.75), Inches(1.16), Inches(11.6), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    text(s, '出会いから、接客・販売・継続・仲間づくりまで。バラバラのツールでこぼれていた導線が、一本につながる。',
         Inches(0.75), Inches(2.0), Inches(11.8), Inches(0.5), font=F_SANS, size=13, color=INK_DIM, align='left')
    steps = [
        ('LUME', '集客・収益化', 'プロフィールで出会い、その場で売る', C_LUM),
        ('IRIS', '発信', 'SNSで見つけてもらい続ける', C_IRI),
        ('CRYSTAL', '接客', 'サイト訪問者を24時間おもてなし', C_CRY),
        ('RESONANCE', '関係を続ける', 'LINEで一人ひとりに返信し続ける', C_RES),
        ('GUILD', '仲間にする', 'ファンを共創コミュニティへ', C_GLD),
        ('PRISM', '司令塔', 'すべての情報を束ね、次の一手を出す', C_PRI),
    ]
    n = len(steps)
    cw = Inches(1.92); gap = Inches(0.08)
    x0 = Inches(0.75); y = Inches(2.95); ch = Inches(3.1)
    for i, (nm, role, d, col) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y, cw, Emu(45000), fill=col)
        circle(s, x + cw / 2, y + Inches(0.75), Inches(0.28), fill=col, alpha=0.18)
        circle(s, x + cw / 2, y + Inches(0.75), Inches(0.09), fill=col)
        text(s, nm, x + Inches(0.1), y + Inches(1.2), cw - Inches(0.2), Inches(0.35),
             font=F_DIDOT, size=13, color=INK, bold=True, align='center')
        text(s, role, x + Inches(0.1), y + Inches(1.6), cw - Inches(0.2), Inches(0.35),
             font=F_GOTHIC, size=12, color=col, bold=True, align='center')
        text(s, d, x + Inches(0.16), y + Inches(2.05), cw - Inches(0.32), Inches(0.9),
             font=F_SANS, size=9.5, color=INK_DIM, align='center', line_spacing=1.25)
    footer(s)


def slide_matrix():
    s = _slide(); bg(s)
    eyebrow(s, 'WHO NEEDS WHAT')
    text(s, '業種別・導入マップ', Inches(0.75), Inches(1.16), Inches(11), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    rows = [
        ('店舗・サロン・整体', 'Resonance ＋ Crystal', 'LINE返信の自動化とサイト接客を24時間AIが代行'),
        ('オンラインサロン・コミュニティ', 'Guild', '提案・投票・貢献の可視化で「関わりたくなる」場に'),
        ('インフルエンサー・クリエイター', 'Iris ＋ Lume', '投稿制作から収益化(販売・予約・投げ銭)まで一気通貫'),
        ('経営者・士業・フリーランス', 'Prism', '知識を蓄えるほど賢くなる専属AI秘書'),
        ('高級不動産・ホテル・接客業', 'Crystal', 'ブランドの言葉で24時間、見込み客を逃さない'),
        ('SNS運用代行会社', 'Iris (Agency)', 'クライアント別ワークスペースで複数運用を一元化'),
    ]
    x0 = Inches(0.75); y0 = Inches(2.05); rh = Inches(0.76)
    cols = [Inches(3.7), Inches(3.3), Inches(4.85)]
    tw = sum(cols, Emu(0)) if False else Inches(3.7) + Inches(3.3) + Inches(4.85)
    # header
    rect(s, x0, y0, tw, rh, fill=GOLD, alpha=0.14)
    hd = ['業種・立場', 'おすすめ', '効く理由']
    cx = x0
    for i, h in enumerate(hd):
        text(s, h, cx + Inches(0.25), y0, cols[i] - Inches(0.4), rh, font=F_GOTHIC, size=13, color=GOLD,
             bold=True, anchor='middle')
        cx += cols[i]
    for r, (a1, a2, a3) in enumerate(rows):
        y = y0 + rh * (r + 1)
        cx = x0
        vals = [a1, a2, a3]
        for i, v in enumerate(vals):
            f = F_GOTHIC if i == 1 else F_SANS
            col = INK if i < 2 else INK_DIM
            b = (i == 1)
            text(s, v, cx + Inches(0.25), y, cols[i] - Inches(0.4), rh, font=f, size=12.5 if i != 1 else 13,
                 color=GOLD_HI if i == 1 else col, bold=b, anchor='middle', line_spacing=1.2)
            cx += cols[i]
        rect(s, x0, y, tw, Emu(9000), fill=LINE)
    footer(s)


def slide_pricing_all():
    s = _slide(); bg(s)
    eyebrow(s, 'PRICING OVERVIEW')
    text(s, '全6サービス 料金一覧', Inches(0.75), Inches(1.16), Inches(11), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    rows = [
        ('サービス', 'エントリー', 'ミドル', 'ハイエンド', None),
        ('Resonance', 'Solo ¥1,980', 'Pro ¥6,980', 'Business ¥14,800', C_RES),
        ('Guild', 'Free ¥0', 'Community ¥980', 'Pro ¥3,980', C_GLD),
        ('Iris', 'Lite ¥2,980', 'Standard ¥6,980', 'Agency ¥29,800', C_IRI),
        ('Lume', '無料 (7日間)', 'Pro ¥1,480', 'Business ¥3,480', C_LUM),
        ('Prism', 'Starter ¥2,980', 'Standard ¥9,800', 'Exclusive ¥29,800', C_PRI),
        ('Crystal', '―', 'Standard ¥29,800', 'Luxury ¥49,800', C_CRY),
    ]
    x0 = Inches(0.75); y0 = Inches(2.0); rh = Inches(0.62)
    cols = [Inches(2.5), Inches(3.05), Inches(3.05), Inches(3.25)]
    tw = Inches(2.5) + Inches(3.05) + Inches(3.05) + Inches(3.25)
    for r, row in enumerate(rows):
        y = y0 + rh * r
        head = (r == 0)
        if head:
            rect(s, x0, y, tw, rh, fill=GOLD, alpha=0.14)
        cx = x0
        for i in range(4):
            v = row[i]
            if head:
                text(s, v, cx + Inches(0.22), y, cols[i] - Inches(0.35), rh, font=F_GOTHIC, size=12.5,
                     color=GOLD, bold=True, anchor='middle')
            else:
                if i == 0:
                    circle(s, cx + Inches(0.28), y + rh / 2, Inches(0.055), fill=row[4])
                    text(s, v, cx + Inches(0.48), y, cols[i] - Inches(0.55), rh, font=F_DIDOT, size=15,
                         color=INK, bold=True, anchor='middle')
                else:
                    text(s, v, cx + Inches(0.22), y, cols[i] - Inches(0.35), rh, font=F_SANS, size=12.5,
                         color=INK if i == 3 else INK_DIM, bold=(i == 3), anchor='middle')
            cx += cols[i]
        if not head:
            rect(s, x0, y, tw, Emu(9000), fill=LINE)
    text(s, '※ Crystal 以外はすべて無料または7日間トライアルからお試しいただけます（Crystalのみ初期費用あり）。表示は月額・税込。',
         Inches(0.75), Inches(6.75), Inches(11.8), Inches(0.4), font=F_SANS, size=10.5, color=INK_FAINT, align='left')
    footer(s)


def slide_start():
    s = _slide(); bg(s)
    eyebrow(s, 'HOW TO START')
    text(s, '導入は、驚くほど簡単です。', Inches(0.75), Inches(1.16), Inches(11), Inches(0.8),
         font=F_MINCHO, size=30, color=INK, bold=True, align='left')
    steps = [
        ('01', 'まず無料で試す', '無料プランまたは7日間トライアルで、実際の画面と機能をすぐ体験。クレジットカードは不要です。'),
        ('02', 'ブランドに合わせる', 'ブランド名・トーン・会社案内などを貼るだけ。その場でAIの応対があなた仕様に変わります。'),
        ('03', '貼って、すぐ稼働', 'リンクを貼る・タグを1行入れるだけ。専門知識やHTML編集は不要、最短1分で動き始めます。'),
    ]
    cw = Inches(3.85); gap = Inches(0.26); x0 = Inches(0.75); y = Inches(2.35); h = Inches(3.7)
    for i, (num, t, d) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, h, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, num, x + Inches(0.38), y + Inches(0.35), Inches(2), Inches(1.0),
             font=F_NUM, size=42, color=GOLD, bold=True)
        rect(s, x + Inches(0.4), y + Inches(1.35), Inches(0.6), Emu(24000), fill=GOLD_HI)
        text(s, t, x + Inches(0.4), y + Inches(1.55), cw - Inches(0.8), Inches(0.7),
             font=F_GOTHIC, size=20, color=INK, bold=True)
        text(s, d, x + Inches(0.4), y + Inches(2.3), cw - Inches(0.8), Inches(1.3),
             font=F_SANS, size=12, color=INK_DIM, line_spacing=1.4)
        if i < 2:
            text(s, '→', x + cw + Inches(0.01), y + h / 2 - Inches(0.3), gap, Inches(0.6),
                 font=F_NUM, size=20, color=GOLD, align='center')
    footer(s)


def slide_cta():
    s = _slide(); bg(s)
    glow(s, Inches(SLIDE_W / 2), Inches(SLIDE_H / 2), Inches(3.4), GOLD, 0.09)
    eyebrow(s, 'GET STARTED')
    text(s, 'まずは、無料でお試しください。', Inches(0.75), Inches(1.75), Inches(11.8), Inches(1.1),
         font=F_MINCHO, size=40, color=INK, bold=True, align='left')
    text(s, '気になる1つから、あるいは事業の接点に合わせて複数を組み合わせて。\nご不明点・お見積り・デモのご相談は、メール1通で。24時間以内にご返信します。',
         Inches(0.75), Inches(3.05), Inches(11.4), Inches(1.1), font=F_SANS, size=15, color=INK_DIM,
         align='left', line_spacing=1.45)
    grad_bar(s, Inches(0.78), Inches(4.5), Inches(2.4), Inches(0.05))
    rect(s, Inches(0.75), Inches(4.95), Inches(6.6), Inches(1.05), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'お問い合わせ', Inches(1.1), Inches(5.15), Inches(3), Inches(0.3),
         font=F_GOTHIC, size=11, color=GOLD, bold=True)
    text(s, 'core.guild.inc@gmail.com', Inches(1.1), Inches(5.45), Inches(6), Inches(0.45),
         font=F_NUM, size=20, color=INK, bold=True)
    footer(s)


def slide_closing():
    s = _slide(); bg(s)
    glow(s, Inches(SLIDE_W / 2), Inches(3.1), Inches(3.9), GOLD, 0.07)
    text(s, 'CORE  CONTINUUM', Inches(0.5), Inches(2.7), Inches(SLIDE_W - 1), Inches(1.0),
         font=F_DIDOT, size=46, color=INK, bold=True, align='center', spacing=300)
    text(s, '事業のあらゆる接点に、途切れない知性を。', Inches(0.5), Inches(3.85), Inches(SLIDE_W - 1), Inches(0.7),
         font=F_MINCHO, size=21, color=GOLD_HI, italic=True, align='center')
    grad_bar(s, Inches(SLIDE_W / 2 - 1.3), Inches(4.85), Inches(2.6), Inches(0.05))
    footer(s)


# ── 組み立て ───────────────────────────
def build():
    slide_cover()
    slide_problem_intro()
    slide_thesis()
    for p in PRODUCTS:
        slide_divider(p)
        slide_problem(p)
        slide_caps(p)
        slide_transform(p)
        slide_offer(p)
    slide_why_now()
    slide_journey()
    slide_matrix()
    slide_pricing_all()
    slide_start()
    slide_cta()
    slide_closing()


build()
out = os.path.expanduser('~/Desktop/CORE_Continuum_6サービス_営業資料.pptx')
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides._sldIdLst))
