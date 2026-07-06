#!/usr/bin/env python3
# ============================================================
# CORE Continuum — 6サービス 営業資料 v3 (flagship)
# 出力: ~/Desktop/CORE_Continuum_6サービス_営業資料.pptx
# 構成: intro 3 + 各製品6枚 (扉/課題/機能I/機能II/変化/料金) ×6 + closing 7 = 46枚
# デザイン: 深黒 × ゴールド × 製品アクセント / Didot × ヒラギノ明朝
#           機能はナンバリングカード、扉は製品モチーフのベクター描画
# ============================================================
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ── カラー ───────────────────────────
BG        = RGBColor(0x07, 0x07, 0x0B)
BG_PANEL  = RGBColor(0x11, 0x11, 0x17)
BG_PANEL2 = RGBColor(0x16, 0x15, 0x1C)
BG_NUM    = RGBColor(0x1F, 0x1D, 0x26)
INK       = RGBColor(0xF7, 0xF3, 0xEA)
INK_DIM   = RGBColor(0xAF, 0xA7, 0x99)
INK_FAINT = RGBColor(0x60, 0x5A, 0x51)
GOLD      = RGBColor(0xC9, 0xA2, 0x4B)
GOLD_HI   = RGBColor(0xEA, 0xD1, 0x96)
LINE      = RGBColor(0x2B, 0x28, 0x23)
RED_SOFT  = RGBColor(0xC2, 0x5B, 0x54)

C_RES = RGBColor(0x43, 0xBD, 0xA5)
C_GLD = RGBColor(0x72, 0x8F, 0xFF)
C_IRI = RGBColor(0xEA, 0x5F, 0xA2)
C_LUM = RGBColor(0xF2, 0xA0, 0x3D)
C_PRI = RGBColor(0xB8, 0x86, 0xE8)
C_CRY = RGBColor(0x82, 0xB3, 0xEA)

F_DIDOT  = 'Didot'
F_OPTIMA = 'Optima'
F_MINCHO = 'Hiragino Mincho ProN'
F_GOTHIC = 'Hiragino Kaku Gothic StdN'
F_SANS   = 'Hiragino Sans'
F_NUM    = 'Avenir Next'

SLIDE_W, SLIDE_H = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]
_pg = {'n': 0}


def _slide():
    return prs.slides.add_slide(BLANK)


def _alpha(sh, alpha):
    sp = sh.fill._xPr
    sf = sp.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            a = srgb.find(qn('a:alpha'))
            if a is None:
                a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))


def bg(slide, color=BG):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    slide.shapes._spTree.remove(sh._element)
    slide.shapes._spTree.insert(2, sh._element)


def shape(slide, st, x, y, w, h, fill=None, line=None, line_w=None, alpha=None, rot=None):
    sh = slide.shapes.add_shape(st, x, y, w, h)
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if rot is not None:
        sh.rotation = rot
    return sh


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None, alpha=None, rounded=False, radius=0.08):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = shape(slide, st, x, y, w, h, fill, line, line_w, alpha)
    if rounded:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def circle(slide, cx, cy, r, fill=None, alpha=None, line=None, line_w=None):
    return shape(slide, MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2, fill, line, line_w, alpha)


def text(slide, s, x, y, w, h, font=F_SANS, size=14, color=INK, bold=False, italic=False,
         align='left', anchor='top', spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    al = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT, 'center': PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    for i, ln in enumerate(s.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic
        if spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(spacing))
    return tb


def grad_bar(slide, x, y, w, h, c1=GOLD, c2=GOLD_HI):
    seg = w // 2
    rect(slide, x, y, seg, h, fill=c1)
    rect(slide, x + seg, y, w - seg, h, fill=c2)


def footer(slide):
    _pg['n'] += 1
    text(slide, 'CORE  CONTINUUM', Inches(0.62), Inches(SLIDE_H - 0.42), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=8, color=INK_FAINT, spacing=250)
    text(slide, f"{_pg['n']:02d}", Inches(SLIDE_W - 1.5), Inches(SLIDE_H - 0.42), Inches(0.9), Inches(0.3),
         font=F_NUM, size=9, color=INK_FAINT, align='right')


def eyebrow(slide, en, accent=GOLD, no=None, y=Inches(0.6)):
    x = Inches(0.75)
    if no:
        text(slide, no, x, y - Inches(0.02), Inches(0.6), Inches(0.35),
             font=F_NUM, size=13, color=accent, bold=True)
        x = Inches(1.36)
    rect(slide, x, y + Inches(0.12), Inches(0.34), Emu(20000), fill=accent)
    text(slide, en, x + Inches(0.46), y, Inches(8), Inches(0.35),
         font=F_OPTIMA, size=12, color=accent, bold=True, spacing=340)


def running_head(slide, name, accent):
    circle(slide, Inches(SLIDE_W - 1.55), Inches(0.77), Inches(0.05), fill=accent)
    text(slide, name.upper(), Inches(SLIDE_W - 4.2), Inches(0.6), Inches(2.45), Inches(0.32),
         font=F_DIDOT, size=13, color=INK_DIM, align='right', spacing=220)


# ── 製品モチーフ (ベクター描画) ───────────────────
def motif_res(s, a):
    """チャットバブル2つ + 返信の点"""
    x, y = Inches(9.1), Inches(1.5)
    rect(s, x, y, Inches(2.9), Inches(1.05), fill=a, alpha=0.16, rounded=True, radius=0.35)
    text(s, 'ご予約できますか？', x + Inches(0.3), y, Inches(2.4), Inches(1.05),
         font=F_SANS, size=13, color=INK_DIM, anchor='middle')
    x2, y2 = Inches(9.85), Inches(2.85)
    rect(s, x2, y2, Inches(3.05), Inches(1.5), fill=a, alpha=0.32, rounded=True, radius=0.28)
    text(s, 'ありがとうございます。\n木曜は14時と16時が\nご案内できます。', x2 + Inches(0.3), y2, Inches(2.6), Inches(1.5),
         font=F_SANS, size=11.5, color=INK, anchor='middle', line_spacing=1.25)
    text(s, 'AIが、あなたの声で下書き', Inches(9.85), Inches(4.45), Inches(3.05), Inches(0.3),
         font=F_SANS, size=9, color=a, align='right')
    for i in range(3):
        circle(s, Inches(9.35) + i * Inches(0.22), Inches(2.72), Inches(0.035), fill=a, alpha=0.7)


def motif_gld(s, a):
    """ハッシュ連鎖ブロック"""
    y = Inches(2.35)
    labels = ['提案', '投票', '決定']
    for i in range(3):
        x = Inches(8.85) + i * Inches(1.42)
        rect(s, x, y, Inches(1.12), Inches(1.12), fill=a, alpha=0.14 + 0.09 * i, rounded=True, line=a, line_w=0.75)
        text(s, labels[i], x, y + Inches(0.18), Inches(1.12), Inches(0.4),
             font=F_GOTHIC, size=12, color=INK, bold=True, align='center')
        text(s, '#' + ('a3f', '9c1', 'e57')[i], x, y + Inches(0.62), Inches(1.12), Inches(0.3),
             font=F_NUM, size=10, color=a, align='center')
        if i < 2:
            rect(s, x + Inches(1.12), y + Inches(0.52), Inches(0.30), Emu(16000), fill=a, alpha=0.8)
    text(s, 'ハッシュ連鎖 ― 書き換えれば、必ず壊れて検出', Inches(8.85), Inches(3.75), Inches(4.1), Inches(0.3),
         font=F_SANS, size=9, color=a)


def motif_iri(s, a):
    """3×3 フィードグリッド"""
    x0, y0 = Inches(9.5), Inches(1.55)
    cell = Inches(0.86); gap = Inches(0.12)
    alphas = [0.38, 0.20, 0.28, 0.16, 0.44, 0.22, 0.30, 0.18, 0.36]
    for i in range(9):
        r, c = divmod(i, 3)
        rect(s, x0 + c * (cell + gap), y0 + r * (cell + gap), cell, cell,
             fill=a, alpha=alphas[i], rounded=True, radius=0.14)
    text(s, '公開前に、フィードの統一感を', Inches(9.1), Inches(4.6), Inches(3.4), Inches(0.3),
         font=F_SANS, size=9, color=a, align='right')


def motif_lum(s, a):
    """リンクカード + 購入カード"""
    x = Inches(9.2); w = Inches(3.3)
    for i, lbl in enumerate(['新曲を聴く', 'ライブ予約']):
        y = Inches(1.55) + i * Inches(0.78)
        rect(s, x, y, w, Inches(0.62), fill=a, alpha=0.14, rounded=True, radius=0.3, line=a, line_w=0.5)
        text(s, lbl, x, y, w, Inches(0.62), font=F_SANS, size=11.5, color=INK_DIM, align='center', anchor='middle')
    y = Inches(3.2)
    rect(s, x, y, w, Inches(1.15), fill=a, alpha=0.34, rounded=True, radius=0.18)
    text(s, '楽曲データを購入', x + Inches(0.3), y + Inches(0.2), Inches(2.2), Inches(0.35),
         font=F_GOTHIC, size=12.5, color=INK, bold=True)
    text(s, '¥1,500 — その場で決済', x + Inches(0.3), y + Inches(0.6), Inches(2.7), Inches(0.35),
         font=F_NUM, size=11, color=INK, bold=True)
    text(s, 'リンクの先ではなく、この場で売る', Inches(9.2), Inches(4.55), Inches(3.3), Inches(0.3),
         font=F_SANS, size=9, color=a, align='right')


def motif_pri(s, a):
    """プリズム三角 + 分光"""
    shape(s, MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(9.7), Inches(1.9), Inches(1.9), Inches(1.7),
          fill=a, alpha=0.22, line=a, line_w=1.0)
    rect(s, Inches(8.7), Inches(2.75), Inches(1.0), Emu(16000), fill=INK_DIM, alpha=0.7)
    for i, (col, dy) in enumerate([(GOLD_HI, -0.28), (a, 0.0), (C_CRY, 0.28)]):
        bar = rect(s, Inches(11.6), Inches(2.72) + Inches(dy), Inches(1.5), Emu(16000), fill=col, alpha=0.85)
        bar.rotation = -8 + i * 8
    text(s, '1つの入力を、7つの専門知性へ', Inches(9.0), Inches(4.1), Inches(3.5), Inches(0.3),
         font=F_SANS, size=9, color=a, align='right')


def motif_cry(s, a):
    """クリスタル(ダイヤ) + 波紋"""
    circle(s, Inches(10.7), Inches(2.85), Inches(1.5), fill=None, line=a, line_w=0.5)
    circle(s, Inches(10.7), Inches(2.85), Inches(1.1), fill=None, line=a, line_w=0.75)
    shape(s, MSO_SHAPE.DIAMOND, Inches(10.05), Inches(2.2), Inches(1.3), Inches(1.3),
          fill=a, alpha=0.30, line=a, line_w=1.2)
    shape(s, MSO_SHAPE.DIAMOND, Inches(10.33), Inches(2.48), Inches(0.74), Inches(0.74),
          fill=a, alpha=0.42)
    text(s, '「いらっしゃいませ」を、24時間', Inches(9.2), Inches(4.55), Inches(3.2), Inches(0.3),
         font=F_SANS, size=9, color=a, align='right')


MOTIFS = {'RES': motif_res, 'GLD': motif_gld, 'IRI': motif_iri,
          'LUM': motif_lum, 'PRI': motif_pri, 'CRY': motif_cry}


# ══════════════════════════════════════════════
#  製品データ (コード実査に基づく事実のみ)
# ══════════════════════════════════════════════
PRODUCTS = [
    dict(
        key='RES', name='Resonance', accent=C_RES, no='01', en='LINE ENGAGEMENT AI',
        area='LINE公式アカウント運用AI',
        promise='公式LINEの返信を、AIがあなたの声で。\n月200通の無料枠は、1通も使わずに。',
        metrics=[('0通', '無料枠の消費 ― 自動返信は別枠の返信トークン'),
                 ('24h', '休まない応対 ― 深夜も、施術中も'),
                 ('7日間', '無料で全機能 ― カード登録は不要')],
        problem_head='「あとで返そう」が、\nいちばん高くついている。',
        problem_sub='LINE公式アカウントを使う店舗・教室・個人事業主が、毎日ぶつかっている現実。',
        pains=[
            ('熱いうちに返せない', '施術中・レッスン中・営業時間外に届いた問い合わせ。気づいた頃には相手の熱が冷め、予約は他店に流れている。'),
            ('月200通の壁が怖い', 'LINE公式の無料配信枠は月200通。使い切ったら課金という不安で、送りたい告知を我慢してしまう。'),
            ('一斉配信が刺さらない', '全員に同じ文面では開かれない。名前も文脈もない告知は、既読スルーとブロックを増やすだけ。'),
            ('約束がトークに埋もれる', '「来週の火曜でお願いします」。その一言がトークの奥に流れ、カレンダーへの転記を忘れて信頼を失う。'),
        ],
        caps=[
            ('AI自動返信 ― 無料枠ゼロ消費', '受信の瞬間、直近8件の会話・相手の口調・あなたの人格データを読み、本人になりきった1〜3文を即返信。返信トークン方式のため月200通の枠を1通も使わない。'),
            ('口調の自動追従', '相手が敬語なら敬語で、くだけた相手にはやわらかく、絵文字を使う人には軽く添えて。判断に迷う相手には敬語に倒し、失礼のリスクを消す。'),
            ('人格コア「あなたの核」', '自己紹介文や事業説明を貼るだけで、AIが人柄・目的・強み・客層・ゴールの5項目を抽出。以後すべての返信が「あなたの声」で書かれる。'),
            ('NGワード誤送信ガード', '「完治」「必ず」「返金」など登録した禁止語がAI文面に含まれたら、送信せず安全な受領文に自動差し替え。薬機法・景表法リスクを機械で止める。'),
            ('アポ自動捕捉', '会話に「予約・体験・明日・◯時」等の日程シグナルが出ると、約束の種類・名前・希望日時を自動で抜き出しアポ一覧へ。トークに埋もれない。'),
            ('Googleカレンダー ＋ Meet自動発行', '日時が確定すればカレンダーへ自動登録。オンライン相談なら Google Meet のURLまで自動発行して先方に案内できる。'),
            ('返信キュー ― 優先度つき', '未返信のお客様を「予約・料金・来店」など商談に近い順に自動で並べ、上位はAIが下書き済み。丁寧/フレンドリー/予約うながしの3トーンを1タップで切替。'),
            ('パーソナライズ一斉配信', '「全員・新規・活発・ごぶさた」のセグメント×興味タグで対象を絞り、1通の告知文を一人ずつ向けの案内にAIが書き分け。日時・場所などの事実は改変しない。'),
            ('熱量スコア', 'やり取りの頻度と鮮度から、お客様一人ひとりの関心度を0〜100で自動算出。高い順に色分けされ、「今日誰に声をかけるべきか」が一目で分かる。'),
            ('ファンカルテ', '会話履歴からその人の人柄(2文)・興味(最大4つ)・刺さる言葉(最大4つ)・関係を深める次の一手(最大3つ)をAIが自動でまとめる。'),
            ('はじめの7日間 ステップ配信', '友だち追加から1〜30日目の好きなタイミングに最大5通を自動送信。新規客を放置しない。これも返信トークンで、無料枠を消費しない。'),
            ('走り書き→整える', '「木曜14時なら空き。駐車場あり」の走り書きを、事実を一切足さずに送れる文面へ整形。やわらかく/きちんと/短くの3トーンから選ぶだけ。'),
        ],
        before_after=[
            ('問い合わせに気づけず翌日返信。客は冷めて他店へ', '24時間、相手の口調と会話の続きを踏まえた返信が即着信'),
            ('「全員こんにちは」の告知が既読スルーとブロックを生む', '興味タグで絞り、一人ずつ名前と文脈を踏まえた案内に書き分け'),
            ('「火曜で」の約束をトークで見失い、転記も忘れる', '日時を自動抽出→カレンダー登録→Meet発行まで自動で完了'),
        ],
        diff=[
            '自動返信・ステップ配信が月200通の無料枠を1通も使わない独自設計',
            '定型ボットではなく、あなたの口調まで再現する人格コア',
            'AIの言い間違いを送信前に止めるNGワードガード(業法リスク対策)',
            '会話→アポ抽出→カレンダー＆Meetまで、手を触れずに一気通貫',
        ],
        ideal='美容室・サロン・整体・音楽教室・カフェなど、公式LINEでお客様と1対1のやり取りが発生するすべての店舗・個人事業主。',
        pricing=[('Free', '¥0', '名簿・トーク・熱量スコア・1アカウント'),
                 ('Solo', '¥1,980/月', 'AI返信 5通/日・返信キュー・カルテ'),
                 ('Pro', '¥6,980/月', 'AI返信 20通/日・3アカウント・配信/リッチメニュー解放'),
                 ('Business', '¥14,800/月', 'AI返信 無制限・10アカウント・優先サポート')],
        price_note='全プラン7日間無料(カード不要)。いつでも解約可。',
    ),
    dict(
        key='GLD', name='Guild', accent=C_GLD, no='02', en='COMMUNITY GOVERNANCE OS',
        area='DAO型コミュニティ運営OS',
        promise='決めて、残して、報いる。\n改ざんできない透明性が、コミュニティを自走させる。',
        metrics=[('検証可', '改ざんゼロを、各メンバーの端末で確認できる'),
                 ('+200', '声が採用された人に届く謝礼トークン'),
                 ('1/3', '可決の定足数 ― 放置しても議題が腐らない')],
        problem_head='「また同じ話で、\n揉めていませんか。」',
        problem_sub='オンラインサロン・コミュニティ・サークル運営者が、毎月すり減っている現実。',
        pains=[
            ('決め事が流れて残らない', 'DiscordやLINEでなんとなく決まり、「なぜそう決めたか」が後から追えない。数ヶ月後、同じ議題がまた蒸し返される。'),
            ('声の大きい人が場を支配する', '発言力が一部に偏り、静かな多数の想いは埋もれる。多数決では「本気の1票」と「なんとなくの1票」が同じ重さになってしまう。'),
            ('貢献した人ほど、静かに去る', 'イベントを手伝い、質問に答え、場を温めた人。その貢献が何も記録されず報われないから、熱量の高い人から順に冷めて抜けていく。'),
            ('「言っても変わらない」という諦め', '出した要望が採用されたのか、握りつぶされたのか分からない。運営への信頼が積み上がらず、参加は「消費」で止まる。'),
        ],
        caps=[
            ('提案 → クアドラティック投票', '賛成・反対・棄権に加え「票の強さ」を選べる方式。強い票ほど持ち票を多く消費する仕組みで、声の大きさではなく本気度が結果に反映される。'),
            ('締切で自動集計・自動確定', '期限が来れば自動で集計し、可決／不成立を確定。可決条件は「賛成の加重が反対を上回り、かつメンバーの1/3以上が投票」。放置で議題が腐らない。'),
            ('改ざん検知のハッシュ連鎖台帳', '提案・投票・コメント・参加、すべての操作を暗号ハッシュで数珠つなぎに記録。過去の1件でも書き換えれば連鎖が壊れ、即座に検出される。'),
            ('メンバー自身のブラウザで独立検証', '台帳の生データを各自の端末に取り寄せ、その場でハッシュを再計算。運営のサーバーを一切信用せずに「改ざんなし・全件一致」を誰でも確認できる。'),
            ('貢献ポイントの自動計算', '提案+5、投票+1、コメント+1、提案が可決された作者に+10、決定をやり切った担当に+25。すべて実際の行動から自動計算され、嘘の数字は入り込めない。'),
            ('称号とプロフィール装飾の解放', '持ち票の残高で「仲間(100)→中核(300)→立役者(700)→殿堂(1500)」と称号が解放。名前の装飾も貢献25/60/120/250で段階的に輝いていく。'),
            ('声の採用 → 謝礼トークン', '「こうしたら良くなる」の声が採用されると、投稿者に既定200の謝礼トークンが届き、台帳に永久記録。見送りの判断も記録される＝握りつぶせない。'),
            ('参加48時間の歓迎設計', '新メンバーには48時間だけ、運営の歓迎メッセージと「まず読む1件・反応できる1件」をホーム最上部に表示。最初の一歩を必ず踏ませる。'),
            ('2問だけのオンボーディング', '「何をしたい?」「関心テーマは?」を選ぶだけ(入力ゼロ)。答えに合わせてホームの表示が変わり、初日から自分の居場所が見える。'),
            ('チャット・募集(LFG)', 'チャンネル制チャット(リアクション・未読バッジつき)と、人数枠・開始時刻つきの募集機能。ゲーム会も勉強会も、ワンタップで人が集まる。'),
            ('決定の実行トラッキング', '可決した提案を「未着手→進行中→完了」で追跡。担当者と期限を設定でき、完了させた人には貢献+25。「決めたのにやってない」を無くす。'),
            ('公開ロードマップ ＋ AI書記', '届いた声を「次やる/検討中/実装済み/見送り」の4列で全員に公開、応援投票つき。可決された決定はAI書記が自動要約して記録に残す。'),
        ],
        before_after=[
            ('「来月の予算どうする?」が雑談に流れ、誰も覚えていない', '提案→投票→期限で自動確定。AI書記が要約し、台帳に永久記録'),
            ('手伝ってくれる人が報われず、熱量の高い人から抜けていく', '貢献が点数・称号・順位表で見え、採用された声には謝礼が届く'),
            ('「運営が裏で決めている」という不信がくすぶる', '全決定がハッシュ台帳に残り、誰でも自分の端末で改ざんゼロを検証'),
        ],
        diff=[
            '改ざん検知×メンバー自身による独立検証 ― Skool/Circle/Discordのどれにもない',
            'クアドラティック投票で「本気度」が結果に反映される',
            '貢献が謝礼トークンという実額の価値に直結する',
            '不採用の判断まで台帳に記録＝「握りつぶし」が構造的に不可能',
        ],
        ideal='オンラインサロン・ファンコミュニティ・DAO・サークル/部活・習い事教室・社内の有志組織。「みんなで決める場」を持つすべての運営者。',
        pricing=[('Free', '¥0', '15人まで・1ギルド・台帳/AI書記つき'),
                 ('Community', '¥980/月', '人数無制限・ギルド5つ・謝礼トークン'),
                 ('Pro', '¥3,980/月', 'ギルド無制限・実行トラッキング・監査エクスポート'),
                 ('Enterprise', '個別見積', 'SSO・組織改革コンサル・個別サポート')],
        price_note='年払いは2ヶ月分おトク。まず無料の1ギルドから。',
    ),
    dict(
        key='IRI', name='Iris', accent=C_IRI, no='03', en='INSTAGRAM GROWTH AI',
        area='Instagram運用支援AI',
        promise='分析、企画、台本、字幕、案件、交渉。\nインフルエンサーの「仕事」を、丸ごと引き受ける。',
        metrics=[('10×8種', '企画・台本テンプレ ― 白紙から考える朝を終わらせる'),
                 ('30', '同時運用アカウント ― 代行会社のための量産設計'),
                 ('SRT', '字幕を自動出力 ― CapCut / Editsに貼るだけ')],
        problem_head='「今日、何を投稿しよう」を\n毎朝繰り返していないか。',
        problem_sub='クリエイターと運用代行会社が、利益を削られ続けている現実。',
        pains=[
            ('毎回ネタ切れ、白紙から企画', '投稿のたびにゼロから考える。この見えない工数が、クリエイターの時間と運用代行の利益を静かに食いつぶしている。'),
            ('分析はある。次の一手がない', 'インサイトの数字は見える。でも「で、明日何をすれば伸びるのか」に変換できず、結局いつも通りの投稿に戻る。'),
            ('台本づくりが属人化して量産できない', '代行会社はクライアントごとに「その人らしさ」を再現する必要がある。企画者とライターの人件費が、受注を増やすほど重くのしかかる。'),
            ('案件DMの真贋も相場も読めない', '届いたコラボ依頼が詐欺か本物か。報酬は適正か。判断に時間を取られ、怪しい案件に不安を抱えたまま返信している。'),
        ],
        caps=[
            ('企画テンプレ10種 ― 白紙ゼロ', '悩み解決・ビフォーアフター・◯選まとめ・よくある誤解・失敗談・比較・ルーティン・Q&A・チェックリスト等。タップするとジャンルと人物像に合わせ、具体的なネタへ自動展開。'),
            ('台本テンプレ8種 ― 撮影指示まで', '1つの企画から、冒頭フック3案・5〜8カットの画角/セリフ/テロップ/編集指示・BGM・投稿本文600字・ハッシュタグ14個・機材リストまで自動生成。撮影者が本人の指示なしで動ける密度。'),
            ('SRT字幕の自動出力', '台本の各カットから標準SRT形式の字幕を自動生成し、クリップボードへ。CapCut・CapCut Web・Edits(Meta公式)の「字幕を読み込む」に貼るだけで字幕入れが終わる。'),
            ('アカウント分析 ― 5軸スコア', 'バイオ・世界観の統一感・独自性・エンゲージメント効率・商業性を各0〜100で採点。強み弱み、競合参考、フィード/リール/ストーリー別の推定報酬幅、30日プランまで一度に出す。'),
            ('最適投稿時間の提案', '曜日×5つの時間帯(朝/昼/夕/夜/深夜)で「狙い目」を提案。実績が4投稿に満たない間は「一般的な目安」と正直に明記し、数字を捏造しない。'),
            ('9マス・グリッドプレビュー', '予約済み投稿を3×3のフィード形式で公開前に一望。ドラッグ(PC)やタップ+矢印(スマホ)で並べ替えると、実際の予約時刻が入れ替わる。世界観づくりが創作になる。'),
            ('案件管理 ― 受注から報告まで', '問い合わせ→下書き→投稿→レポートのライフサイクルを一元管理。届いたDMのスクショを1〜3枚投げるだけで、ブランド名・報酬・締切・依頼内容を自動抽出。'),
            ('詐欺リスクの自動検知', '無償依頼・過度な急かし・暗号資産/MLM・個人口座への振込・不自然な高額報酬などの危険シグナルを自動検知して警告。安心して返信できる。'),
            ('案件精査AI ― 相場つき', '案件の安全度と魅力度を各0〜100で採点し、受ける/検討/交渉/断る/ブロックを判定。日本市場の相場観(フォロワー単価)に基づく推定報酬レンジも提示。'),
            ('DM・交渉ドラフト', '丁寧/フレンドリー/プロ/熱量高めの4トーンで返信文を生成し、相手の返信パターンを3通り予測。ワンタップでInstagramのDM画面へ。'),
            ('ブランドマッチ ＋ メディアキット', '企業リストとあなたのフォロワー・ER・ジャンルの相性を0〜100で採点し、応募メールを自動生成。企業送付用の美しいメディアキット(HTML1枚)も実データのみで出力。'),
            ('ナレッジ ＋ 人格コア', '過去の生成物・メモ・投稿を蓄積し、最新の文脈を全AI呼び出しに自動注入。使うほど「その人らしい」企画・台本になっていく。代行会社の複数クライアント運用でも人格が混ざらない。'),
        ],
        before_after=[
            ('毎朝30分、白紙とにらめっこしてネタを絞り出す', 'テンプレをタップ→その人らしい企画が、打鍵ゼロで並ぶ'),
            ('企画書と口頭指示だけでは、撮影者・編集者が動けない', 'カット割り・画角・テロップ・SRT字幕まで揃って「貼るだけ」'),
            ('怪しいDMに悩み、相場も分からず勘で交渉', 'スクショを投げると抽出・詐欺判定・相場・交渉文まで即完成'),
        ],
        diff=[
            '分析→企画→台本→字幕→投稿→案件→交渉→請求の縦一気通貫(Later/Buffer/CapCutは各工程止まり)',
            'クライアント別ワークスペース×人格注入で代行の量産に耐える設計',
            '実績がない指標は「目安」と明記する誠実な数字(honest numbers)',
            '企業リスト連動のブランドマッチという自前の案件導線',
        ],
        ideal='フォロワー5千〜10万のクリエイター、そして本命はSNS運用代行会社。企画者・台本ライターの工数をゼロに近づける。',
        pricing=[('Lite', '¥2,980/月', '入門 ― AI相談30回/月・案件管理無制限'),
                 ('Standard', '¥6,980/月', '人気 ― AI相談無制限・IG解析 月10回'),
                 ('Pro', '¥12,800/月', '企画/台本スタジオ・連携5アカウント'),
                 ('Agency', '¥29,800/月', 'スタジオ無制限・30アカウント・チーム30名・専任支援')],
        price_note='7日間無料・カード登録不要・自動課金なし。',
    ),
    dict(
        key='LUM', name='Lume', accent=C_LUM, no='04', en='LINK-IN-BIO COMMERCE',
        area='プロフィール収益化ページ',
        promise='リンクを「置くだけ」で終わらせない。\nそのプロフィールから、その場で売る。',
        metrics=[('30秒', '3問に答えるだけで、ほぼ完成のページが立ち上がる'),
                 ('4種', '収益化カード ― 購入・予約・投げ銭・メンバー登録'),
                 ('5,000件', '見込み客リスト ― メールを資産化してCSVへ')],
        problem_head='フォロワー1万人。\n売上、ゼロ円。',
        problem_sub='クリエイター・店舗・個人事業主のプロフィールリンクで、毎日起きている現実。',
        pains=[
            ('リンク集は、売上につながらない', 'Linktreeやlit.linkはリンクを並べるだけ。訪問者は外へ流れていき、プロフィールそのものは1円も生まない。'),
            ('どのリンクが効いているか見えない', '無料ツールで分かるのはクリック総数だけ。一番下のリンクが実は誰にも押されていない、という事実に気づけない。'),
            ('一度きりの訪問者を、資産にできない', 'せっかく来てくれた人の連絡先を受け取る仕組みがなく、二度と届けられないまま去られていく。'),
            ('数字を見ても、次の一手が分からない', '仮に分析があっても「で、どう直せば?」は教えてくれない。改善も継続も、結局は根性頼み。'),
        ],
        caps=[
            ('30秒AIオンボーディング', '表示名→ひとこと活動→使っているSNS(9種から選択)の3問だけで、リンクとテーマが入った「ほぼ完成」のページを自動生成。白紙から作らせない。'),
            ('ライブプレビュー編集', 'PCは左に設定・右にスマホ実寸プレビューの2画面。文字を打つそばからリアルタイム反映され、保存ボタンも再読み込みも要らない。'),
            ('テーマ12種 ＋ 切替演出', 'Aurora(光るダーク)・Light・Photo(写真背景)・Sunset・Mono など12種。切り替えは「ふわっと」フェードで気持ちよく、カスタム8色×フォント4種も。'),
            ('その場で収益化カード ― 4種', '「購入」「予約」「投げ銭・応援」「メンバー登録」のカードを価格・画像つきで設置。Stripe/BASE/note等、あなた自身の決済URLへ直結し、手数料の上乗せなし。'),
            ('メール登録フォーム → CSV', '訪問者のメールアドレスを最大5,000件まで蓄積(重複は自動除去)し、CSVで書き出し。SNSのアルゴリズムに左右されない「自分のリスト」を持てる。'),
            ('期間限定リンク ― 自動公開/非表示', 'セールやイベントの告知に、開始と終了の日時を入れるだけ。時間が来れば自動で現れ、終われば自動で消える。貼り替えも消し忘れもゼロ。'),
            ('クリックヒートマップ', 'どのリンクがどれだけ押されたかを、熱(赤=熱い/青=冷たい)で直感的に可視化。「置いてあるだけのリンク」が一目で分かる。'),
            ('流入元クロス分析', '「Instagramから来た人はECを押し、TikTokから来た人はYouTubeへ行く」― 流入元×リンクの掛け合わせまで分解。流入元は9種を自動判定。'),
            ('成長トレンド ― 8週グラフ', '日別の実訪問データを8週間の推移グラフに。「先週比 +◯%」「◯週連続で伸びています」を毎週表示し、伸びている実感が続く。'),
            ('週次AIレポート「今週の通信簿」', '毎週、数字を読んだAIが見出し＋まとめ＋具体的な次の一手(1〜3個)を日本語で執筆。分析を「読む楽しみ」に変える。'),
            ('コーチのワンタップ最適化', '「よく押される順に並べ替え」「0クリックのリンクを見直し」「この時間帯に投稿を」― 提案から実行までワンタップ。'),
            ('選べるストリーク ― 罰なし', '継続の約束を2/5/7/14日から自分で選ぶ。途切れても赤字も減点もなく「また今日から」だけ。数字は実データから計算し、嘘の連続日数は表示しない。'),
        ],
        before_after=[
            ('Linktreeにリンクを並べたが、フォロワー1万人でも売上ゼロ', 'プロフィールに「購入 ¥1,500」カードを置き、その場で売れる'),
            ('どのリンクが効いているか分からず、勘で並べている', 'ヒートマップで一目。コーチが並べ替えをワンタップで提案・実行'),
            ('セール告知を毎回手で貼り、終わったら消し忘れる', '開始・終了時刻を入れるだけで、自動で現れて自動で消える'),
        ],
        diff=[
            '販売・予約・投げ銭・登録がプロフィール内で完結(競合は「リンクを貼るだけ」)',
            'ヒートマップ＋流入元クロス分析(競合はクリック総数のみ)',
            '週次AIレポートとワンタップ最適化で「次の一手」まで面倒を見る',
            '外部読み込みゼロ・iOS標準フォントの爆速表示',
        ],
        ideal='ミュージシャン・作家・フォトグラファーなどの販売系クリエイター、予約を受けたいサロン/教室、メールリストを育てたい個人事業主・士業。',
        pricing=[('無料', '¥0 (7日間)', 'リンク無制限・基本テーマ・クリック総数'),
                 ('Pro', '¥1,480/月', '全テーマ・ヒートマップ・流入元分析・独自ドメイン・AIレポート'),
                 ('Business', '¥3,480/月', 'Proの全機能＋チーム管理・EC/予約連携・チーム分析')],
        price_note='7日間無料・カード不要。定価から値下げ中(Pro 旧¥2,980)。',
    ),
    dict(
        key='PRI', name='Prism', accent=C_PRI, no='05', en='PERSONAL AGENT OS',
        area='経営者の専属AIエージェント',
        promise='あなたの資料・数字・タスクを全部覚えている専属秘書。\n知識を足すほど、提案が賢くなる。',
        metrics=[('30分→3分', '会議の議事録 ― 録音を投げるだけで構造化'),
                 ('14役', 'CXO役員AI ― 専門ごとに分業して働く'),
                 ('朝・夜', '自律巡回 ― 聞かれる前に「今日の一手」を出す')],
        problem_head='雑務に週12時間。\n戦略は、いつ考える。',
        problem_sub='一人で全部を回す経営者・士業・フリーランスの、時間の現実。',
        pains=[
            ('雑務が戦略の時間を食い尽くす', '議事録、提案書、請求書、経費。採用するほどではない仕事の山に、時給換算で月20万円分の経営者の時間が消えている。'),
            ('情報が散らばり、次の一手に変わらない', '会議メモ、PDF資料、売上の数字。フォルダと頭の中に散在したまま眠り、意思決定の材料として活きていない。'),
            ('AIに毎回、ゼロから説明している', 'ChatGPTは便利だが、あなたの事業も数字も覚えていない。毎回の文脈説明に疲れ、答えも一般論で終わる。'),
            ('自分の分身が、どこにもいない', '経営も営業も財務も広報も自分。相談相手も、先回りして動いてくれる右腕も、いない。'),
        ],
        caps=[
            ('ナレッジベース自動分析', 'テキスト・URL・PDF・画像を放り込むと、要約(3-5行)・洞察・戦略・アクション・リスクに自動で構造化。本文3万字クラスの資料まで一度に処理。'),
            ('財務データの自動抽出', '資料の中の売上・経費・キャッシュフロー・MRRを自動検知し、千円/百万円/年額の表記ゆれを月額の円に揃えて抽出。数字が意思決定の形になる。'),
            ('プロアクティブ提案「今日の一手」', '朝と夜に自律巡回し、あなたの人格・未完了タスク・蓄積知識・連携サービスの実データを織り込んで、「事業・お金・判断」に絞った具体アクションを先回りで提示。'),
            ('根拠チップ ― 嘘数字の禁止', '提案には「◯件の知識をもとに提案しました」と実際に使った知識の件数だけを表示。0件なら演出ごと非表示。水増しが構造的にできない誠実設計。'),
            ('アクション分解', '提案文を「明日から動ける」命令形のステップ3〜5個に分解し、今日/今週/今月の工数ラベルつきで提示。「意識する」のような抽象論は出さない。'),
            ('知識の先回り提案', '溜まった資料からAIが「この知識、こう活かせます」と3案を先出し。承認すれば成果物の本文まで書き上げ、1行の指示で書き直しにも応じる。'),
            ('会議要約 ― 録音を投げるだけ', 'Zoom/Meet/Teamsの録音(9形式対応)をそのまま投入→文字起こし→決定事項・アクション・リスクに構造化→ナレッジへ自動保存。30分の作業が3分に。'),
            ('クイックキャプチャ', '画面左下に常設された小さな入力口。移動中の思いつき・商談後の気づきを、作業を止めずに1〜2タップで知識化。「あとで書こう」で消える気づきを拾い切る。'),
            ('コマンドパレット Cmd+K', '議事録・スライド・交渉コーチ・CRM・P&L・請求書・経費OCRなど50超のコマンドを1本の入力に集約。よく使う操作を学習して先頭に出す。'),
            ('常設チャットDock', '画面下にいつでも呼べる対話バー。「待機」ボタンで小さな帯に畳めば、下の画面を見ながら作業できる。畳んだ状態は次回も記憶。'),
            ('CXO 14役員 ＋ 業界パック7業種', 'CEO/CFO/CMO/CTO…14の専門AI役員に仕事を委譲。飲食・美容・IT・不動産・EC・医療・教育の業界パックがKPI・悩み・打ち手・用語を最初から注入。'),
            ('司令塔ループ ― 4製品連動', 'Lume(集める)→Iris(読む)→Prism(考える)→Resonance(届ける)。CORE製品群が1本の信号でつながり、ファン一人への次の一言まで自動で連携。'),
        ],
        before_after=[
            ('会議の録音を前に、議事録を書く気力が湧かない', '録音を投げるだけで決定事項とアクションに構造化(30分→3分)'),
            ('提案書もメールも、毎回白紙から書き始める', '蓄積した知識を根拠に、下書きが数分で立ち上がる(20分→2分)'),
            ('ChatGPTに毎回、事業の文脈をゼロから説明する', '資料も数字もタスクも覚えた専属AIが、先回りで一手を出す'),
        ],
        diff=[
            '知識を足すほど提案の根拠が増える自己強化ループ(汎用AIは毎回ゼロ)',
            '使った知識の件数を可視化する根拠チップ ― 水増し不可能な誠実設計',
            '質問を待たず朝夜に自律巡回する「先回りの秘書」',
            'ChatGPT+Notion AI+管理ツールのバラバラな月額を1つに束ねる',
        ],
        ideal='一人で全部を回す経営者・一人社長・士業・高単価フリーランス、1〜10名の法人代表。「右腕が欲しいが採用するほどではない」層。',
        pricing=[('Starter', '¥2,980/月', '3人格・商談/議事録/スライドAI・Cmd+K'),
                 ('Standard', '¥9,800/月', '7人格・提案書/契約書/財務AI・Gmail秘書'),
                 ('Exclusive', '¥29,800/月', '全機能・API/Webhook・チーム5名・戦略コーチ')],
        price_note='AIの利用料込み(自前のAPIキー不要)。7日間無料。',
    ),
    dict(
        key='CRY', name='Crystal', accent=C_CRY, no='06', en='AI CONCIERGE',
        area='AIコンシェルジュ(接客エージェント)',
        promise='話しかけるだけの、24時間コンシェルジュ。\nブランドの言葉のまま、取りこぼしをゼロに。',
        metrics=[('最短1分', '設置完了 ― リンクを貼るだけ、HTML不要'),
                 ('24往復', '会話全文がそのままメールで届く'),
                 ('定額', '従量課金なし ― 世界勢は1件$0.99〜年$68,000')],
        problem_head='深夜2時の問い合わせは、\n今夜も無人のまま。',
        problem_sub='高級不動産・ホテル・サロン・高単価サービス業で、毎晩起きている機会損失。',
        pains=[
            ('営業時間外の見込み客を取りこぼす', '問い合わせの熱が最も高いのは「調べたその瞬間」。翌朝の返信では、比較検討の候補から静かに外されている。'),
            ('有人対応のコストが青天井', 'チャットの人員も電話番も、時間を延ばすほど人件費が積み上がる。世界のAI接客は従量課金で、使うほど請求が膨らむ。'),
            ('サイトに来ても、何も残さず去る', '訪問者は迷い、質問できず、離脱する。誰が来て何を知りたかったのか、痕跡すら残らない。'),
            ('ブランドの「言葉づかい」が守れない', '外注チャットや汎用ボットの応対は、せっかく磨いたブランドの世界観を一瞬で安っぽくする。'),
        ],
        caps=[
            ('24時間・365日の応対', '深夜も休日も、ブランドの言葉づかいのまま自動応対。問い合わせの取りこぼしを構造的にゼロへ。'),
            ('声で話す・声で返す', 'マイクに話しかけるだけで自動文字起こし(無音2.6秒で自動送信)。美しい字幕と合成音声で、その場で答える。iPhoneでも動作。'),
            ('ナレッジ貼るだけ学習', '会社案内・サービス説明・料金表を貼るだけ(最大4,000字)。AIはその内容を最優先の根拠に答え、書いていないことは創作しない。'),
            ('FAQ自動生成', '貼り付けた文章から、訪問者が実際に尋ねそうな質問と回答を最大6組、AIが自動で起こす(回答は2文以内・原文にない情報は作らない)。'),
            ('商談・来店の日程獲得', '会話の中で関心が高まった瞬間を検知し、連絡先の入力カードをその場でスッと差し出す。熱いうちに次のアポへ。'),
            ('見込み客の見極め ― AI SDR', '「予算月3万円以上・導入は3ヶ月以内」のような有望条件を設定すると、尋問にならないよう会話の流れで1つずつ自然に確認。'),
            ('会話まるごとメール通知', 'お名前・連絡先に加え、会話の全文(最大24往復)と有望条件との合致をまとめて、あなたのメールへ。折り返すだけで商談が始まる。'),
            ('多言語の自動応対', '英語・中国語など、お客様の言語を自動で見分けて同じ言語でお迎え。インバウンドの取りこぼしも消える。'),
            ('先に話しかける接客', '迷っている訪問者へ、設定した秒数(3〜60秒)の後にそっと一言。「話しかけられたから聞いてみた」から商談が生まれる。'),
            ('ブランド人格の調整', '正統派(高級ホテル)・親しみやすい・簡潔なプロの3人格に、呼び名・一人称・アクセント色(金/白金/紫/深紅/翡翠+自由色)まで、世界観ごと調律。'),
            ('予約ページへの橋渡し', '予約URLを設定しておくと、会話の流れに合わせて予約ボタンを自動で差し出す。応対から予約完了まで途切れない。'),
            ('設置は3通り ― HTML不要も', '①専用リンク(QRコード自動生成)を貼るだけ ②サイトにタグ1行 ③メール1通で設置代行(初期費用に同梱)。最短1分で働き始める。'),
        ],
        before_after=[
            ('営業時間外は無人。翌朝には見込み客の熱が冷めている', '深夜でもブランドの言葉で即応対し、その場で連絡先を確保'),
            ('訪問者が迷って離脱。誰が何を知りたかったかも分からない', '数秒後にそっと声かけ→会話→会話全文があなたのメールに届く'),
            ('チャットボット導入にはHTML改修とエンジニアが必要', '専用リンクをInstagram/LINE/QRに貼るだけ。最短1分で稼働'),
        ],
        diff=[
            '月額定額・従量課金なし(Intercom Fin=1件$0.99、Qualified=年約$68,000〜)',
            '声で話し声で返す接客と多言語対応(多くの競合はテキストのみ)',
            'サーバー設定ゼロ ― 専用リンク/QR/タグ1行、エンジニア不要',
            '応対・見極め(SDR)・日程獲得・通知が1つに統合されている',
        ],
        ideal='高級不動産・ホテル/旅館・美容クリニック・サロン・士業・高単価サービス業。「接客の質がブランドそのもの」の業種。',
        pricing=[('Standard', '¥29,800/月', '初期¥98,000(設置代行つき) ― 全12機能・メール通知'),
                 ('Luxury', '¥49,800/月', '初期¥298,000 ― 専任チューニング・複数サイト・優先対応')],
        price_note='購入前に、あなたのブランド設定のまま無料で試せる。',
    ),
]


# ══════════════════════════════════════════════
#  INTRO
# ══════════════════════════════════════════════
def slide_cover():
    s = _slide(); bg(s)
    circle(s, Inches(2.6), Inches(1.9), Inches(3.2), fill=GOLD, alpha=0.07)
    circle(s, Inches(11.6), Inches(6.4), Inches(2.9), fill=GOLD_HI, alpha=0.05)
    circle(s, Inches(11.1), Inches(1.3), Inches(1.6), fill=C_PRI, alpha=0.08)
    # 6色の点列 (製品の予告)
    for i, c in enumerate([C_RES, C_GLD, C_IRI, C_LUM, C_PRI, C_CRY]):
        circle(s, Inches(0.95) + i * Inches(0.34), Inches(1.05), Inches(0.05), fill=c)
    text(s, 'CORE  CONTINUUM', Inches(0.9), Inches(1.5), Inches(11), Inches(0.5),
         font=F_DIDOT, size=20, color=GOLD, bold=True, spacing=520)
    text(s, '事業のあらゆる接点に、\n途切れない知性を。', Inches(0.86), Inches(2.3), Inches(11.8), Inches(2.5),
         font=F_MINCHO, size=54, color=INK, bold=True, line_spacing=1.12)
    grad_bar(s, Inches(0.92), Inches(5.05), Inches(2.6), Inches(0.05))
    text(s, '接客・コミュニティ・SNS・収益化・秘書・顧客対応 ―― 6つの接点それぞれに、専属のAIを。',
         Inches(0.9), Inches(5.4), Inches(11.6), Inches(0.5), font=F_GOTHIC, size=15.5, color=INK_DIM)
    text(s, 'RESONANCE · GUILD · IRIS · LUME · PRISM · CRYSTAL',
         Inches(0.9), Inches(6.1), Inches(11.6), Inches(0.4),
         font=F_OPTIMA, size=13, color=GOLD_HI, spacing=220)
    text(s, '統合ご提案資料 ／ CORE（設立準備中）', Inches(0.9), Inches(6.78), Inches(11), Inches(0.35),
         font=F_SANS, size=10.5, color=INK_FAINT)
    footer(s)


def slide_problem_intro():
    s = _slide(); bg(s)
    eyebrow(s, 'THE PROBLEM')
    text(s, '売上は、6つの「すき間」から\n静かにこぼれ落ちている。', Inches(0.75), Inches(1.12), Inches(12), Inches(1.8),
         font=F_MINCHO, size=33, color=INK, bold=True, line_spacing=1.16)
    holes = [
        ('返信の遅れ', '営業時間外・作業中の問い合わせを取りこぼし、熱が冷めてから返す。'),
        ('刺さらない発信', '全員に同じ告知。開かれず、むしろブロックされていく。'),
        ('離脱する訪問者', 'サイトやプロフィールに来ても、何も残さず去っていく。'),
        ('眠る情報資産', '会議・資料・数字が散らばったまま、次の一手に変わらない。'),
        ('報われない貢献', 'コミュニティを支える人が可視化されず、熱い人から抜ける。'),
        ('青天井の人件費', '対応量を増やすには、採用・教育・外注の固定費が積み上がる。'),
    ]
    cw, ch, gx, gy = Inches(3.83), Inches(1.62), Inches(0.22), Inches(0.24)
    x0, y0 = Inches(0.75), Inches(3.15)
    for i, (t, d) in enumerate(holes):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y + Inches(0.24), Emu(34000), Inches(1.14), fill=RED_SOFT)
        text(s, f'{i+1:02d}', x + Inches(0.3), y + Inches(0.24), Inches(0.8), Inches(0.3),
             font=F_NUM, size=11, color=RED_SOFT, bold=True)
        text(s, t, x + Inches(0.3), y + Inches(0.5), cw - Inches(0.5), Inches(0.42),
             font=F_GOTHIC, size=16.5, color=INK, bold=True)
        text(s, d, x + Inches(0.3), y + Inches(0.97), cw - Inches(0.55), Inches(0.6),
             font=F_SANS, size=10.5, color=INK_DIM, line_spacing=1.28)
    text(s, 'どれも「人が足りない」せいではない。接点ごとの専属AIがいなかっただけ。',
         Inches(0.75), Inches(6.85), Inches(11.8), Inches(0.4), font=F_MINCHO, size=14, color=GOLD_HI, italic=True)
    footer(s)


def slide_thesis():
    s = _slide(); bg(s)
    circle(s, Inches(11.0), Inches(3.6), Inches(3.1), fill=GOLD, alpha=0.06)
    eyebrow(s, 'THE ANSWER')
    text(s, 'すき間は6つ。\nならば、専属のAIも6つ。', Inches(0.75), Inches(1.35), Inches(11.8), Inches(2.0),
         font=F_MINCHO, size=40, color=INK, bold=True, line_spacing=1.15)
    text(s, '一つの万能ツールですべてを薄くこなすのではなく、接点ごとに深く磨いた専門AIを置く。\nそれぞれが一流の仕事をし、裏側ではひとつの思想 ― 「数字に嘘をつかない」「ブランドの声を守る」― でつながる。',
         Inches(0.75), Inches(3.55), Inches(11.6), Inches(1.2), font=F_SANS, size=14.5, color=INK_DIM, line_spacing=1.5)
    # 6製品ミニ帯
    items = [('Resonance', 'LINE', C_RES), ('Guild', 'コミュニティ', C_GLD), ('Iris', 'Instagram', C_IRI),
             ('Lume', '収益化', C_LUM), ('Prism', '専属秘書', C_PRI), ('Crystal', '接客', C_CRY)]
    x0, y = Inches(0.75), Inches(5.15)
    w = Inches(1.92); gap = Inches(0.06)
    for i, (nm, role, c) in enumerate(items):
        x = x0 + i * (w + gap)
        rect(s, x, y, w, Inches(1.0), fill=c, alpha=0.12, rounded=True)
        rect(s, x, y, w, Emu(34000), fill=c)
        text(s, nm, x, y + Inches(0.16), w, Inches(0.35), font=F_DIDOT, size=13, color=INK, bold=True, align='center')
        text(s, role, x, y + Inches(0.55), w, Inches(0.3), font=F_SANS, size=10, color=c, align='center')
    text(s, '“ 全部入り ” ではなく、“ 全部それぞれに一流 ” を。', Inches(0.75), Inches(6.5), Inches(11.4), Inches(0.55),
         font=F_MINCHO, size=19, color=GOLD_HI, italic=True)
    footer(s)


# ══════════════════════════════════════════════
#  製品 6枚組
# ══════════════════════════════════════════════
def slide_divider(p):
    s = _slide(); bg(s)
    a = p['accent']
    circle(s, Inches(11.2), Inches(5.9), Inches(2.6), fill=a, alpha=0.06)
    circle(s, Inches(1.5), Inches(6.6), Inches(1.8), fill=a, alpha=0.05)
    # 巨大番号 (背景)
    text(s, p['no'], Inches(0.62), Inches(0.35), Inches(4), Inches(2.4),
         font=F_NUM, size=140, color=BG_NUM, bold=True)
    circle(s, Inches(0.98), Inches(3.15), Inches(0.06), fill=a)
    text(s, p['en'], Inches(1.2), Inches(2.99), Inches(8), Inches(0.4),
         font=F_OPTIMA, size=13, color=a, bold=True, spacing=380)
    text(s, p['name'], Inches(0.68), Inches(3.35), Inches(8.4), Inches(1.35),
         font=F_DIDOT, size=74, color=INK, bold=True)
    text(s, p['promise'], Inches(0.76), Inches(4.78), Inches(8.2), Inches(1.15),
         font=F_MINCHO, size=17.5, color=INK_DIM, line_spacing=1.35)
    # モチーフ (右)
    MOTIFS[p['key']](s, a)
    # 下部メトリクス
    rect(s, Inches(0.75), Inches(6.12), Inches(11.85), Emu(9000), fill=LINE)
    for i, (num, lab) in enumerate(p['metrics']):
        x = Inches(0.75) + i * Inches(4.0)
        text(s, num, x, Inches(6.3), Inches(3.8), Inches(0.5),
             font=F_NUM, size=24, color=a, bold=True)
        text(s, lab, x, Inches(6.82), Inches(3.85), Inches(0.35),
             font=F_SANS, size=9, color=INK_FAINT)
    footer(s)


def slide_problem(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'THE REALITY', accent=a, no=p['no'])
    text(s, p['problem_head'], Inches(0.75), Inches(1.18), Inches(11.6), Inches(1.6),
         font=F_MINCHO, size=31, color=INK, bold=True, line_spacing=1.18)
    text(s, p['problem_sub'], Inches(0.78), Inches(2.82), Inches(11), Inches(0.4),
         font=F_SANS, size=12.5, color=INK_DIM)
    cw, ch, gx, gy = Inches(5.85), Inches(1.72), Inches(0.3), Inches(0.26)
    x0, y0 = Inches(0.75), Inches(3.35)
    for i, (t, d) in enumerate(p['pains']):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y + Inches(0.26), Emu(38000), Inches(1.2), fill=RED_SOFT)
        text(s, t, x + Inches(0.4), y + Inches(0.26), cw - Inches(0.7), Inches(0.42),
             font=F_GOTHIC, size=16, color=INK, bold=True)
        text(s, d, x + Inches(0.4), y + Inches(0.74), cw - Inches(0.75), Inches(0.9),
             font=F_SANS, size=10.8, color=INK_DIM, line_spacing=1.32)
    footer(s)


def _caps_slide(p, caps, part, start_no):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, f'WHAT IT DOES — {part}', accent=a, no=p['no'])
    title = '具体的に、できること。' if part == 'I' else 'まだ、ある。'
    text(s, title, Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    cw, ch = Inches(5.88), Inches(1.52)
    gx, gy = Inches(0.26), Inches(0.18)
    x0, y0 = Inches(0.75), Inches(2.05)
    for i, (name, desc) in enumerate(caps):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        # 番号バッジ
        rect(s, x + Inches(0.22), y + Inches(0.22), Inches(0.52), Inches(0.42),
             fill=a, alpha=0.16, rounded=True, radius=0.25)
        text(s, f'{start_no + i:02d}', x + Inches(0.22), y + Inches(0.24), Inches(0.52), Inches(0.38),
             font=F_NUM, size=13, color=a, bold=True, align='center', anchor='middle')
        text(s, name, x + Inches(0.9), y + Inches(0.2), cw - Inches(1.1), Inches(0.42),
             font=F_GOTHIC, size=13.5, color=INK, bold=True)
        text(s, desc, x + Inches(0.9), y + Inches(0.62), cw - Inches(1.16), Inches(0.85),
             font=F_SANS, size=9.6, color=INK_DIM, line_spacing=1.26)
    footer(s)


def slide_caps1(p):
    _caps_slide(p, p['caps'][:6], 'I', 1)


def slide_caps2(p):
    _caps_slide(p, p['caps'][6:12], 'II', 7)


def slide_transform(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'BEFORE  /  AFTER', accent=a, no=p['no'])
    text(s, '導入すると、日常がこう変わる。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    text(s, 'BEFORE ― いままで', Inches(0.9), Inches(2.12), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=11, color=RED_SOFT, bold=True, spacing=200)
    text(s, 'AFTER ― これから', Inches(7.5), Inches(2.12), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=11, color=a, bold=True, spacing=200)
    y0 = Inches(2.56); rh = Inches(1.36)
    bw = Inches(5.72)
    for i, (bfr, aft) in enumerate(p['before_after']):
        y = y0 + i * (rh + Inches(0.14))
        rect(s, Inches(0.75), y, bw, rh, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, bfr, Inches(1.08), y, bw - Inches(0.65), rh, font=F_SANS, size=12.5, color=INK_DIM,
             anchor='middle', line_spacing=1.32)
        text(s, '→', Inches(6.52), y, Inches(0.85), rh, font=F_NUM, size=22, color=a,
             align='center', anchor='middle')
        rect(s, Inches(7.38), y, bw, rh, fill=a, alpha=0.13, rounded=True)
        rect(s, Inches(7.38), y, Emu(40000), rh, fill=a)
        text(s, aft, Inches(7.72), y, bw - Inches(0.65), rh, font=F_GOTHIC, size=12.5, color=INK,
             bold=True, anchor='middle', line_spacing=1.32)
    footer(s)


def slide_offer(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p['name'], a)
    eyebrow(s, 'WHY  &  PRICING', accent=a, no=p['no'])
    text(s, '選ばれる理由と、料金。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    lx, lw = Inches(0.75), Inches(6.15)
    text(s, '他にはない強み', lx, Inches(2.06), lw, Inches(0.35),
         font=F_GOTHIC, size=13.5, color=a, bold=True)
    for i, d in enumerate(p['diff']):
        y = Inches(2.52) + i * Inches(0.74)
        circle(s, lx + Inches(0.08), y + Inches(0.14), Inches(0.05), fill=a)
        text(s, d, lx + Inches(0.32), y, lw - Inches(0.4), Inches(0.68),
             font=F_SANS, size=12, color=INK, line_spacing=1.3)
    iy = Inches(2.52) + len(p['diff']) * Inches(0.74) + Inches(0.12)
    rect(s, lx, iy, lw, Inches(1.15), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'こんな方に', lx + Inches(0.3), iy + Inches(0.16), lw - Inches(0.5), Inches(0.3),
         font=F_GOTHIC, size=10.5, color=a, bold=True)
    text(s, p['ideal'], lx + Inches(0.3), iy + Inches(0.46), lw - Inches(0.6), Inches(0.65),
         font=F_SANS, size=11, color=INK_DIM, line_spacing=1.3)
    # 料金
    rx, rw = Inches(7.25), Inches(5.35)
    rect(s, rx, Inches(2.06), rw, Inches(4.55), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'PRICING', rx + Inches(0.35), Inches(2.3), rw - Inches(0.7), Inches(0.32),
         font=F_OPTIMA, size=11, color=a, bold=True, spacing=300)
    plans = p['pricing']
    py0 = Inches(2.74); prow = Inches(3.35) / max(len(plans), 1)
    for i, (plan, price, note) in enumerate(plans):
        y = py0 + prow * i
        if i > 0:
            rect(s, rx + Inches(0.35), y, rw - Inches(0.7), Emu(9000), fill=LINE)
        text(s, plan, rx + Inches(0.35), y + Inches(0.09), Inches(2.4), Inches(0.4),
             font=F_GOTHIC, size=13.5, color=INK, bold=True)
        text(s, price, rx + Inches(0.35), y + Inches(0.07), rw - Inches(0.7), Inches(0.42),
             font=F_NUM, size=17, color=GOLD_HI, bold=True, align='right')
        text(s, note, rx + Inches(0.35), y + Inches(0.5), rw - Inches(0.7), Inches(0.35),
             font=F_SANS, size=9, color=INK_DIM)
    text(s, p['price_note'], rx + Inches(0.35), Inches(6.18), rw - Inches(0.7), Inches(0.35),
         font=F_SANS, size=9, color=a)
    footer(s)


# ══════════════════════════════════════════════
#  CLOSING
# ══════════════════════════════════════════════
def slide_why_now():
    s = _slide(); bg(s)
    circle(s, Inches(11.0), Inches(3.2), Inches(2.9), fill=GOLD, alpha=0.06)
    eyebrow(s, 'WHY NOW')
    text(s, '世界はもう、AIの接客に\n人件費以上を払っている。', Inches(0.75), Inches(1.15), Inches(11.8), Inches(1.7),
         font=F_MINCHO, size=32, color=INK, bold=True, line_spacing=1.16)
    cards = [
        ('Intercom Fin', '$0.99', '問い合わせを解決するごとに従量課金。使うほど請求が膨らむ。'),
        ('Qualified', '$68,000〜 / 年', 'AI営業担当(SDR)の年間コスト。それでも導入が進む。'),
        ('Sierra', 'Fortune 500', '世界の大企業がAI接客へ本気の投資を始めている。'),
    ]
    cw, ch, gx = Inches(3.86), Inches(2.0), Inches(0.26)
    x0, y0 = Inches(0.75), Inches(3.25)
    for i, (t, big, d) in enumerate(cards):
        x = x0 + i * (cw + gx)
        rect(s, x, y0, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, t, x + Inches(0.32), y0 + Inches(0.26), cw - Inches(0.6), Inches(0.35),
             font=F_DIDOT, size=15, color=INK, bold=True)
        text(s, big, x + Inches(0.32), y0 + Inches(0.68), cw - Inches(0.6), Inches(0.5),
             font=F_NUM, size=22, color=GOLD_HI, bold=True)
        text(s, d, x + Inches(0.32), y0 + Inches(1.28), cw - Inches(0.62), Inches(0.65),
             font=F_SANS, size=10, color=INK_DIM, line_spacing=1.3)
    text(s, 'CORE Continuum は、同じ水準の知性を ―― 日本語ネイティブの品質で、月額定額で。',
         Inches(0.75), Inches(5.75), Inches(11.8), Inches(0.6),
         font=F_MINCHO, size=18, color=GOLD_HI, italic=True)
    footer(s)


def slide_journey():
    s = _slide(); bg(s)
    eyebrow(s, 'ONE JOURNEY, SIX AGENTS')
    text(s, '顧客の一生を、途切れなく。', Inches(0.75), Inches(1.14), Inches(11.6), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    text(s, '出会いから、接客・販売・関係の継続・仲間づくりまで。バラバラのツールの「継ぎ目」でこぼれていた顧客が、一本の導線でつながる。',
         Inches(0.75), Inches(1.95), Inches(11.9), Inches(0.5), font=F_SANS, size=12.5, color=INK_DIM)
    steps = [
        ('IRIS', '見つけてもらう', 'SNS発信を企画から量産', C_IRI),
        ('LUME', '受け止めて売る', 'プロフィールでその場で収益化', C_LUM),
        ('CRYSTAL', 'おもてなしする', 'サイト訪問者を24時間接客', C_CRY),
        ('RESONANCE', '関係を続ける', 'LINEで一人ひとりに返信', C_RES),
        ('GUILD', '仲間にする', 'ファンを共創コミュニティへ', C_GLD),
        ('PRISM', '全体を指揮する', '情報を束ね、次の一手を提案', C_PRI),
    ]
    cw = Inches(1.93); gap = Inches(0.06)
    x0, y, ch = Inches(0.75), Inches(2.75), Inches(3.35)
    for i, (nm, role, d, col) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y, cw, Emu(42000), fill=col)
        text(s, f'{i+1:02d}', x, y + Inches(0.28), cw, Inches(0.45),
             font=F_NUM, size=20, color=col, bold=True, align='center')
        text(s, nm, x + Inches(0.06), y + Inches(0.92), cw - Inches(0.12), Inches(0.35),
             font=F_DIDOT, size=12.5, color=INK, bold=True, align='center')
        text(s, role, x + Inches(0.06), y + Inches(1.36), cw - Inches(0.12), Inches(0.35),
             font=F_GOTHIC, size=11.5, color=col, bold=True, align='center')
        text(s, d, x + Inches(0.14), y + Inches(1.85), cw - Inches(0.28), Inches(1.2),
             font=F_SANS, size=9.5, color=INK_DIM, align='center', line_spacing=1.3)
        if i < 5:
            text(s, '›', x + cw - Inches(0.05), y + Inches(1.3), Inches(0.18), Inches(0.5),
                 font=F_NUM, size=15, color=INK_FAINT, align='center')
    text(s, '単品でも強い。組み合わせると、こぼれない。', Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.5),
         font=F_MINCHO, size=16, color=GOLD_HI, italic=True)
    footer(s)


def slide_matrix():
    s = _slide(); bg(s)
    eyebrow(s, 'WHO NEEDS WHAT')
    text(s, '業種別・導入マップ', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    rows = [
        ('店舗・サロン・整体・教室', 'Resonance ＋ Crystal', 'LINE返信とサイト接客を、24時間AIが二正面で代行', (C_RES, C_CRY)),
        ('オンラインサロン・コミュニティ', 'Guild', '透明な意思決定と貢献の可視化で「関わりたくなる」場に', (C_GLD,)),
        ('クリエイター・インフルエンサー', 'Iris ＋ Lume', '発信の量産から、プロフィールでのその場収益化まで', (C_IRI, C_LUM)),
        ('経営者・士業・フリーランス', 'Prism', '知識を蓄えるほど賢くなる、24時間の専属AI秘書', (C_PRI,)),
        ('高級不動産・ホテル・接客業', 'Crystal (Luxury)', 'ブランドの言葉を守ったまま、深夜の見込み客を逃さない', (C_CRY,)),
        ('SNS運用代行会社', 'Iris (Agency)', 'クライアント別ワークスペース×30アカウントで量産を利益に', (C_IRI,)),
    ]
    x0, y0, rh = Inches(0.75), Inches(1.98), Inches(0.7)
    cols = [Inches(3.6), Inches(3.15), Inches(5.1)]
    tw = Inches(3.6) + Inches(3.15) + Inches(5.1)
    rect(s, x0, y0, tw, rh, fill=GOLD, alpha=0.13)
    for i, h in enumerate(['業種・立場', 'おすすめの組み合わせ', '効く理由']):
        cx = x0 + sum([Inches(0)] + [cols[j] for j in range(i)], Emu(0)) if i else x0
        cx = x0
        for j in range(i):
            cx += cols[j]
        text(s, h, cx + Inches(0.25), y0, cols[i] - Inches(0.4), rh, font=F_GOTHIC, size=12.5,
             color=GOLD, bold=True, anchor='middle')
    for r, (a1, a2, a3, dots) in enumerate(rows):
        y = y0 + rh * (r + 1)
        text(s, a1, x0 + Inches(0.25), y, cols[0] - Inches(0.4), rh, font=F_SANS, size=12,
             color=INK, anchor='middle')
        dx = x0 + cols[0] + Inches(0.25)
        for k, dc in enumerate(dots):
            circle(s, dx + Inches(0.06) + k * Inches(0.24), y + rh / 2, Inches(0.05), fill=dc)
        text(s, a2, dx + Inches(0.16) + len(dots) * Inches(0.24), y, cols[1] - Inches(0.7), rh,
             font=F_GOTHIC, size=12.5, color=GOLD_HI, bold=True, anchor='middle')
        text(s, a3, x0 + cols[0] + cols[1] + Inches(0.25), y, cols[2] - Inches(0.4), rh,
             font=F_SANS, size=11.5, color=INK_DIM, anchor='middle', line_spacing=1.2)
        rect(s, x0, y, tw, Emu(9000), fill=LINE)
    footer(s)


def slide_pricing_all():
    s = _slide(); bg(s)
    eyebrow(s, 'PRICING OVERVIEW')
    text(s, '全6サービス 料金一覧', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    rows = [
        ('Resonance', 'Solo ¥1,980', 'Pro ¥6,980', 'Business ¥14,800', C_RES),
        ('Guild', 'Free ¥0', 'Community ¥980', 'Pro ¥3,980', C_GLD),
        ('Iris', 'Lite ¥2,980', 'Standard ¥6,980', 'Agency ¥29,800', C_IRI),
        ('Lume', '無料 (7日間)', 'Pro ¥1,480', 'Business ¥3,480', C_LUM),
        ('Prism', 'Starter ¥2,980', 'Standard ¥9,800', 'Exclusive ¥29,800', C_PRI),
        ('Crystal', '―', 'Standard ¥29,800', 'Luxury ¥49,800', C_CRY),
    ]
    x0, y0, rh = Inches(0.75), Inches(2.02), Inches(0.63)
    cols = [Inches(2.55), Inches(3.0), Inches(3.0), Inches(3.3)]
    tw = sum([Inches(2.55), Inches(3.0), Inches(3.0), Inches(3.3)], Emu(0))
    rect(s, x0, y0, tw, rh, fill=GOLD, alpha=0.13)
    for i, h in enumerate(['サービス', 'エントリー', 'ミドル', 'ハイエンド']):
        cx = x0
        for j in range(i):
            cx += cols[j]
        text(s, h, cx + Inches(0.22), y0, cols[i] - Inches(0.35), rh, font=F_GOTHIC, size=12,
             color=GOLD, bold=True, anchor='middle')
    for r, row in enumerate(rows):
        y = y0 + rh * (r + 1)
        circle(s, x0 + Inches(0.3), y + rh / 2, Inches(0.055), fill=row[4])
        text(s, row[0], x0 + Inches(0.5), y, cols[0] - Inches(0.6), rh,
             font=F_DIDOT, size=14.5, color=INK, bold=True, anchor='middle')
        for i in range(1, 4):
            cx = x0
            for j in range(i):
                cx += cols[j]
            text(s, row[i], cx + Inches(0.22), y, cols[i] - Inches(0.35), rh, font=F_SANS,
                 size=12, color=INK if i == 3 else INK_DIM, bold=(i == 3), anchor='middle')
        rect(s, x0, y, tw, Emu(9000), fill=LINE)
    text(s, '※ 表示は月額・税込。Crystal 以外は無料または7日間トライアルから。Crystalのみ初期費用(設置代行つき)あり。',
         Inches(0.75), Inches(6.6), Inches(11.8), Inches(0.4), font=F_SANS, size=10, color=INK_FAINT)
    footer(s)


def slide_start():
    s = _slide(); bg(s)
    eyebrow(s, 'HOW TO START')
    text(s, '導入は、驚くほど簡単です。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    steps = [
        ('01', 'まず、無料で触る', '無料プランまたは7日間トライアルで、実際の画面をそのまま体験。クレジットカードの登録は要りません。'),
        ('02', 'ブランドを注ぐ', 'ブランド名・トーン・会社案内などの文章を貼るだけ。その場でAIの言葉づかいが、あなた仕様に変わります。'),
        ('03', '貼って、任せる', 'リンクを貼る、タグを1行入れる ― それだけで稼働開始。設置の丸投げ(代行)もメール1通で受け付けます。'),
    ]
    cw, gap = Inches(3.88), Inches(0.24)
    x0, y, h = Inches(0.75), Inches(2.3), Inches(3.6)
    for i, (num, t, d) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, h, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, num, x + Inches(0.38), y + Inches(0.34), Inches(2), Inches(0.9),
             font=F_NUM, size=40, color=GOLD, bold=True)
        rect(s, x + Inches(0.42), y + Inches(1.3), Inches(0.6), Emu(24000), fill=GOLD_HI)
        text(s, t, x + Inches(0.42), y + Inches(1.52), cw - Inches(0.84), Inches(0.6),
             font=F_GOTHIC, size=19, color=INK, bold=True)
        text(s, d, x + Inches(0.42), y + Inches(2.2), cw - Inches(0.84), Inches(1.25),
             font=F_SANS, size=11.5, color=INK_DIM, line_spacing=1.42)
        if i < 2:
            text(s, '→', x + cw - Inches(0.02), y + h / 2 - Inches(0.3), gap + Inches(0.06), Inches(0.6),
                 font=F_NUM, size=18, color=GOLD, align='center')
    text(s, 'すべてのサービスが、買う前に「あなたの設定のまま」試せます。',
         Inches(0.75), Inches(6.3), Inches(11.8), Inches(0.45),
         font=F_MINCHO, size=15, color=GOLD_HI, italic=True)
    footer(s)


def slide_cta():
    s = _slide(); bg(s)
    circle(s, Inches(SLIDE_W / 2), Inches(SLIDE_H / 2), Inches(3.3), fill=GOLD, alpha=0.07)
    eyebrow(s, 'GET STARTED')
    text(s, 'まずは1つ、\n無料で試してください。', Inches(0.75), Inches(1.6), Inches(11.8), Inches(1.9),
         font=F_MINCHO, size=38, color=INK, bold=True, line_spacing=1.18)
    text(s, '気になるサービスから1つでも、事業の接点に合わせた組み合わせでも。\nデモ・お見積り・導入のご相談はメール1通で ― 24時間以内にご返信します。',
         Inches(0.75), Inches(3.7), Inches(11.4), Inches(1.0), font=F_SANS, size=14.5, color=INK_DIM, line_spacing=1.45)
    grad_bar(s, Inches(0.78), Inches(4.95), Inches(2.4), Inches(0.05))
    rect(s, Inches(0.75), Inches(5.35), Inches(6.7), Inches(1.05), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'お問い合わせ', Inches(1.1), Inches(5.55), Inches(3), Inches(0.3),
         font=F_GOTHIC, size=10.5, color=GOLD, bold=True)
    text(s, 'core.guild.inc@gmail.com', Inches(1.1), Inches(5.85), Inches(6), Inches(0.45),
         font=F_NUM, size=19, color=INK, bold=True)
    footer(s)


def slide_closing():
    s = _slide(); bg(s)
    circle(s, Inches(SLIDE_W / 2), Inches(3.3), Inches(4.6), fill=GOLD, alpha=0.035)
    for i, c in enumerate([C_RES, C_GLD, C_IRI, C_LUM, C_PRI, C_CRY]):
        circle(s, Inches(SLIDE_W / 2 - 0.85) + i * Inches(0.34), Inches(2.25), Inches(0.05), fill=c)
    text(s, 'CORE  CONTINUUM', Inches(0.5), Inches(2.75), Inches(SLIDE_W - 1), Inches(1.0),
         font=F_DIDOT, size=46, color=INK, bold=True, align='center', spacing=300)
    text(s, '事業のあらゆる接点に、途切れない知性を。', Inches(0.5), Inches(3.95), Inches(SLIDE_W - 1), Inches(0.6),
         font=F_MINCHO, size=20, color=GOLD_HI, italic=True, align='center')
    grad_bar(s, Inches(SLIDE_W / 2 - 1.3), Inches(4.85), Inches(2.6), Inches(0.05))
    footer(s)


# ── 組み立て ───────────────────────────
slide_cover()
slide_problem_intro()
slide_thesis()
for p in PRODUCTS:
    slide_divider(p)
    slide_problem(p)
    slide_caps1(p)
    slide_caps2(p)
    slide_transform(p)
    slide_offer(p)
slide_why_now()
slide_journey()
slide_matrix()
slide_pricing_all()
slide_start()
slide_cta()
slide_closing()

out = os.path.expanduser('~/Desktop/CORE_Continuum_6サービス_営業資料.pptx')
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides._sldIdLst))
