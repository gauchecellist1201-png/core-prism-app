#!/usr/bin/env python3
# ============================================================
# CORE Continuum — 6サービス 営業資料 v4 (final)
# 出力: ~/Desktop/CORE_Continuum_6サービス_営業資料.pptx
# v4: ①実際のサービスロゴ(PNG)を透過処理して全面使用
#     ②文章をやさしい言葉に全面書き直し(専門用語なし・短文)
#     構成: intro3 + 各製品6枚×6 + closing7 = 46枚
# ============================================================
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ══════════════════════════════════════════════
#  ロゴ素材の前処理 (黒背景 → 透過、アイコン切り出し)
# ══════════════════════════════════════════════
ASSET = '/tmp/deck_assets_v4'
os.makedirs(ASSET, exist_ok=True)

LOGO_SRC = {
    'RES': '/Users/naokiide/core-manual-build/logos/resonance.png',
    'GLD': '/Users/naokiide/core-manual-build/logos/guild.png',
    'IRI': '/Users/naokiide/core-manual-build/logos/iris.png',
    'LUM': '/Users/naokiide/core-manual-build/logos/lume.png',
    'PRI': '/Users/naokiide/core-manual-build/logos/prism.png',
    'CRY': '/Users/naokiide/core-prism-app/public/crystal-512.png',
}

def make_transparent(im, thr):
    im = im.convert('RGBA')
    out = []
    for r, g, b, a in im.getdata():
        out.append((r, g, b, 0) if max(r, g, b) < thr else (r, g, b, a))
    im.putdata(out)
    return im

WORDMARK = {}
ICON = {}
for key, src in LOGO_SRC.items():
    thr = 42 if key == 'CRY' else 26
    im = make_transparent(Image.open(src), thr)
    bb = im.getbbox()
    if bb:
        im = im.crop(bb)
    wp = f'{ASSET}/{key}_wm.png'
    im.save(wp)
    WORDMARK[key] = (wp, im.width / im.height)  # (path, aspect)
    # アイコン部 (左端の正方形近傍) — Crystalは全体がアイコン
    if key == 'CRY':
        ICON[key] = (wp, 1.0)
    else:
        h = im.height
        ic = im.crop((0, 0, min(int(h * 1.18), im.width), h))
        bb2 = ic.getbbox()
        if bb2:
            ic = ic.crop(bb2)
        ip = f'{ASSET}/{key}_icon.png'
        ic.save(ip)
        ICON[key] = (ip, ic.width / ic.height)

# ══════════════════════════════════════════════
#  デッキ基盤
# ══════════════════════════════════════════════
BG        = RGBColor(0x07, 0x07, 0x0B)
BG_PANEL  = RGBColor(0x11, 0x11, 0x17)
BG_NUM    = RGBColor(0x1F, 0x1D, 0x26)
INK       = RGBColor(0xF7, 0xF3, 0xEA)
INK_DIM   = RGBColor(0xB2, 0xAA, 0x9C)
INK_FAINT = RGBColor(0x60, 0x5A, 0x51)
GOLD      = RGBColor(0xC9, 0xA2, 0x4B)
GOLD_HI   = RGBColor(0xEA, 0xD1, 0x96)
LINE      = RGBColor(0x2B, 0x28, 0x23)
RED_SOFT  = RGBColor(0xC2, 0x5B, 0x54)

C_RES = RGBColor(0x35, 0xC2, 0x8F)
C_GLD = RGBColor(0x4F, 0xD4, 0xD4)
C_IRI = RGBColor(0xE8, 0x5B, 0x9E)
C_LUM = RGBColor(0xF5, 0x9B, 0x2D)
C_PRI = RGBColor(0xB8, 0x86, 0xE8)
C_CRY = RGBColor(0xA9, 0xC3, 0xE8)

F_DIDOT  = 'Didot'
F_OPTIMA = 'Optima'
F_MINCHO = 'Hiragino Mincho ProN'
F_GOTHIC = 'Hiragino Kaku Gothic StdN'
F_SANS   = 'Hiragino Sans'
F_NUM    = 'Avenir Next'

SLIDE_W, SLIDE_H = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)
BLANK = prs.slide_layouts[6]
_pg = {'n': 0}


def _slide():
    return prs.slides.add_slide(BLANK)


def _alpha(sh, alpha):
    sp = sh.fill._xPr
    sf = sp.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            a = srgb.find(qn('a:alpha'))
            if a is None:
                a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', str(int(alpha * 100000)))


def bg(slide, color=BG):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.shadow.inherit = False
    slide.shapes._spTree.remove(sh._element)
    slide.shapes._spTree.insert(2, sh._element)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=None, alpha=None, rounded=False, radius=0.08):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(st, x, y, w, h)
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 0.75)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    if rounded:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    return sh


def circle(slide, cx, cy, r, fill=None, alpha=None, line=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    if fill is not None:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
        if alpha is not None:
            _alpha(sh, alpha)
    else:
        sh.fill.background()
    if line is not None:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, s, x, y, w, h, font=F_SANS, size=14, color=INK, bold=False, italic=False,
         align='left', anchor='top', spacing=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = {'top': MSO_ANCHOR.TOP, 'bottom': MSO_ANCHOR.BOTTOM}.get(anchor, MSO_ANCHOR.MIDDLE)
    al = {'left': PP_ALIGN.LEFT, 'right': PP_ALIGN.RIGHT, 'center': PP_ALIGN.CENTER}.get(align, PP_ALIGN.LEFT)
    for i, ln in enumerate(s.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = al
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic
        if spacing is not None:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(spacing))
    return tb


def pic(slide, key, x, y, h, kind='wm'):
    """ロゴ画像を高さ指定で配置し、実幅(Inches)を返す"""
    path, aspect = WORDMARK[key] if kind == 'wm' else ICON[key]
    slide.shapes.add_picture(path, x, y, height=Inches(h))
    return h * aspect


def grad_bar(slide, x, y, w, h):
    seg = w // 2
    rect(slide, x, y, seg, h, fill=GOLD)
    rect(slide, x + seg, y, w - seg, h, fill=GOLD_HI)


def footer(slide):
    _pg['n'] += 1
    text(slide, 'CORE  CONTINUUM', Inches(0.62), Inches(SLIDE_H - 0.42), Inches(4), Inches(0.3),
         font=F_OPTIMA, size=8, color=INK_FAINT, spacing=250)
    text(slide, f"{_pg['n']:02d}", Inches(SLIDE_W - 1.5), Inches(SLIDE_H - 0.42), Inches(0.9), Inches(0.3),
         font=F_NUM, size=9, color=INK_FAINT, align='right')


def eyebrow(slide, en, accent=GOLD, no=None, y=Inches(0.6)):
    x = Inches(0.75)
    if no:
        text(slide, no, x, y - Inches(0.02), Inches(0.6), Inches(0.35),
             font=F_NUM, size=13, color=accent, bold=True)
        x = Inches(1.36)
    rect(slide, x, y + Inches(0.12), Inches(0.34), Emu(20000), fill=accent)
    text(slide, en, x + Inches(0.46), y, Inches(8), Inches(0.35),
         font=F_OPTIMA, size=12, color=accent, bold=True, spacing=340)


def running_head(slide, p):
    """右上に実ロゴ(小)"""
    h = 0.30 if p['key'] != 'CRY' else 0.34
    path, aspect = WORDMARK[p['key']]
    w = h * aspect
    slide.shapes.add_picture(path, Inches(SLIDE_W - 0.7 - w), Inches(0.58), height=Inches(h))
    if p['key'] == 'CRY':
        text(slide, 'Crystal', Inches(SLIDE_W - 2.6), Inches(0.58), Inches(1.4), Inches(0.34),
             font=F_DIDOT, size=15, color=INK_DIM, align='right', anchor='middle')


# ══════════════════════════════════════════════
#  製品データ — やさしい言葉で
# ══════════════════════════════════════════════
PRODUCTS = [
    dict(
        key='RES', name='Resonance', accent=C_RES, no='01', en='LINE AI',
        area='公式LINEの自動返信AI',
        promise='LINEに来たメッセージに、AIがあなたの代わりに返事。\nしかも、LINEの無料枠(月200通)は使いません。',
        metrics=[('0通', '無料枠を使わない ― 何通返しても追加料金なし'),
                 ('24時間', '夜中でも、仕事中でも、すぐ返事'),
                 ('7日間', '無料で全部試せる ― カード登録なし')],
        problem_head='「あとで返そう」が、\nいちばん高くついている。',
        problem_sub='公式LINEを使うお店・教室・ひとり事業主の、毎日の困りごと。',
        pains=[
            ('返事が遅れて、お客様が離れる', '仕事中や夜に来た問い合わせ。返す頃には、もう他のお店に決めている。'),
            ('月200通が気になって、送れない', '無料で送れるのは月200通まで。「使い切ったら怖い」で、お知らせをがまんしてしまう。'),
            ('全員に同じ文面は、もう読まれない', '一斉送信は開かれない。それどころか、ブロックが増えていく。'),
            ('予約の約束を、うっかり忘れる', '「来週の火曜で」がトークに埋もれて、カレンダーに書き忘れる。'),
        ],
        caps=[
            ('AIが返事を書いて、送る', '会話の流れとあなたの話し方をまねて、すぐに返事。LINEの無料枠は1通も使いません。'),
            ('相手に合わせた話し方', '敬語の人には敬語で。くだけた人にはやわらかく。迷ったら丁寧な言葉に寄せます。'),
            ('あなたの「らしさ」を覚える', '自己紹介の文章を貼るだけ。AIがあなたの言葉づかいや考え方で話すようになります。'),
            ('危ない言葉は、送る前に止める', '「必ず」「返金」など決めておいた言葉が入ったら送信ストップ。言い過ぎを防ぎます。'),
            ('予約の話を見逃さない', '会話に日にちの話が出たら、自動で「予約リスト」に載せます。埋もれません。'),
            ('カレンダーにも自動で登録', '日時が決まればGoogleカレンダーへ。オンライン相談ならMeetのURLも自動で用意。'),
            ('返事待ちを、大事な順に', '予約や料金の問い合わせが上に並ぶ。AIの下書きつきで、確認して押すだけ。'),
            ('ひとりずつ違う文面で、一斉送信', '同じお知らせでも、相手の名前ややり取りに合わせて文面を書き分けます。'),
            ('お客様の「熱さ」が点数で見える', 'よくやり取りする人ほど高い点数に。今日誰に声をかけるべきか、すぐ分かる。'),
            ('お客様カルテを自動で作る', 'その人の人柄・好きなこと・響く言葉を、AIがメモにまとめてくれます。'),
            ('友だち追加後の自動フォロー', '追加から7日間、決めたタイミングで自動でごあいさつ。これも無料枠ゼロ。'),
            ('走り書きを、送れる文章に', '「木曜14時 空きあり」のメモを、そのまま送れる丁寧な文章に整えます。'),
        ],
        before_after=[
            ('返事が翌日になり、お客様は他のお店へ', '夜中でもAIがすぐ返事。チャンスを逃さない'),
            ('一斉送信が読まれず、ブロックが増える', 'ひとりずつ違う文面で、ちゃんと読まれる'),
            ('予約の約束を忘れて、信頼を失う', '自動でカレンダー登録。Meetの発行まで'),
        ],
        diff=[
            '自動返信がLINEの無料枠(月200通)を使わない',
            'あなたの話し方をまねる。ロボットっぽくならない',
            '危ない言葉は、送る前に自動で止まる',
            '会話から予約、カレンダー登録まで、ぜんぶ自動',
        ],
        ideal='美容室・サロン・整体・教室・カフェなど、公式LINEでお客様とやり取りする、すべてのお店とひとり事業主に。',
        pricing=[('Free', '¥0', 'お客様名簿・トーク画面・熱さの点数'),
                 ('Solo', '¥1,980/月', 'AI返事 1日5通・返事待ちリスト'),
                 ('Pro', '¥6,980/月', 'AI返事 1日20通・3アカウント・一斉送信'),
                 ('Business', '¥14,800/月', 'AI返事 無制限・10アカウント')],
        price_note='どのプランも7日間無料。カード登録なし。いつでも解約OK。',
    ),
    dict(
        key='GLD', name='Guild', accent=C_GLD, no='02', en='COMMUNITY OS',
        area='コミュニティ運営ツール',
        promise='みんなで決めて、記録して、頑張った人に報いる。\nコミュニティ運営が、これひとつでラクになります。',
        metrics=[('改ざん不可', '決めたことの記録は、あとから書き換えられない'),
                 ('+200', '良い提案をした人に届く、お礼のトークン'),
                 ('15人まで', '無料で始められる')],
        problem_head='「また同じ話で、\n揉めていませんか。」',
        problem_sub='オンラインサロン・サークル・コミュニティ運営者の、毎月の困りごと。',
        pains=[
            ('決めたことが、流れて消える', 'チャットでなんとなく決まって、あとから「そんな話あったっけ?」。また同じ議論が始まる。'),
            ('声の大きい人だけで決まる', '静かな人の本気の意見が埋もれる。多数決では「本気の1票」も「なんとなくの1票」も同じ。'),
            ('頑張る人ほど、報われずに辞める', '場を支えてくれる人の働きが、誰にも見えていない。熱心な人から順に去っていく。'),
            ('「言っても無駄」とあきらめられる', '出した意見が採用されたのか、流されたのか分からない。運営への信頼が育たない。'),
        ],
        caps=[
            ('提案して、みんなで投票', '「こうしませんか?」を出すと、みんなが投票。本気の議題には強い票を入れられます。'),
            ('締め切りが来たら、自動で決定', '期限が来れば自動で集計して結果が出る。ほったらかしでも議題がたまりません。'),
            ('記録は、書き換えられない', '決めたことは特別な仕組みで鎖のようにつながって記録。1つでも改ざんすると必ずバレます。'),
            ('誰でも「本物の記録」と確認できる', 'メンバー自身が自分のスマホで、記録が書き換えられていないことを確かめられます。'),
            ('頑張りが、ポイントになる', '提案・投票・コメント。やった分だけ自動でポイントが貯まる。ズルはできません。'),
            ('ポイントで、称号が上がる', '貯めると「仲間」「中核」「殿堂」など称号が解放。名前も少しずつ輝いていきます。'),
            ('良い提案には、お礼が届く', '意見が採用されたら、その人にお礼のトークン。「見送り」でも理由がちゃんと残ります。'),
            ('新しい人を、ちゃんと歓迎', '参加から48時間は、歓迎メッセージと「まずやること」を一番上に案内します。'),
            ('質問2つで、自分の居場所', '「何をしたい?」「興味は?」を選ぶだけで、画面がその人向けに変わります。'),
            ('チャットも、人集めも', 'ふだんの雑談も、イベントの募集も、この中で完結。既読バッジや締切管理つき。'),
            ('「決めたのにやってない」をなくす', '決まったことの進み具合(未着手/進行中/完了)を、みんなで見える化。'),
            ('意見の行き先を、全員に公開', 'もらった声を「次やる/検討中/実装済み/見送り」の一覧で公開。応援投票もできます。'),
        ],
        before_after=[
            ('決めたことが流れて、また同じ話に', '提案→投票→自動で決定。記録はずっと残る'),
            ('頑張る人が報われず、辞めていく', '頑張りがポイントと称号で見える。お礼も届く'),
            ('「運営が勝手に決めてる」と言われる', '記録は改ざんできない。誰でも確かめられる'),
        ],
        diff=[
            '記録の改ざんが、仕組みの上でできない(他のツールにはない)',
            '本気の意見がちゃんと通る、新しい投票のかたち',
            '頑張った人に、お礼のトークンが実際に届く',
            '「見送り」の理由まで残る。意見を握りつぶせない',
        ],
        ideal='オンラインサロン・ファンコミュニティ・サークルや部活・習い事教室・社内の有志活動。「みんなで決める場」があるすべての運営者に。',
        pricing=[('Free', '¥0', '15人まで・1つのギルド・記録も投票も全部'),
                 ('Community', '¥980/月', '人数無制限・ギルド5つ・お礼トークン'),
                 ('Pro', '¥3,980/月', 'ギルド無制限・進捗管理・記録の書き出し'),
                 ('Enterprise', '相談', '会社での導入・専任サポート')],
        price_note='年払いなら2ヶ月分おトク。まず無料の1ギルドから。',
    ),
    dict(
        key='IRI', name='Iris', accent=C_IRI, no='03', en='INSTAGRAM AI',
        area='インスタ運用のAI相棒',
        promise='「今日、何を投稿しよう」をなくします。\n企画も台本も字幕も、AIが用意。案件の相談まで。',
        metrics=[('18種類', '企画と台本の型 ― 選ぶだけで完成'),
                 ('30個', '運用代行向け ― 同時に管理できるアカウント数'),
                 ('1タップ', '字幕データを書き出して、動画アプリに貼るだけ')],
        problem_head='毎朝、白紙とにらめっこ\nしていませんか。',
        problem_sub='クリエイターと、インスタ運用代行会社の毎日の困りごと。',
        pains=[
            ('「何を投稿しよう」で手が止まる', '毎回ゼロから考えるのが、いちばん大変。ネタ切れの朝がつらい。'),
            ('数字を見ても、直し方が分からない', '分析アプリはあっても、「じゃあ明日どうすれば?」は教えてくれない。'),
            ('台本づくりに時間がかかりすぎる', '運用代行なら、この手間がそのまま人件費に。受注を増やすほど苦しくなる。'),
            ('届いた案件が、本物か分からない', '詐欺っぽいDM、安すぎる報酬。判断に迷って、返事も遅れる。'),
        ],
        caps=[
            ('企画を、10種類の型から自動作成', '「悩み解決」「比較」「Q&A」など型を選ぶだけ。あなたのジャンルに合わせたネタが出てきます。'),
            ('台本も、8種類の型から自動作成', '最初のセリフ、カメラの割り方、テロップ、投稿文まで一式。撮る人がそのまま動けます。'),
            ('字幕データまで、自動で作る', '台本から字幕ファイルを書き出し。CapCutなどの動画アプリに貼るだけで字幕完成。'),
            ('アカウントをAIが採点', 'プロフィールや投稿を5つの観点で点数化。どこを直せば伸びるか、すぐ分かります。'),
            ('投稿にいい時間を教えてくれる', '曜日×時間帯で「狙い目」を提案。データが足りない時は正直に「目安」と表示します。'),
            ('投稿前に、9マスで並びを確認', 'インスタのプロフィール画面と同じ形で見え方をチェック。並べ替えもできます。'),
            ('案件のやり取りを、まとめて管理', '問い合わせから投稿、報告まで、案件の進み具合をひと目で。DMの画像から自動で読み取り。'),
            ('怪しい案件は、自動で警告', '「無料でやって」「暗号資産」など危ないサインを自動チェック。安心して返事できます。'),
            ('相場も教えてくれる', '「この案件、受けていい?」に、目安の金額つきで答えてくれます。'),
            ('返信文も、AIが書く', '丁寧・フレンドリーなど4種類の文体で、すぐ送れる返信案を用意します。'),
            ('企業とのマッチング', 'あなたに合いそうな企業を探して、応募のメールまで自動で作ります。'),
            ('あなたの「らしさ」を学習', '使うほど、あなたらしい企画・台本に。代行なら、クライアントごとに学習します。'),
        ],
        before_after=[
            ('毎朝30分、投稿のネタに悩む', '型を選ぶだけで、企画がすぐ並ぶ'),
            ('台本がなくて、撮影がグダグダに', 'セリフもカット割りも字幕も、一式そろう'),
            ('怪しい案件に、勘で返事している', '危険チェックと相場つきで、安心して返せる'),
        ],
        diff=[
            '分析から投稿、案件のやり取りまで、これ1つで全部',
            '運用代行は30アカウントまで一括管理。台本は作り放題',
            '数字をごまかさない。足りない時は「目安」と正直に表示',
            '使うほど「その人らしさ」を学習していく',
        ],
        ideal='フォロワーを増やしたいクリエイター、副業インフルエンサー。そして特に、インスタ運用代行の会社に。台本づくりの人件費が要らなくなります。',
        pricing=[('Lite', '¥2,980/月', '入門 ― AI相談 月30回・案件管理'),
                 ('Standard', '¥6,980/月', '人気 ― AI相談無制限・アカウント分析'),
                 ('Pro', '¥12,800/月', '企画・台本づくり・5アカウント'),
                 ('Agency', '¥29,800/月', '代行向け ― 30アカウント・台本作り放題')],
        price_note='7日間無料。カード登録なし。勝手に課金されません。',
    ),
    dict(
        key='LUM', name='Lume', accent=C_LUM, no='04', en='PROFILE PAGE',
        area='売れるプロフィールページ',
        promise='SNSに貼る「あなたのページ」を30秒で。\n見るだけのページではなく、その場で売れるページです。',
        metrics=[('30秒', '質問3つに答えるだけで、ページがほぼ完成'),
                 ('4種類', '「買う・予約・応援・登録」ボタンを置ける'),
                 ('5,000件', 'お客様のメールアドレスを貯めておける')],
        problem_head='フォロワー1万人。\n売上、ゼロ円。',
        problem_sub='SNSのプロフィールリンクで、毎日起きていること。',
        pains=[
            ('リンクを並べても、売上にならない', 'よくあるリンクまとめは、他のサイトへの案内板。ページ自体は1円も生まない。'),
            ('どのリンクが押されてるか、分からない', '見えるのは合計クリック数だけ。誰にも押されていないリンクに気づけない。'),
            ('来てくれた人と、二度と繋がれない', '連絡先を受け取る仕組みがなく、そのまま帰られてしまう。'),
            ('何を直せばいいか、分からない', '数字を眺めても、次にやることは分からない。結局そのまま放置。'),
        ],
        caps=[
            ('30秒で、ページが完成', '名前・ひとこと・使っているSNSの3つに答えるだけ。ほぼ完成の状態から始められます。'),
            ('編集した瞬間、見た目に反映', '保存ボタンはありません。文字を打ったそばから、プレビューに映ります。'),
            ('デザインは12種類から選ぶだけ', '切り替えもふわっときれい。色や文字の形も、好みに変えられます。'),
            ('「買う・予約・応援」ボタンを置ける', '商品の販売、予約の受付、投げ銭。ぜんぶページの中で完結します。'),
            ('メールアドレスを集められる', '見に来た人の連絡先を最大5,000件まで保存。ファイルで取り出せます。'),
            ('期間限定のリンク', 'セールやイベントの告知に。開始と終了の日時を入れるだけで、自動で出て自動で消えます。'),
            ('押された場所が、色で分かる', 'よく押されるリンクは赤く、押されないリンクは青く。無駄が一目瞭然。'),
            ('どこから来た人か、分かる', '「インスタから来た人は商品ページを見る」まで分かる。SNSごとの相性が見えます。'),
            ('伸びているかが、グラフで分かる', '8週間の推移と「先週より+◯%」。頑張りが数字で見えるから続きます。'),
            ('毎週、AIが通信簿をくれる', '今週の結果と「次にやるといいこと」を、やさしい日本語で教えてくれます。'),
            ('直すのは、ワンタップ', '「人気順に並べ替えましょう」などの提案を、ボタン1つでそのまま実行。'),
            ('続ける約束を、自分で選べる', '2日・5日・7日・14日から自分で選ぶ。途切れても責められません。「また今日から」。'),
        ],
        before_after=[
            ('フォロワー1万人。でも売上はゼロ', 'ページに「買う」ボタン。その場で売れる'),
            ('どのリンクが無駄か、分からないまま', '色で一目瞭然。直すのもワンタップ'),
            ('告知の貼り替えを、毎回手作業で', '日時を入れるだけ。自動で出て、自動で消える'),
        ],
        diff=[
            'ページの中で売れる(よくあるリンク集は「貼るだけ」)',
            'どのリンクが効いているか、色ですぐ分かる',
            '毎週AIが「次の一手」を教えてくれる',
            '表示がとても速い。待たせない',
        ],
        ideal='ミュージシャン・作家・写真家などの作品を売る人、予約を受けたいサロンや教室、お客様リストを育てたいひとり事業主に。',
        pricing=[('無料', '¥0 (7日間)', 'リンク置き放題・基本デザイン'),
                 ('Pro', '¥1,480/月', '全デザイン・色分け分析・自分のドメイン'),
                 ('Business', '¥3,480/月', 'Proの全部＋チームで管理・EC連携')],
        price_note='7日間無料・カード登録なし。',
    ),
    dict(
        key='PRI', name='Prism', accent=C_PRI, no='05', en='AI SECRETARY',
        area='あなた専属のAI秘書',
        promise='資料もメモも数字も、ぜんぶ覚えてくれるAI秘書。\n朝になると「今日やること」を先に教えてくれます。',
        metrics=[('30分→3分', '会議の議事録 ― 録音を渡すだけ'),
                 ('14人', '専門家AI ― 財務・営業などプロの視点で答える'),
                 ('朝・夜', '聞かなくても、向こうから提案してくれる')],
        problem_head='雑務に追われて、\n考える時間がない。',
        problem_sub='ひとりで全部を回す経営者・士業・フリーランスの、時間の現実。',
        pains=[
            ('雑務ばかりで、1日が終わる', '議事録、請求書、提案書。気づけば夕方。大事なことを考える時間が残らない。'),
            ('メモも資料も、そのまま眠っている', 'せっかく貯めた情報が、次の行動につながらない。フォルダの肥やしに。'),
            ('AIに毎回、イチから説明するのが面倒', 'ChatGPTは便利。でも、あなたの仕事も数字も覚えてくれない。'),
            ('相談相手が、いない', '経営も営業もお金のことも、ぜんぶひとりで決めている。'),
        ],
        caps=[
            ('資料を放り込むだけで、整理される', 'メモ・PDF・画像を入れると、要点と「やるべきこと」に自動でまとまります。'),
            ('数字も、自動で拾う', '資料の中の売上や経費を見つけて、月ごとの数字にそろえてくれます。'),
            ('朝、「今日の一手」が届く', '覚えた情報から「今日はこれをやりましょう」を先回りで提案してくれます。'),
            ('提案の根拠が見える', '「あなたの知識◯件をもとに提案」と表示。数字のごまかしは一切ありません。'),
            ('やることを、小さく分解', '大きな提案を「明日からできる手順」に分けてくれる。動き出せます。'),
            ('眠っている資料を、活かす提案', '「この資料、こう使えます」を、AIの方から言ってきてくれます。'),
            ('会議の録音を、投げるだけ', '議事録・決まったこと・宿題に自動でまとめる。30分の作業が3分に。'),
            ('思いつきを、2タップでメモ', '画面の隅の入力口に放り込むだけ。あとでAIがちゃんと活かします。'),
            ('何でも、1つの入力欄から', 'ボタン1つで呼び出して、議事録も請求書も検索も、ぜんぶここから。'),
            ('チャットは、畳める', '使わない時は小さく畳んで、画面を広く。じゃまになりません。'),
            ('14人の専門家AIに相談できる', '財務・営業・マーケなど、専門ごとのAI役員に相談したり、仕事を任せたり。'),
            ('他のCOREサービスとつながる', 'Lume・Iris・Resonanceの情報がつながり、お客様への次の一言まで提案。'),
        ],
        before_after=[
            ('議事録づくりに30分かかる', '録音を投げるだけ。3分で完成'),
            ('提案書を、白紙から書き始める', '覚えた知識から、下書きがすぐ出てくる'),
            ('AIに毎回、イチから説明する', 'あなたの仕事を覚えたAIが、先回りで動く'),
        ],
        diff=[
            '使うほど賢くなる(ふつうのAIは毎回ゼロから)',
            '提案の根拠を正直に見せる。数字を盛らない',
            '聞かなくても、朝になると先に動いてくれる',
            'バラバラな月額ツールを、これ1つにまとめられる',
        ],
        ideal='ひとりで全部を回す経営者・ひとり社長・士業・フリーランス、少人数の会社の代表に。「右腕が欲しいけど、雇うほどではない」方へ。',
        pricing=[('Starter', '¥2,980/月', '3人の専門家AI・議事録・スライド作成'),
                 ('Standard', '¥9,800/月', '7人の専門家AI・提案書・メール秘書'),
                 ('Exclusive', '¥29,800/月', '全機能・チーム5人・専任サポート')],
        price_note='AIの利用料込み。7日間無料。',
    ),
    dict(
        key='CRY', name='Crystal', accent=C_CRY, no='06', en='AI CONCIERGE',
        area='24時間のAI接客係',
        promise='お店やサイトに来た人を、AIが24時間おもてなし。\n話しかけるだけで答えて、連絡先まで受け取ります。',
        metrics=[('1分', '設置にかかる時間 ― リンクを貼るだけ'),
                 ('24時間', '深夜も休日も、あなたのお店の言葉で応対'),
                 ('定額', '月額だけ。使った分の追加料金なし')],
        problem_head='深夜2時の問い合わせは、\n今夜も無人のまま。',
        problem_sub='ホテル・不動産・サロン・高級サービス業で、毎晩起きている機会損失。',
        pains=[
            ('夜の問い合わせに、誰も出られない', 'お客様が調べ物をするのは夜。翌朝の返事では、もう他と契約している。'),
            ('人を雇うほどでは、ない', 'チャット担当を置く余裕はない。でも、取りこぼしは痛い。'),
            ('サイトに来ても、何も聞かずに帰る', '迷ったまま離脱。誰が来て、何を知りたかったのかも分からない。'),
            ('安っぽいボットは、ブランドに合わない', 'ちぐはぐな定型文の応対が、せっかくの世界観を壊してしまう。'),
        ],
        caps=[
            ('24時間、いつでも応対', '深夜も休日も、あなたのお店の言葉づかいで答え続けます。'),
            ('声で話せる', 'お客様がマイクに話しかけると、きれいな字幕と声で答えます。スマホでもOK。'),
            ('会社案内を貼るだけで、賢くなる', '貼った内容だけを根拠に答えます。勝手な作り話は、しません。'),
            ('よくある質問を、自動で作る', '貼った文章から、Q&Aを最大6組自動で用意。手で書かなくていい。'),
            ('商談の約束を、取りつける', '話が盛り上がった瞬間に、連絡先の入力欄をすっと差し出します。'),
            ('見込みのあるお客様を、見分ける', '「ご予算」「時期」などの条件を、会話の中でさりげなく確認してくれます。'),
            ('会話が、まるごとメールで届く', 'お名前・連絡先・会話の全文があなたのメールに。折り返すだけで商談に。'),
            ('外国語でも、大丈夫', '英語や中国語で話しかけられたら、同じ言葉で答えます。'),
            ('AIから、先に声をかける', '迷っている人に「ご質問はありますか?」とひとこと。会話が始まります。'),
            ('人格と色を、選べる', '「正統派」「親しみ」「簡潔」の3タイプ。色も呼び名も、お店に合わせて。'),
            ('予約ページへ、そのまま案内', '会話の流れで予約ボタンを出して、そのまま予約完了まで。'),
            ('設置は1分。貼るだけ', '専用リンクを貼るだけ。QRコードも自動で作成。丸投げの設置代行もあります。'),
        ],
        before_after=[
            ('夜の問い合わせは、朝まで放置', '深夜でもその場で応対。連絡先まで受け取る'),
            ('訪問者が、黙って帰っていく', 'AIから声をかけて、会話がはじまる'),
            ('ボット導入は、業者に頼んで数週間', 'リンクを貼るだけ。今日から働き始める'),
        ],
        diff=[
            '月額だけ。使った分の追加料金なし(海外の有名AIは1回ごとに課金)',
            '声で話せて、外国語にも答えられる',
            '工事もエンジニアも不要。リンクを貼るだけ',
            '応対から商談の約束まで、これ1つで',
        ],
        ideal='高級不動産・ホテルや旅館・美容クリニック・サロン・士業など、「接客の質」がそのままブランドになる業種に。',
        pricing=[('Standard', '¥29,800/月', '初期¥98,000(設置代行つき) ― 全機能'),
                 ('Luxury', '¥49,800/月', '初期¥298,000 ― 専任の磨き込み・複数サイト')],
        price_note='買う前に、あなたのお店の設定のまま無料で試せます。',
    ),
]


# ══════════════════════════════════════════════
#  INTRO
# ══════════════════════════════════════════════
def slide_cover():
    s = _slide(); bg(s)
    circle(s, Inches(2.6), Inches(1.7), Inches(3.0), fill=GOLD, alpha=0.06)
    circle(s, Inches(11.8), Inches(5.8), Inches(2.7), fill=GOLD_HI, alpha=0.05)
    text(s, 'CORE  CONTINUUM', Inches(0.9), Inches(0.95), Inches(11), Inches(0.5),
         font=F_DIDOT, size=20, color=GOLD, bold=True, spacing=520)
    text(s, '事業のあらゆる接点に、\n途切れない知性を。', Inches(0.86), Inches(1.75), Inches(11.8), Inches(2.4),
         font=F_MINCHO, size=52, color=INK, bold=True, line_spacing=1.12)
    grad_bar(s, Inches(0.92), Inches(4.35), Inches(2.6), Inches(0.05))
    text(s, '返信・接客・発信・販売・事務・コミュニティ ―― 6つの仕事それぞれに、専属のAIを。',
         Inches(0.9), Inches(4.68), Inches(11.6), Inches(0.5), font=F_GOTHIC, size=15, color=INK_DIM)
    # 実ロゴ 6つ (2行×3列)
    xs = [Inches(0.9), Inches(5.15), Inches(9.4)]
    ys = [Inches(5.55), Inches(6.35)]
    order = ['RES', 'GLD', 'IRI', 'LUM', 'PRI', 'CRY']
    for i, k in enumerate(order):
        r, c = divmod(i, 3)
        h = 0.34 if k != 'CRY' else 0.5
        y = ys[r] if k != 'CRY' else ys[r] - Inches(0.08)
        w = pic(s, k, xs[c], y, h)
        if k == 'CRY':
            text(s, 'Crystal', xs[c] + Inches(w + 0.12), ys[r] - Inches(0.03), Inches(1.6), Inches(0.42),
                 font=F_DIDOT, size=19, color=INK, bold=True, anchor='middle')
    footer(s)


def slide_problem_intro():
    s = _slide(); bg(s)
    eyebrow(s, 'THE PROBLEM')
    text(s, '売上がこぼれていく場所は、\n6つある。', Inches(0.75), Inches(1.1), Inches(12), Inches(1.8),
         font=F_MINCHO, size=34, color=INK, bold=True, line_spacing=1.16)
    holes = [
        ('返事が遅れる', '夜や作業中の問い合わせに、出られない。'),
        ('お知らせが読まれない', '全員同じ文面は、もう開かれない。'),
        ('見に来た人が帰ってしまう', 'サイトに来ても、何もせずに離脱。'),
        ('情報が散らかったまま', 'メモも資料も、活かされずに眠る。'),
        ('頑張る人が辞めていく', 'コミュニティの貢献が、誰にも見えない。'),
        ('人を増やすと、お金がかかる', '対応を増やす＝人件費が増える。'),
    ]
    cw, ch, gx, gy = Inches(3.83), Inches(1.56), Inches(0.22), Inches(0.24)
    x0, y0 = Inches(0.75), Inches(3.2)
    for i, (t, d) in enumerate(holes):
        r, c = divmod(i, 3)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y + Inches(0.24), Emu(34000), Inches(1.08), fill=RED_SOFT)
        text(s, f'{i+1:02d}', x + Inches(0.3), y + Inches(0.24), Inches(0.8), Inches(0.3),
             font=F_NUM, size=11, color=RED_SOFT, bold=True)
        text(s, t, x + Inches(0.3), y + Inches(0.5), cw - Inches(0.5), Inches(0.42),
             font=F_GOTHIC, size=16, color=INK, bold=True)
        text(s, d, x + Inches(0.3), y + Inches(0.97), cw - Inches(0.55), Inches(0.55),
             font=F_SANS, size=11, color=INK_DIM, line_spacing=1.28)
    text(s, 'どれも「人が足りない」せいではありません。その仕事の専属AIが、いなかっただけ。',
         Inches(0.75), Inches(6.82), Inches(11.8), Inches(0.4), font=F_MINCHO, size=14, color=GOLD_HI, italic=True)
    footer(s)


def slide_thesis():
    s = _slide(); bg(s)
    circle(s, Inches(11.0), Inches(3.4), Inches(3.0), fill=GOLD, alpha=0.05)
    eyebrow(s, 'THE ANSWER')
    text(s, 'だから、6つの専属AIを\nつくりました。', Inches(0.75), Inches(1.3), Inches(11.8), Inches(2.0),
         font=F_MINCHO, size=40, color=INK, bold=True, line_spacing=1.15)
    text(s, '1つの「何でもできるツール」ではなく、それぞれの仕事の専用AI。\nだから深いところまで任せられます。必要なものを、1つからでも。月980円から始められます。',
         Inches(0.75), Inches(3.5), Inches(11.6), Inches(1.1), font=F_SANS, size=14.5, color=INK_DIM, line_spacing=1.5)
    # 実ロゴの6連ストリップ
    items = [('RES', 'LINEの返事', C_RES), ('GLD', 'コミュニティ', C_GLD), ('IRI', 'インスタ発信', C_IRI),
             ('LUM', 'ページで販売', C_LUM), ('PRI', 'AI秘書', C_PRI), ('CRY', '24時間接客', C_CRY)]
    x0, y = Inches(0.75), Inches(4.95)
    w = Inches(1.92); gap = Inches(0.06)
    for i, (k, role, c) in enumerate(items):
        x = x0 + i * (w + gap)
        rect(s, x, y, w, Inches(1.35), fill=c, alpha=0.10, rounded=True)
        rect(s, x, y, w, Emu(30000), fill=c)
        ih = 0.34 if k != 'CRY' else 0.52
        path, aspect = ICON[k]
        iw = ih * aspect
        s.shapes.add_picture(path, x + (w - Inches(iw)) / 2, y + Inches(0.22 if k != 'CRY' else 0.13),
                             height=Inches(ih))
        text(s, role, x, y + Inches(0.78), w, Inches(0.4),
             font=F_GOTHIC, size=11.5, color=INK, bold=True, align='center')
    text(s, '“ 全部入り ” ではなく、“ 全部それぞれに一流 ” を。', Inches(0.75), Inches(6.6), Inches(11.4), Inches(0.5),
         font=F_MINCHO, size=18, color=GOLD_HI, italic=True)
    footer(s)


# ══════════════════════════════════════════════
#  製品 6枚組
# ══════════════════════════════════════════════
def slide_divider(p):
    s = _slide(); bg(s)
    a = p['accent']
    circle(s, Inches(11.3), Inches(1.9), Inches(2.6), fill=a, alpha=0.07)
    circle(s, Inches(1.5), Inches(6.5), Inches(1.9), fill=a, alpha=0.05)
    text(s, p['no'], Inches(0.62), Inches(0.3), Inches(4), Inches(2.3),
         font=F_NUM, size=130, color=BG_NUM, bold=True)
    circle(s, Inches(0.98), Inches(2.85), Inches(0.06), fill=a)
    text(s, p['en'] + '   ―   ' + p['area'], Inches(1.2), Inches(2.69), Inches(10), Inches(0.4),
         font=F_OPTIMA, size=13, color=a, bold=True, spacing=280)
    # 実ロゴ (大)
    if p['key'] == 'CRY':
        w = pic(s, 'CRY', Inches(0.72), Inches(3.18), 1.35)
        text(s, 'Crystal', Inches(0.72 + w + 0.3), Inches(3.18), Inches(6), Inches(1.35),
             font=F_DIDOT, size=64, color=INK, bold=True, anchor='middle')
    else:
        pic(s, p['key'], Inches(0.72), Inches(3.3), 1.05)
    text(s, p['promise'], Inches(0.76), Inches(4.75), Inches(11.6), Inches(1.15),
         font=F_MINCHO, size=18.5, color=INK_DIM, line_spacing=1.4)
    rect(s, Inches(0.75), Inches(6.1), Inches(11.85), Emu(9000), fill=LINE)
    for i, (num, lab) in enumerate(p['metrics']):
        x = Inches(0.75) + i * Inches(4.0)
        text(s, num, x, Inches(6.28), Inches(3.8), Inches(0.5),
             font=F_NUM, size=24, color=a, bold=True)
        text(s, lab, x, Inches(6.8), Inches(3.85), Inches(0.35),
             font=F_SANS, size=9.5, color=INK_FAINT)
    footer(s)


def slide_problem(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p)
    eyebrow(s, 'THE REALITY', accent=a, no=p['no'])
    text(s, p['problem_head'], Inches(0.75), Inches(1.18), Inches(11.6), Inches(1.6),
         font=F_MINCHO, size=31, color=INK, bold=True, line_spacing=1.18)
    text(s, p['problem_sub'], Inches(0.78), Inches(2.82), Inches(11), Inches(0.4),
         font=F_SANS, size=12.5, color=INK_DIM)
    cw, ch, gx, gy = Inches(5.85), Inches(1.72), Inches(0.3), Inches(0.26)
    x0, y0 = Inches(0.75), Inches(3.35)
    for i, (t, d) in enumerate(p['pains']):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y + Inches(0.26), Emu(38000), Inches(1.2), fill=RED_SOFT)
        text(s, t, x + Inches(0.4), y + Inches(0.28), cw - Inches(0.7), Inches(0.44),
             font=F_GOTHIC, size=15.5, color=INK, bold=True)
        text(s, d, x + Inches(0.4), y + Inches(0.78), cw - Inches(0.75), Inches(0.85),
             font=F_SANS, size=11.5, color=INK_DIM, line_spacing=1.35)
    footer(s)


def _caps_slide(p, caps, part, start_no):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p)
    eyebrow(s, f'WHAT IT DOES ({part}/2)', accent=a, no=p['no'])
    title = 'できること。' if part == '1' else 'できること、まだあります。'
    text(s, title, Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    cw, ch = Inches(5.88), Inches(1.52)
    gx, gy = Inches(0.26), Inches(0.18)
    x0, y0 = Inches(0.75), Inches(2.05)
    for i, (name, desc) in enumerate(caps):
        r, c = divmod(i, 2)
        x = x0 + c * (cw + gx); y = y0 + r * (ch + gy)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x + Inches(0.22), y + Inches(0.22), Inches(0.52), Inches(0.42),
             fill=a, alpha=0.16, rounded=True, radius=0.25)
        text(s, f'{start_no + i:02d}', x + Inches(0.22), y + Inches(0.24), Inches(0.52), Inches(0.38),
             font=F_NUM, size=13, color=a, bold=True, align='center', anchor='middle')
        text(s, name, x + Inches(0.9), y + Inches(0.21), cw - Inches(1.1), Inches(0.42),
             font=F_GOTHIC, size=13.5, color=INK, bold=True)
        text(s, desc, x + Inches(0.9), y + Inches(0.66), cw - Inches(1.16), Inches(0.8),
             font=F_SANS, size=10.2, color=INK_DIM, line_spacing=1.3)
    footer(s)


def slide_transform(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p)
    eyebrow(s, 'BEFORE  /  AFTER', accent=a, no=p['no'])
    text(s, '入れると、毎日がこう変わる。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    text(s, 'いままで', Inches(0.9), Inches(2.14), Inches(4), Inches(0.3),
         font=F_GOTHIC, size=12, color=RED_SOFT, bold=True)
    text(s, 'これから', Inches(7.55), Inches(2.14), Inches(4), Inches(0.3),
         font=F_GOTHIC, size=12, color=a, bold=True)
    y0 = Inches(2.58); rh = Inches(1.36)
    bw = Inches(5.72)
    for i, (bfr, aft) in enumerate(p['before_after']):
        y = y0 + i * (rh + Inches(0.14))
        rect(s, Inches(0.75), y, bw, rh, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, bfr, Inches(1.08), y, bw - Inches(0.65), rh, font=F_SANS, size=13, color=INK_DIM,
             anchor='middle', line_spacing=1.32)
        text(s, '→', Inches(6.52), y, Inches(0.85), rh, font=F_NUM, size=22, color=a,
             align='center', anchor='middle')
        rect(s, Inches(7.38), y, bw, rh, fill=a, alpha=0.13, rounded=True)
        rect(s, Inches(7.38), y, Emu(40000), rh, fill=a)
        text(s, aft, Inches(7.72), y, bw - Inches(0.65), rh, font=F_GOTHIC, size=13, color=INK,
             bold=True, anchor='middle', line_spacing=1.32)
    footer(s)


def slide_offer(p):
    s = _slide(); bg(s)
    a = p['accent']
    running_head(s, p)
    eyebrow(s, 'WHY  &  PRICE', accent=a, no=p['no'])
    text(s, '選ばれる理由と、料金。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    lx, lw = Inches(0.75), Inches(6.15)
    text(s, 'ここが違います', lx, Inches(2.06), lw, Inches(0.35),
         font=F_GOTHIC, size=13.5, color=a, bold=True)
    for i, d in enumerate(p['diff']):
        y = Inches(2.52) + i * Inches(0.74)
        circle(s, lx + Inches(0.08), y + Inches(0.14), Inches(0.05), fill=a)
        text(s, d, lx + Inches(0.32), y, lw - Inches(0.4), Inches(0.68),
             font=F_SANS, size=12.5, color=INK, line_spacing=1.3)
    iy = Inches(2.52) + len(p['diff']) * Inches(0.74) + Inches(0.12)
    rect(s, lx, iy, lw, Inches(1.15), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'こんな方に', lx + Inches(0.3), iy + Inches(0.16), lw - Inches(0.5), Inches(0.3),
         font=F_GOTHIC, size=10.5, color=a, bold=True)
    text(s, p['ideal'], lx + Inches(0.3), iy + Inches(0.46), lw - Inches(0.6), Inches(0.65),
         font=F_SANS, size=11, color=INK_DIM, line_spacing=1.3)
    rx, rw = Inches(7.25), Inches(5.35)
    rect(s, rx, Inches(2.06), rw, Inches(4.55), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, '料金', rx + Inches(0.35), Inches(2.3), rw - Inches(0.7), Inches(0.32),
         font=F_GOTHIC, size=12, color=a, bold=True)
    plans = p['pricing']
    py0 = Inches(2.74); prow = Inches(3.35) / max(len(plans), 1)
    for i, (plan, price, note) in enumerate(plans):
        y = py0 + prow * i
        if i > 0:
            rect(s, rx + Inches(0.35), y, rw - Inches(0.7), Emu(9000), fill=LINE)
        text(s, plan, rx + Inches(0.35), y + Inches(0.09), Inches(2.6), Inches(0.4),
             font=F_GOTHIC, size=13.5, color=INK, bold=True)
        text(s, price, rx + Inches(0.35), y + Inches(0.07), rw - Inches(0.7), Inches(0.42),
             font=F_NUM, size=17, color=GOLD_HI, bold=True, align='right')
        text(s, note, rx + Inches(0.35), y + Inches(0.5), rw - Inches(0.7), Inches(0.35),
             font=F_SANS, size=9.5, color=INK_DIM)
    text(s, p['price_note'], rx + Inches(0.35), Inches(6.18), rw - Inches(0.7), Inches(0.35),
         font=F_SANS, size=9.5, color=a)
    footer(s)


# ══════════════════════════════════════════════
#  CLOSING
# ══════════════════════════════════════════════
def slide_why_now():
    s = _slide(); bg(s)
    circle(s, Inches(11.0), Inches(3.2), Inches(2.9), fill=GOLD, alpha=0.05)
    eyebrow(s, 'WHY NOW')
    text(s, '世界ではもう、AIの接客に\nこれだけ払っています。', Inches(0.75), Inches(1.15), Inches(11.8), Inches(1.7),
         font=F_MINCHO, size=32, color=INK, bold=True, line_spacing=1.16)
    cards = [
        ('Intercom Fin', '1回 $0.99', '問い合わせに答えるたびにお金がかかる仕組み。使うほど請求が増える。'),
        ('Qualified', '年 $68,000〜', 'AI営業係の年間費用。日本円で約1,000万円。それでも売れている。'),
        ('Sierra', '大企業が続々導入', '世界の名だたる企業が、AI接客に本気でお金をかけ始めている。'),
    ]
    cw, ch, gx = Inches(3.86), Inches(2.0), Inches(0.26)
    x0, y0 = Inches(0.75), Inches(3.25)
    for i, (t, big, d) in enumerate(cards):
        x = x0 + i * (cw + gx)
        rect(s, x, y0, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, t, x + Inches(0.32), y0 + Inches(0.26), cw - Inches(0.6), Inches(0.35),
             font=F_DIDOT, size=15, color=INK, bold=True)
        text(s, big, x + Inches(0.32), y0 + Inches(0.68), cw - Inches(0.6), Inches(0.5),
             font=F_NUM, size=22, color=GOLD_HI, bold=True)
        text(s, d, x + Inches(0.32), y0 + Inches(1.26), cw - Inches(0.62), Inches(0.68),
             font=F_SANS, size=10.5, color=INK_DIM, line_spacing=1.3)
    text(s, 'CORE Continuum は、同じ水準の仕事を ―― 日本語の品質で、月額だけで。',
         Inches(0.75), Inches(5.75), Inches(11.8), Inches(0.6),
         font=F_MINCHO, size=18, color=GOLD_HI, italic=True)
    footer(s)


def slide_journey():
    s = _slide(); bg(s)
    eyebrow(s, 'ONE JOURNEY, SIX AGENTS')
    text(s, 'お客様との出会いから、ファンになるまで。', Inches(0.75), Inches(1.14), Inches(11.6), Inches(0.75),
         font=F_MINCHO, size=27, color=INK, bold=True)
    text(s, '知ってもらう→買ってもらう→また来てもらう→仲間になってもらう。6つのAIが、途切れずにバトンをつなぎます。',
         Inches(0.75), Inches(1.95), Inches(11.9), Inches(0.5), font=F_SANS, size=12.5, color=INK_DIM)
    steps = [
        ('IRI', '知ってもらう', 'インスタ発信をAIが量産', C_IRI),
        ('LUM', '買ってもらう', 'プロフィールページでその場で販売', C_LUM),
        ('CRY', 'おもてなし', 'サイトに来た人を24時間接客', C_CRY),
        ('RES', 'また来てもらう', 'LINEでひとりずつに返事', C_RES),
        ('GLD', '仲間になってもらう', 'ファンをコミュニティに', C_GLD),
        ('PRI', '全体をまとめる', '情報を束ねて、次の一手を提案', C_PRI),
    ]
    cw = Inches(1.93); gap = Inches(0.06)
    x0, y, ch = Inches(0.75), Inches(2.75), Inches(3.35)
    for i, (k, role, d, col) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, ch, fill=BG_PANEL, line=LINE, rounded=True)
        rect(s, x, y, cw, Emu(42000), fill=col)
        text(s, f'{i+1:02d}', x, y + Inches(0.24), cw, Inches(0.42),
             font=F_NUM, size=18, color=col, bold=True, align='center')
        ih = 0.34 if k != 'CRY' else 0.55
        path, aspect = ICON[k]
        iw = ih * aspect
        s.shapes.add_picture(path, x + (cw - Inches(iw)) / 2, y + Inches(0.82 if k != 'CRY' else 0.72),
                             height=Inches(ih))
        text(s, role, x + Inches(0.06), y + Inches(1.42), cw - Inches(0.12), Inches(0.4),
             font=F_GOTHIC, size=12, color=col, bold=True, align='center')
        text(s, d, x + Inches(0.14), y + Inches(1.95), cw - Inches(0.28), Inches(1.2),
             font=F_SANS, size=9.5, color=INK_DIM, align='center', line_spacing=1.3)
        if i < 5:
            text(s, '›', x + cw - Inches(0.05), y + Inches(1.3), Inches(0.18), Inches(0.5),
                 font=F_NUM, size=15, color=INK_FAINT, align='center')
    text(s, '1つだけでも、強い。組み合わせると、こぼれない。', Inches(0.75), Inches(6.45), Inches(11.8), Inches(0.5),
         font=F_MINCHO, size=16, color=GOLD_HI, italic=True)
    footer(s)


def slide_matrix():
    s = _slide(); bg(s)
    eyebrow(s, 'WHO NEEDS WHAT')
    text(s, 'あなたには、どれが合う?', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    rows = [
        ('お店・サロン・整体・教室', 'Resonance ＋ Crystal', 'LINEの返事とサイトの接客を、AIがまとめて引き受け', (C_RES, C_CRY)),
        ('オンラインサロン・コミュニティ', 'Guild', 'みんなで決めて、頑張りが見える場に変わる', (C_GLD,)),
        ('クリエイター・インフルエンサー', 'Iris ＋ Lume', '投稿づくりから、ページでの販売までひとつながり', (C_IRI, C_LUM)),
        ('経営者・士業・フリーランス', 'Prism', 'あなたの仕事を覚えた秘書が、24時間そばに', (C_PRI,)),
        ('高級不動産・ホテル・接客業', 'Crystal', 'ブランドの言葉づかいのまま、夜の来客を逃さない', (C_CRY,)),
        ('インスタ運用代行の会社', 'Iris (Agency)', '30アカウントまとめて管理。台本づくりの人件費がゼロに', (C_IRI,)),
    ]
    x0, y0, rh = Inches(0.75), Inches(1.98), Inches(0.7)
    cols = [Inches(3.6), Inches(3.15), Inches(5.1)]
    tw = Inches(3.6) + Inches(3.15) + Inches(5.1)
    rect(s, x0, y0, tw, rh, fill=GOLD, alpha=0.13)
    for i, h in enumerate(['あなたは…', 'おすすめ', 'どう変わる?']):
        cx = x0
        for j in range(i):
            cx += cols[j]
        text(s, h, cx + Inches(0.25), y0, cols[i] - Inches(0.4), rh, font=F_GOTHIC, size=12.5,
             color=GOLD, bold=True, anchor='middle')
    for r, (a1, a2, a3, dots) in enumerate(rows):
        y = y0 + rh * (r + 1)
        text(s, a1, x0 + Inches(0.25), y, cols[0] - Inches(0.4), rh, font=F_SANS, size=12,
             color=INK, anchor='middle')
        dx = x0 + cols[0] + Inches(0.25)
        for k, dc in enumerate(dots):
            circle(s, dx + Inches(0.06) + k * Inches(0.24), y + rh / 2, Inches(0.05), fill=dc)
        text(s, a2, dx + Inches(0.16) + len(dots) * Inches(0.24), y, cols[1] - Inches(0.7), rh,
             font=F_GOTHIC, size=12.5, color=GOLD_HI, bold=True, anchor='middle')
        text(s, a3, x0 + cols[0] + cols[1] + Inches(0.25), y, cols[2] - Inches(0.4), rh,
             font=F_SANS, size=11.5, color=INK_DIM, anchor='middle', line_spacing=1.2)
        rect(s, x0, y, tw, Emu(9000), fill=LINE)
    footer(s)


def slide_pricing_all():
    s = _slide(); bg(s)
    eyebrow(s, 'PRICING')
    text(s, '料金一覧 ― ぜんぶ月額。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    rows = [
        ('RES', 'Resonance', 'Solo ¥1,980', 'Pro ¥6,980', 'Business ¥14,800', C_RES),
        ('GLD', 'Guild', '無料 ¥0', 'Community ¥980', 'Pro ¥3,980', C_GLD),
        ('IRI', 'Iris', 'Lite ¥2,980', 'Standard ¥6,980', 'Agency ¥29,800', C_IRI),
        ('LUM', 'Lume', '無料 (7日間)', 'Pro ¥1,480', 'Business ¥3,480', C_LUM),
        ('PRI', 'Prism', 'Starter ¥2,980', 'Standard ¥9,800', 'Exclusive ¥29,800', C_PRI),
        ('CRY', 'Crystal', '―', 'Standard ¥29,800', 'Luxury ¥49,800', C_CRY),
    ]
    x0, y0, rh = Inches(0.75), Inches(2.0), Inches(0.63)
    cols = [Inches(2.75), Inches(2.95), Inches(2.95), Inches(3.2)]
    tw = Inches(2.75) + Inches(2.95) + Inches(2.95) + Inches(3.2)
    rect(s, x0, y0, tw, rh, fill=GOLD, alpha=0.13)
    for i, h in enumerate(['サービス', 'はじめる', 'ふつうに使う', 'しっかり使う']):
        cx = x0
        for j in range(i):
            cx += cols[j]
        text(s, h, cx + Inches(0.22), y0, cols[i] - Inches(0.35), rh, font=F_GOTHIC, size=12,
             color=GOLD, bold=True, anchor='middle')
    for r, row in enumerate(rows):
        y = y0 + rh * (r + 1)
        k = row[0]
        ih = 0.22 if k != 'CRY' else 0.34
        path, aspect = ICON[k]
        s.shapes.add_picture(path, x0 + Inches(0.24), y + (rh - Inches(ih)) / 2, height=Inches(ih))
        text(s, row[1], x0 + Inches(0.72), y, cols[0] - Inches(0.8), rh,
             font=F_DIDOT, size=14.5, color=INK, bold=True, anchor='middle')
        for i in range(2, 5):
            cx = x0
            for j in range(i - 1):
                cx += cols[j]
            text(s, row[i], cx + Inches(0.22), y, cols[i - 1] - Inches(0.35), rh, font=F_SANS,
                 size=12, color=INK if i == 4 else INK_DIM, bold=(i == 4), anchor='middle')
        rect(s, x0, y, tw, Emu(9000), fill=LINE)
    text(s, '※ 税込・月額。Crystal 以外は無料か7日間のお試しから。Crystalのみ初期費用(設置の代行つき)があります。',
         Inches(0.75), Inches(6.55), Inches(11.8), Inches(0.4), font=F_SANS, size=10, color=INK_FAINT)
    footer(s)


def slide_start():
    s = _slide(); bg(s)
    eyebrow(s, 'HOW TO START')
    text(s, 'はじめ方は、3つだけ。', Inches(0.75), Inches(1.14), Inches(11), Inches(0.75),
         font=F_MINCHO, size=28, color=INK, bold=True)
    steps = [
        ('01', '無料で触ってみる', 'カード登録はいりません。実際の画面を、そのまま試せます。'),
        ('02', 'あなたの情報を貼る', '会社案内や自己紹介の文章を貼るだけ。AIがあなた仕様になります。'),
        ('03', 'リンクを貼って、完了', '難しい設定はありません。困ったら、設置の代行もあります。'),
    ]
    cw, gap = Inches(3.88), Inches(0.24)
    x0, y, h = Inches(0.75), Inches(2.3), Inches(3.5)
    for i, (num, t, d) in enumerate(steps):
        x = x0 + i * (cw + gap)
        rect(s, x, y, cw, h, fill=BG_PANEL, line=LINE, rounded=True)
        text(s, num, x + Inches(0.38), y + Inches(0.34), Inches(2), Inches(0.9),
             font=F_NUM, size=40, color=GOLD, bold=True)
        rect(s, x + Inches(0.42), y + Inches(1.3), Inches(0.6), Emu(24000), fill=GOLD_HI)
        text(s, t, x + Inches(0.42), y + Inches(1.52), cw - Inches(0.84), Inches(0.6),
             font=F_GOTHIC, size=19, color=INK, bold=True)
        text(s, d, x + Inches(0.42), y + Inches(2.18), cw - Inches(0.84), Inches(1.2),
             font=F_SANS, size=12, color=INK_DIM, line_spacing=1.42)
        if i < 2:
            text(s, '→', x + cw - Inches(0.02), y + h / 2 - Inches(0.3), gap + Inches(0.06), Inches(0.6),
                 font=F_NUM, size=18, color=GOLD, align='center')
    text(s, 'どのサービスも、買う前に「あなたの設定のまま」試せます。',
         Inches(0.75), Inches(6.2), Inches(11.8), Inches(0.45),
         font=F_MINCHO, size=15, color=GOLD_HI, italic=True)
    footer(s)


def slide_cta():
    s = _slide(); bg(s)
    circle(s, Inches(SLIDE_W / 2), Inches(SLIDE_H / 2), Inches(3.3), fill=GOLD, alpha=0.06)
    eyebrow(s, 'GET STARTED')
    text(s, 'まずは1つ、\n無料で試してください。', Inches(0.75), Inches(1.6), Inches(11.8), Inches(1.9),
         font=F_MINCHO, size=38, color=INK, bold=True, line_spacing=1.18)
    text(s, '気になるサービス1つからでも、組み合わせでも。\nデモ・お見積り・ご相談はメール1通で ―― 24時間以内にお返事します。',
         Inches(0.75), Inches(3.7), Inches(11.4), Inches(1.0), font=F_SANS, size=14.5, color=INK_DIM, line_spacing=1.45)
    grad_bar(s, Inches(0.78), Inches(4.95), Inches(2.4), Inches(0.05))
    rect(s, Inches(0.75), Inches(5.35), Inches(6.7), Inches(1.05), fill=BG_PANEL, line=LINE, rounded=True)
    text(s, 'お問い合わせ', Inches(1.1), Inches(5.55), Inches(3), Inches(0.3),
         font=F_GOTHIC, size=10.5, color=GOLD, bold=True)
    text(s, 'core.guild.inc@gmail.com', Inches(1.1), Inches(5.85), Inches(6), Inches(0.45),
         font=F_NUM, size=19, color=INK, bold=True)
    footer(s)


def slide_closing():
    s = _slide(); bg(s)
    circle(s, Inches(SLIDE_W / 2), Inches(3.3), Inches(4.6), fill=GOLD, alpha=0.035)
    # 実ロゴアイコン6連
    order = ['RES', 'GLD', 'IRI', 'LUM', 'PRI', 'CRY']
    total_w = 0
    ws = []
    for k in order:
        ih = 0.3 if k != 'CRY' else 0.42
        path, aspect = ICON[k]
        ws.append((k, path, ih, ih * aspect))
        total_w += ih * aspect + 0.45
    total_w -= 0.45
    cx = (SLIDE_W - total_w) / 2
    for k, path, ih, iw in ws:
        s.shapes.add_picture(path, Inches(cx), Inches(2.05 if k != 'CRY' else 1.99), height=Inches(ih))
        cx += iw + 0.45
    text(s, 'CORE  CONTINUUM', Inches(0.5), Inches(2.75), Inches(SLIDE_W - 1), Inches(1.0),
         font=F_DIDOT, size=46, color=INK, bold=True, align='center', spacing=300)
    text(s, '事業のあらゆる接点に、途切れない知性を。', Inches(0.5), Inches(3.95), Inches(SLIDE_W - 1), Inches(0.6),
         font=F_MINCHO, size=20, color=GOLD_HI, italic=True, align='center')
    grad_bar(s, Inches(SLIDE_W / 2 - 1.3), Inches(4.85), Inches(2.6), Inches(0.05))
    footer(s)


# ── 組み立て ───────────────────────────
slide_cover()
slide_problem_intro()
slide_thesis()
for p in PRODUCTS:
    slide_divider(p)
    slide_problem(p)
    _caps_slide(p, p['caps'][:6], '1', 1)
    _caps_slide(p, p['caps'][6:12], '2', 7)
    slide_transform(p)
    slide_offer(p)
slide_why_now()
slide_journey()
slide_matrix()
slide_pricing_all()
slide_start()
slide_cta()
slide_closing()

out = os.path.expanduser('~/Desktop/CORE_Continuum_6サービス_営業資料.pptx')
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides._sldIdLst))
